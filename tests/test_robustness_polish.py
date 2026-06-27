"""Phase B/C polish-cluster regression tests (AUDIT_AND_REMEDIATION_PLAN_2026-06-26).

  * B6 — _fit_body hard-truncates with an ellipsis at the shrink-to-fit floor so a
    too-long passage can't silently bleed past its box (it returns short text intact).
  * B7 — the course->PowerPoint flatten prepends "Select all that apply." to a
    multi-select knowledge check (several ✓ marks are ambiguous without the cue).
  * C5 — the source-ingestion readers refuse an over-large file / oversized source
    set / decompression-bomb .odt with an actionable skip note instead of OOM.

(B8 — player a11y — is JS; node --check + the existing tests/test_player.js guard it.)
"""
import zipfile

import authoring
import md_import
import pptx_export
import slide_layouts

_EMU = 914400  # EMU per inch


# --- B6: shrink-to-fit floor truncation --------------------------------------

def test_fit_body_returns_fitting_text_unchanged():
    pt, txt = slide_layouts._fit_body("A short line.", 4 * _EMU, 2 * _EMU, 14, 9)
    assert txt == "A short line."          # fits comfortably -> untouched
    assert 9 <= pt <= 14


def test_fit_body_truncates_overflow_at_floor():
    long = "word " * 400                    # ~2000 chars
    # tiny box + a high floor: even at min_pt the text cannot fit -> must truncate.
    pt, txt = slide_layouts._fit_body(long, _EMU // 2, _EMU // 4, 14, 11)
    assert txt.endswith("…")
    assert len(txt) < len(long)
    cap = slide_layouts._box_capacity_chars(_EMU // 2, _EMU // 4, pt)
    assert len(txt) <= cap                  # never exceeds the box's char capacity


def test_fit_body_empty_text_is_safe():
    pt, txt = slide_layouts._fit_body("", 2 * _EMU, _EMU, 12, 8)
    assert txt == "" and pt == 12


# --- B7: PPTX multi-select "select all that apply" cue -----------------------

def _kc_pptx_text(block):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])     # blank layout
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    pptx_export._render_kc(tb.text_frame, block, pptx_export.NAVY, [True])  # mutable first-para flag
    return tb.text_frame.text


def test_multi_kc_flatten_has_select_all_cue():
    block = md_import._knowledge_check(
        "*Question:* Q\n- A) a\n- B) b\n- C) c\n*Correct Answer:* A, B\n")
    assert block["multi"] is True
    assert "Select all that apply." in _kc_pptx_text(block)


def test_single_kc_flatten_has_no_cue():
    block = md_import._knowledge_check(
        "*Question:* Q\n- A) a\n- B) b\n*Correct Answer:* B\n")
    assert block["multi"] is False
    assert "Select all that apply." not in _kc_pptx_text(block)


# --- C5: ingestion bounds ----------------------------------------------------

def test_read_sources_skips_oversized_file(tmp_path, monkeypatch):
    monkeypatch.setattr(authoring, "_MAX_SOURCE_BYTES", 10)
    (tmp_path / "big.txt").write_text("x" * 100)
    _text, used, skipped = authoring.read_sources(str(tmp_path))
    assert used == []
    assert any("big.txt" in s and "larger than" in s for s in skipped)


def test_read_sources_total_cap_stops_later_files(tmp_path, monkeypatch):
    monkeypatch.setattr(authoring, "_MAX_SOURCE_BYTES", 10_000)
    monkeypatch.setattr(authoring, "_MAX_TOTAL_SOURCE_BYTES", 120)
    (tmp_path / "a.txt").write_text("a" * 100)   # 100B fits (total 100)
    (tmp_path / "b.txt").write_text("b" * 100)   # would push total to 200 > 120
    _text, used, skipped = authoring.read_sources(str(tmp_path))
    assert "a.txt" in used
    assert any("b.txt" in s and "total cap" in s for s in skipped)


def test_odt_decompression_bomb_is_refused(tmp_path, monkeypatch):
    monkeypatch.setattr(authoring, "_MAX_DECOMPRESSED_BYTES", 50)
    p = tmp_path / "bomb.odt"
    with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("content.xml", "<x>" + ("A" * 500) + "</x>")   # 500B decompressed
    assert authoring._odt_to_text(str(p)) == ""                   # refused, no crash


def test_odt_small_reads_normally(tmp_path):
    p = tmp_path / "ok.odt"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("content.xml", "<text:p>hello world</text:p>")
    assert "hello world" in authoring._odt_to_text(str(p))
