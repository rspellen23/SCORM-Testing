"""Workshop template family (region-specs in templates/layouts/workshop-*.json).

These are our own original, brand-agnostic slide types (Step 3 of template
ingestion): agenda, learning objectives, section divider, hands-on activity.
Proves each one loads, registers as a generic ("default") pickable template,
renders the same spec under two different brands (structure is brand-agnostic),
renders a bulleted list from a `list` bind, paints the section divider's
full-bleed brand fill with reversed white text, and renders in the SVG preview.
"""
import pytest

import slide_layouts
import slide_svg
import brand as brandmod

pptx = pytest.importorskip("pptx")  # export extra

WORKSHOP = ["workshop-agenda", "workshop-objectives",
            "workshop-section", "workshop-activity"]


def _pptx_texts(path):
    from pptx import Presentation
    p = Presentation(path)
    return " ".join(sh.text_frame.text for s in p.slides
                    for sh in s.shapes if sh.has_text_frame)


def test_all_workshop_specs_listed_and_loaded():
    names = slide_layouts.list_template_layouts()
    for n in WORKSHOP:
        assert n in names, n
        spec = slide_layouts.load_template_spec(n)
        assert spec["name"] == n and spec["regions"]


def test_workshop_specs_are_generic_pickable_templates():
    info = {r["name"]: r for r in slide_layouts.template_layout_info()}
    for n in WORKSHOP:
        rec = info[n]
        assert rec["category"] == "default"   # generic: group under "Templates", any brand
        assert rec["brand"] is None           # not client-locked
        assert rec["starter"]["template"] == n


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_agenda_renders_bulleted_list_in_any_brand(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"agenda_{brand_name}.pptx")
    content = {"template": "workshop-agenda", "kicker": "Agenda",
               "title": "Today's session",
               "items": ["Welcome and context", "Core concepts",
                         "Hands-on activity", "Wrap and next steps"]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Agenda", "Today's session", "Welcome and context", "Wrap and next steps"):
        assert w in txt


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_objectives_renders_title_subtitle_list(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"obj_{brand_name}.pptx")
    content = {"template": "workshop-objectives", "title": "Learning objectives",
               "subtitle": "By the end you will be able to:",
               "items": ["Identify the dashboard", "Configure a rule",
                         "Troubleshoot a failure"]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Learning objectives", "By the end", "Configure a rule",
              "Troubleshoot a failure"):
        assert w in txt


def test_section_divider_full_fill_with_reverse_white_text(tmp_path):
    # The section divider is a bold full-bleed brand-color fill carrying the
    # kicker + title in REVERSED white text (the per-region `color` override).
    br = brandmod.load_brand("teletracking")
    pal = slide_layouts._palette_of(br)
    out = str(tmp_path / "section.pptx")
    content = {"template": "workshop-section", "kicker": "Section 1",
               "title": "Getting started"}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    from pptx import Presentation
    p = Presentation(out)
    slide = p.slides[0]
    txt = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "Section 1" in txt and "Getting started" in txt
    # every text run on the divider is reversed white, not the dark ink default
    run_colors = {str(r.font.color.rgb)
                  for sh in slide.shapes if sh.has_text_frame
                  for para in sh.text_frame.paragraphs for r in para.runs
                  if r.text.strip()}
    assert run_colors == {str(pal["white"])}, run_colors
    # a full-bleed brand-color fill block exists (auto shape spanning the slide)
    full = [sh for sh in slide.shapes if sh.shape_type == 1
            and sh.width >= int(slide_layouts.W * 0.98)
            and sh.height >= int(slide_layouts.H * 0.98)]
    assert full, "expected a full-bleed background fill"


def test_activity_renders_steps_and_tolerates_empty_image(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "activity.pptx")
    content = {"template": "workshop-activity", "title": "Try it yourself",
               "subtitle": "Create your first transfer request.",
               "steps": ["Open the queue", "Click New request", "Fill the form"]}
    # no `image` key -> the image region degrades to a placeholder card, never crashes
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Try it yourself", "Create your first", "Open the queue", "Fill the form"):
        assert w in txt


@pytest.mark.parametrize("name,content", [
    ("workshop-agenda", {"kicker": "Agenda", "title": "Plan",
                         "items": ["Intro", "Body", "Close"]}),
    ("workshop-objectives", {"title": "Objectives", "subtitle": "You will:",
                             "items": ["Do X", "Do Y"]}),
    ("workshop-section", {"kicker": "Section 2", "title": "Deep dive"}),
    ("workshop-activity", {"title": "Exercise", "subtitle": "Try this",
                           "steps": ["Step one", "Step two"]}),
])
def test_each_workshop_spec_renders_in_svg_preview(name, content):
    br = brandmod.load_brand("teletracking")
    content = dict(content, template=name)
    svg = slide_svg.render_slide_svg("template", content, br)
    assert svg.startswith("<svg")
    # preview wraps long titles into per-word <tspan>s; words stay intact
    assert content["title"].split()[0] in svg
