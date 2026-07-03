"""Read a foreign .pptx and extract its DESIGN DNA — theme palette + fonts, slide
size/aspect, layout inventory, colors-in-use, and a per-slide shape summary — into a
structured profile dict.

This is the FIRST half of template ingestion: a brand-mapping step (theme palette ->
brand.json + a starter blueprint) and an "author our own version of this look" step
both build on this profile. Read-only — it never writes to or modifies the source.

LICENSE NOTE: this extracts a STRUCTURAL PROFILE (colors, fonts, slide size, layout
names, geometry, text snippets) to inform building our OWN original templates. It does
not copy or redistribute the source artwork. General genre/structure isn't protected;
specific creative expression is — keep authored output original.
"""
import json
import os
import re
import zipfile
from collections import Counter

# The standard theme color slots (ECMA-376 a:clrScheme).
_THEME_SLOTS = ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3",
                "accent4", "accent5", "accent6", "hlink", "folHlink")


def _theme(path):
    """Theme color scheme + font scheme from ppt/theme/theme1.xml. python-pptx has no
    clean theme API, so the theme part is parsed directly. Returns ({slot: #hex}, {major,minor})."""
    try:
        z = zipfile.ZipFile(path)
    except (OSError, zipfile.BadZipFile):
        return {}, {}
    names = sorted(n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml", n))
    if not names:
        return {}, {}
    xml = z.read(names[0]).decode("utf-8", "ignore")
    palette = {}
    blk = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
    if blk:
        for slot in _THEME_SLOTS:
            sm = re.search(r"<a:%s>(.*?)</a:%s>" % (slot, slot), blk.group(0), re.S)
            if sm:
                hx = re.search(r'(?:srgbClr val="|lastClr=")([0-9A-Fa-f]{6})', sm.group(1))
                if hx:
                    palette[slot] = "#" + hx.group(1).upper()
    fonts = {}
    for which in ("majorFont", "minorFont"):
        fm = re.search(r'<a:%s>\s*<a:latin typeface="([^"]*)"' % which, xml)
        if fm and fm.group(1):
            fonts[which.replace("Font", "")] = fm.group(1)
    return palette, fonts


def _in(emu):
    return round(emu / 914400.0, 3)


def _shape_fill_hex(shp):
    """An explicit solid RGB fill on a shape, as #RRGGBB, or None (theme/none/picture)."""
    try:
        fill = shp.fill
        if fill.type is None:
            return None
        fc = fill.fore_color
        if fc is None or fc.type is None:
            return None
        return "#" + str(fc.rgb).upper()
    except Exception:
        return None


def _rel_rect(shp, sw, sh):
    """Shape position as a region-spec RELATIVE rect [x, y, w, h] (each a 0..1 fraction of
    slide width/height, matching templates/layouts/*.json). Returns None when any coordinate
    is unset (a placeholder inheriting its geometry from the layout — geometry unknown here)."""
    if not (sw and sh):
        return None
    left, top, width, height = shp.left, shp.top, shp.width, shp.height
    if None in (left, top, width, height):
        return None
    return [round(left / sw, 4), round(top / sh, 4),
            round(width / sw, 4), round(height / sh, 4)]


def _ph_type(shp):
    """The placeholder type name (TITLE, BODY, PICTURE, ...) if the shape is a placeholder."""
    try:
        if shp.is_placeholder:
            return str(shp.placeholder_format.type)
    except Exception:
        pass
    return None


def inspect_pptx(path):
    """Extract a foreign .pptx into a design-DNA profile dict. Raises on an unreadable
    file (the caller decides how to surface it)."""
    from pptx import Presentation
    prs = Presentation(path)
    palette, fonts = _theme(path)
    sw, sh = prs.slide_width, prs.slide_height

    layouts = []
    for master in prs.slide_masters:
        for lo in master.slide_layouts:
            phs = []
            for ph in lo.placeholders:
                try:
                    phs.append(str(ph.placeholder_format.type))
                except Exception:
                    phs.append("?")
            layouts.append({"name": lo.name, "placeholders": phs})

    fill_colors = Counter()
    slides = []
    for i, s in enumerate(prs.slides):
        types = Counter()
        texts = []
        shapes = []
        for shp in s.shapes:
            types[str(shp.shape_type)] += 1
            txt = ""
            if shp.has_text_frame and shp.text_frame.text.strip():
                txt = shp.text_frame.text.strip().replace("\n", " ")
                texts.append(txt[:60])
            hx = _shape_fill_hex(shp)
            if hx:
                fill_colors[hx] += 1
            # Per-shape geometry profile — the bridge to AI role-tagging (step 2c):
            # a relative rect in the same convention as a region-spec, plus the hints
            # (placeholder type, text snippet, fill) a role-tagger reads to assign a role.
            shapes.append({
                "kind": str(shp.shape_type),
                "rect": _rel_rect(shp, sw, sh),
                "placeholder": _ph_type(shp),
                "text": txt[:80],
                "fill": hx,
            })
        slides.append({
            "index": i,
            "layout": s.slide_layout.name,
            "shape_types": dict(types),
            "texts": texts[:6],
            "shapes": shapes[:40],
        })

    return {
        "name": os.path.splitext(os.path.basename(path))[0],
        "source": os.path.basename(path),
        "slide_size_in": [_in(sw), _in(sh)],
        "aspect": round(sw / sh, 3) if sh else None,
        "theme": {"palette": palette, "fonts": fonts},
        "colors_in_use": [c for c, _ in fill_colors.most_common(12)],
        "layout_count": len(layouts),
        "layouts": layouts,
        "slide_count": len(prs.slides),
        "slides": slides,
    }


# --------------------------------------------------------------- step 2c: role tagging
# Take one extracted slide (its `shapes` with relative rects) and turn it into a
# renderable region-spec (templates/layouts/<name>.json). The AI's ONLY job is to
# CLASSIFY each shape into a role from a fixed vocabulary; the geometry, binds, and
# spec assembly are deterministic Python — so the AI never controls coordinates and
# the output is testable with a stubbed runner (no live CLI). Uses the subscription
# claude CLI via authoring.run_cli, never the metered API.

# The role vocabulary the renderer (slide_layouts._render_template) understands. Each
# maps to a content bind + styling; "skip" drops a decorative/empty shape.
_ALLOWED_ROLES = ("title", "kicker", "value", "subtitle", "body",
                  "image", "cards", "background", "skip")

_ROLE_HELP = {
    "title": "the slide's main heading",
    "kicker": "a small eyebrow/label line above the title",
    "value": "a large standout number or metric",
    "subtitle": "a secondary line under the title, or a caption",
    "body": "a paragraph or bullet list of supporting text",
    "image": "a picture, logo, icon, or graphic",
    "cards": "one of several repeating equal boxes in a row (feature/step cards)",
    "background": "a full-bleed color band or backdrop shape",
    "skip": "decorative or empty; leave out of the template",
}

# role -> the content key the assembled region binds to (singular roles).
_ROLE_BIND = {"title": "title", "kicker": "kicker", "value": "value",
              "subtitle": "subtitle", "body": "body", "image": "image"}


def build_roletag_prompt(shapes, context=None):
    """A one-shot classification prompt: list each placeable shape, ask the model to
    label it with exactly one role. Returns ONLY a JSON array — no tools, no prose."""
    lines = ["You are tagging the shapes of ONE presentation slide so its structure can",
             "be reused as a reusable layout. For each shape below, choose exactly ONE",
             "role from this list (by what the shape IS, using its position and text):", ""]
    for r in _ALLOWED_ROLES:
        lines.append(f"  {r}: {_ROLE_HELP[r]}")
    if context:
        lines += ["", f"Slide context: {context}"]
    lines += ["", "Shapes (rect is [x, y, w, h] as fractions of the slide, origin top-left):"]
    for sh in shapes:
        rect = sh.get("rect")
        rs = "inherited" if rect is None else "[" + ", ".join(f"{v:.2f}" for v in rect) + "]"
        snip = (sh.get("text") or "").replace("\n", " ")
        ph = sh.get("placeholder")
        meta = f"kind={sh.get('kind')}" + (f" placeholder={ph}" if ph else "")
        lines.append(f"  #{sh['_i']}: rect={rs} {meta}" + (f' text="{snip}"' if snip else ""))
    lines += ["",
              "Return ONLY a JSON array, one object per shape, like:",
              '  [{"i": 0, "role": "title"}, {"i": 1, "role": "cards"}]',
              "Every shape index must appear exactly once. No prose, no code fence."]
    return "\n".join(lines)


def _extract_json_array(text):
    """Pull the first balanced top-level JSON array out of a model reply (which may wrap
    it in prose or a code fence). Returns the parsed list, or None."""
    if not text:
        return None
    start = text.find("[")
    while start != -1:
        depth = 0
        in_str = False
        esc = False
        for j in range(start, len(text)):
            c = text[j]
            if in_str:
                if esc:
                    esc = False
                elif c == "\\":
                    esc = True
                elif c == '"':
                    in_str = False
                continue
            if c == '"':
                in_str = True
            elif c == "[":
                depth += 1
            elif c == "]":
                depth -= 1
                if depth == 0:
                    try:
                        val = json.loads(text[start:j + 1])
                    except ValueError:
                        break
                    return val if isinstance(val, list) else None
        start = text.find("[", start + 1)
    return None


def parse_roletags(reply, shapes):
    """Map a model reply to a role per shape (aligned to `shapes`). Unmapped shapes,
    unknown roles, and shapes with no geometry (rect is None) fall back to 'skip'."""
    roles = ["skip"] * len(shapes)
    arr = _extract_json_array(reply) or []
    by_i = {}
    for obj in arr:
        if isinstance(obj, dict) and "i" in obj:
            try:
                by_i[int(obj["i"])] = str(obj.get("role", "skip")).strip().lower()
            except (TypeError, ValueError):
                continue
    for idx, sh in enumerate(shapes):
        if sh.get("rect") is None:           # can't place a shape with inherited geometry
            continue
        r = by_i.get(sh.get("_i", idx), "skip")
        roles[idx] = r if r in _ALLOWED_ROLES else "skip"
    return roles


def _bbox(rects):
    """Bounding box [x, y, w, h] over a list of rects."""
    xs = [r[0] for r in rects]
    ys = [r[1] for r in rects]
    x2 = max(r[0] + r[2] for r in rects)
    y2 = max(r[1] + r[3] for r in rects)
    x, y = min(xs), min(ys)
    return [round(x, 4), round(y, 4), round(x2 - x, 4), round(y2 - y, 4)]


def assemble_region_spec(name, shapes, roles, title=None, source=None,
                         category="default", brand=None):
    """Deterministically build a renderable region-spec from per-shape roles + the
    EXTRACTED geometry. Repeated 'cards' shapes merge into one cards region; a repeated
    singular role (two bodies, etc.) gets a numbered bind so binds stay unique.

    `category` tags the picker group: "default" (generic, renders under any brand) or a
    client name (client-specific). `brand` is the client-specific default brand."""
    regions = []
    card_rects = [shapes[i]["rect"] for i, r in enumerate(roles)
                  if r == "cards" and shapes[i].get("rect")]
    bind_used = {}

    def uniq(bind):
        bind_used[bind] = bind_used.get(bind, 0) + 1
        return bind if bind_used[bind] == 1 else f"{bind}{bind_used[bind]}"

    for sh, role in zip(shapes, roles):
        rect = sh.get("rect")
        if role == "skip" or rect is None:
            continue
        if role == "cards":
            continue                          # handled once, below
        if role == "background":
            regions.append({"role": "background", "rect": rect, "color": "primary"})
        elif role == "image":
            regions.append({"role": "image", "rect": rect,
                            "bind": uniq("image"), "fit": "contain"})
        else:
            reg = {"role": role, "rect": rect, "bind": uniq(_ROLE_BIND.get(role, "body"))}
            if role in ("title", "kicker", "value", "subtitle"):
                reg["align"] = "left"
            regions.append(reg)

    if card_rects:
        regions.append({"role": "cards", "rect": _bbox(card_rects), "bind": "cards",
                        "columns": len(card_rects),
                        "item": {"title": "title", "body": "body"}})

    # Order regions top-to-bottom, then left-to-right, for a tidy, readable spec.
    regions.sort(key=lambda r: (round(r["rect"][1], 2), round(r["rect"][0], 2)))
    spec = {"name": name,
            "title": title or name,
            "description": "Ingested layout"
            + (f" from {source}" if source else "") + " (structure only; brand supplies all color/type).",
            "regions": regions}
    if category and category != "default":
        spec["category"] = category
    if brand:
        spec["brand"] = brand
    return spec


def roletag_slide(profile, index=0, name=None, provider="claude", runner=None,
                  timeout=None, category="default", brand=None):
    """Tag one slide of an inspect_pptx() profile into a region-spec. Returns
    (spec, "") on success or (None, error). `runner` defaults to authoring.run_cli;
    injecting a stub keeps this testable without a live CLI. `category`/`brand` tag the
    spec for the picker (a client name + its default brand make it client-specific)."""
    slides = profile.get("slides") or []
    if not (0 <= index < len(slides)):
        return None, f"slide index {index} out of range (0..{len(slides) - 1})"
    raw = slides[index].get("shapes") or []
    shapes = [dict(sh, _i=i) for i, sh in enumerate(raw)]   # tag stable indices
    placeable = [sh for sh in shapes if sh.get("rect") is not None]
    if not placeable:
        return None, "slide has no placeable shapes (all geometry inherited)"
    if runner is None:
        import authoring
        runner = authoring.run_cli
    context = "; ".join(slides[index].get("texts") or [])[:200] or None
    prompt = build_roletag_prompt(placeable, context)
    ok, out, err = runner(provider, prompt, timeout=timeout)
    if not ok:
        return None, err or "role-tag CLI call failed"
    roles = parse_roletags(out, shapes)
    if all(r == "skip" for r in roles):
        return None, "model returned no usable roles"
    return assemble_region_spec(name or profile.get("name") or "ingested",
                                shapes, roles, source=profile.get("source"),
                                category=category, brand=brand), ""
