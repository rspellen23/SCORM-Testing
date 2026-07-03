"""Generic data-driven `template` layout (ingestable slide-types).

A layout is a JSON region-spec (templates/layouts/<name>.json) rendered with the
SELECTED BRAND's tokens + the client's images. Proves: it loads, renders the same
spec under two different brands (brand-agnostic structure), places a client image,
and degrades to a labeled placeholder when the template is missing.
"""
import os

import pytest

import slide_layouts
import authoring
import brand as brandmod

pptx = pytest.importorskip("pptx")  # export extra


def test_specs_listed_and_loaded():
    names = slide_layouts.list_template_layouts()
    assert "bignumber" in names and "three-cards" in names
    spec = slide_layouts.load_template_spec("bignumber")
    assert spec["name"] == "bignumber" and spec["regions"]
    assert slide_layouts.load_template_spec("nope") is None
    # basename-confined
    assert slide_layouts.load_template_spec("../../etc/passwd") is None


def test_template_lint_valid_but_not_in_generation_order():
    ok, n, errs = authoring.lint_deck([{"layout": "template", "content": {"template": "bignumber"}}])
    assert ok and not errs                              # passes lint
    assert "template" not in authoring._LAYOUT_ORDER    # AI never auto-picks it


def _pptx_texts(path):
    from pptx import Presentation
    p = Presentation(path)
    return " ".join(sh.text_frame.text for s in p.slides
                    for sh in s.shapes if sh.has_text_frame)


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_bignumber_renders_in_any_brand(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"big_{brand_name}.pptx")
    content = {"template": "bignumber", "kicker": "Throughput",
               "value": "6×", "title": "Faster bed assignment", "subtitle": "vs. manual"}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    assert "6×" in txt and "Faster bed assignment" in txt    # brand-agnostic structure


def test_three_cards_repeats_and_takes_client_image(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "cards.pptx")
    content = {"template": "three-cards", "title": "Three pillars",
               "image": "icon-building-hospital.svg",      # client/library image slot
               "cards": [{"title": "Sense", "body": "See demand"},
                         {"title": "Plan", "body": "Match supply"},
                         {"title": "Act", "body": "Move the patient"}]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    from pptx import Presentation
    p = Presentation(out)
    txt = " ".join(sh.text_frame.text for s in p.slides
                   for sh in s.shapes if sh.has_text_frame)
    for w in ("Three pillars", "Sense", "Plan", "Act", "Move the patient"):
        assert w in txt
    pics = [sh for s in p.slides for sh in s.shapes if sh.shape_type == 13]
    assert pics                                          # the icon placed via the image region


def test_missing_template_is_labeled_not_a_crash(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "missing.pptx")
    slide_layouts.export_deck(
        [{"layout": "template", "content": {"template": "does-not-exist"}}], out, brand=br)
    assert "Template not found" in _pptx_texts(out)


def test_template_renders_in_svg_preview():
    import slide_svg
    br = brandmod.load_brand("teletracking")
    svg = slide_svg.render_slide_svg(
        "template", {"template": "bignumber", "value": "42%", "title": "Adoption"}, br)
    assert svg.startswith("<svg") and "42%" in svg and "Adoption" in svg
