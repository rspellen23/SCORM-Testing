"""Template ingestion, step 1: the design-DNA extractor (src/pptx_ingest.py).

Reads a foreign .pptx into a profile (theme palette + fonts, slide size, layout
inventory, colors-in-use, per-slide shapes). Uses a generated fixture so the test is
self-contained (no checked-in binary / external file).
"""
import os
import json

import pytest

pptx = pytest.importorskip("pptx")  # the export extra
import pptx_ingest


def _fixture(path):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    prs = Presentation()                         # default Office template (has a theme)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Hello Template"
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x30)
    prs.slides.add_slide(blank)                  # a second, empty slide
    prs.save(path)


@pytest.fixture
def profile(tmp_path):
    f = tmp_path / "sample.pptx"
    _fixture(str(f))
    return pptx_ingest.inspect_pptx(str(f))


def test_identity_and_size(profile):
    assert profile["name"] == "sample" and profile["source"] == "sample.pptx"
    w, h = profile["slide_size_in"]
    assert w > h > 0 and profile["aspect"] > 1          # landscape


def test_theme_palette_and_fonts(profile):
    pal = profile["theme"]["palette"]
    assert "accent1" in pal and pal["accent1"].startswith("#") and len(pal["accent1"]) == 7
    assert {"dk1", "lt1"} <= set(pal)
    fonts = profile["theme"]["fonts"]
    assert fonts.get("major") and fonts.get("minor")    # Office theme has both


def test_layout_inventory(profile):
    assert profile["layout_count"] == len(profile["layouts"]) >= 7
    names = [l["name"] for l in profile["layouts"]]
    assert any("Blank" in n or "Title" in n for n in names)
    assert all("placeholders" in l for l in profile["layouts"])


def test_slides_and_shapes(profile):
    assert profile["slide_count"] == 2
    first = profile["slides"][0]
    assert first["index"] == 0 and "layout" in first and isinstance(first["shape_types"], dict)
    assert any("Hello Template" in t for t in first["texts"])


def test_colors_in_use_captures_explicit_fill(profile):
    assert "#102030" in profile["colors_in_use"]


def test_per_shape_geometry_relative_rects(profile):
    # Step 2b: every shape carries a region-spec-shaped relative rect [x, y, w, h].
    first = profile["slides"][0]
    assert "shapes" in first and len(first["shapes"]) >= 2
    for sh in first["shapes"]:
        assert set(sh) == {"kind", "rect", "placeholder", "text", "fill"}
        r = sh["rect"]
        assert r is None or (len(r) == 4 and all(0.0 <= v <= 1.0 for v in r))

    # The textbox at (1in, 1in) 4x1in on a 13.33x7.5in slide → known relative rect.
    w_in, h_in = profile["slide_size_in"]
    tb = next(s for s in first["shapes"] if "Hello Template" in (s["text"] or ""))
    x, y, w, h = tb["rect"]
    assert x == pytest.approx(1.0 / w_in, abs=0.01)
    assert y == pytest.approx(1.0 / h_in, abs=0.01)
    assert w == pytest.approx(4.0 / w_in, abs=0.01)
    assert h == pytest.approx(1.0 / h_in, abs=0.01)


def test_shape_carries_fill_hint(profile):
    # The filled rectangle's per-shape `fill` mirrors the colors_in_use roll-up.
    shapes = profile["slides"][0]["shapes"]
    assert any(s["fill"] == "#102030" for s in shapes)


def test_unreadable_file_raises(tmp_path):
    bad = tmp_path / "not.pptx"
    bad.write_text("not a pptx")
    with pytest.raises(Exception):
        pptx_ingest.inspect_pptx(str(bad))


def test_cli_writes_profile(tmp_path):
    import cli
    f = tmp_path / "deck.pptx"
    _fixture(str(f))
    out = tmp_path / "profile.json"

    class A:
        pptx = str(f)
        out = None
    A.out = str(out)
    cli.cmd_inspect_pptx(A)
    data = json.loads(out.read_text())
    assert data["slide_count"] == 2 and data["theme"]["palette"]
