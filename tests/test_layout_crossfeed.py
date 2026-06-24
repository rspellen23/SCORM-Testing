"""Audit item 2.7 — the same layout JSON must render correctly as BOTH a course
HTML block and a PPTX slide, in either schema dialect.

comparison: course used `panels[]`, the slide used `columns[]`.
timeline:   course read milestone `html`, the slide read `body`.

Cross-feeding one dialect to the other renderer used to silently produce an empty
block. The shared layouts.normalize_* adapters fix it; this test feeds each
dialect to both renderers and asserts the content actually lands.
"""
import os
import tempfile

import render
import slide_layouts
from pptx import Presentation


def _slide_text(content, layout):
    out = tempfile.mktemp(suffix=".pptx")
    try:
        slide_layouts.export_slide(content, out, layout=layout)
        prs = Presentation(out)
        chunks = []
        for sl in prs.slides:
            for sh in sl.shapes:
                if sh.has_text_frame:
                    chunks.append(sh.text_frame.text)
        return "\n".join(chunks)
    finally:
        if os.path.exists(out):
            os.unlink(out)


# ----------------------------------------------------------------- comparison

_CMP_PANELS = {"title": "Old vs New", "panels": [
    {"heading": "Before", "items": ["Manual step"], "callout": "slow"},
    {"heading": "After", "items": ["Automated"], "callout": "fast"},
]}
_CMP_COLUMNS = {"title": "Old vs New", "columns": _CMP_PANELS["panels"]}


def test_comparison_course_block_renders_both_dialects():
    for d in (_CMP_PANELS, _CMP_COLUMNS):
        html = render.render_block({"type": "comparison", **d})
        assert "Before" in html and "After" in html, (d, html)
        assert "Manual step" in html and "Automated" in html


def test_comparison_slide_renders_both_dialects():
    for d in (_CMP_PANELS, _CMP_COLUMNS):
        txt = {**d, "layout": "comparison"}
        out = _slide_text(d, "comparison")
        assert "Before" in out and "After" in out, out


# ------------------------------------------------------------------- timeline

_TL_HTML = {"title": "Roadmap", "milestones": [
    {"phase": "Q1", "title": "Pilot", "html": "<p>Start small</p>"},
    {"phase": "Q2", "title": "Scale", "html": "<p>Roll out wide</p>"},
]}
_TL_BODY = {"title": "Roadmap", "milestones": [
    {"phase": "Q1", "title": "Pilot", "body": "Start small"},
    {"phase": "Q2", "title": "Scale", "body": "Roll out wide"},
]}


def test_timeline_course_block_renders_both_dialects():
    for d in (_TL_HTML, _TL_BODY):
        html = render.render_block({"type": "timeline", **d})
        assert "Pilot" in html and "Scale" in html, html
        assert "Start small" in html and "Roll out wide" in html, html


def test_timeline_slide_renders_both_dialects():
    for d in (_TL_HTML, _TL_BODY):
        out = _slide_text(d, "timeline")
        # the course-dialect (html) milestone body must now appear on the slide
        assert "Start small" in out and "Roll out wide" in out, out
        assert "Pilot" in out and "Scale" in out, out
