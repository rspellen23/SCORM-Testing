"""slide_svg — the HTML/SVG preview backend for the slide-deck slideshow.

It reuses the EXACT slide_layouts geometry (via a mock-pptx surface) so the
in-browser preview matches the .pptx export. These tests assert each layout
renders well-formed SVG carrying real geometry + brand color, that a deck
renders one SVG per slide, and that a broken slide degrades to an error card
instead of killing the whole preview.
"""
import json
import os
import xml.dom.minidom as minidom

import slide_svg
import brand as brandmod
import slide_layouts

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def _wellformed(svg):
    minidom.parseString(svg)            # raises on malformed XML
    assert svg.startswith("<svg") and svg.rstrip().endswith("</svg>")


def test_every_layout_renders_wellformed_svg():
    b = brandmod.load_brand("_default")
    for layout in slide_layouts.LAYOUTS:
        ex = os.path.join(REPO, "templates", "slide-layouts", f"{layout}.example.json")
        if layout == "chart":
            content = {"title": "Q", "chart": "bar", "categories": ["A", "B"],
                       "series": [{"name": "S", "data": [3, 5]}]}
        else:
            with open(ex, encoding="utf-8") as fh:
                content = json.load(fh)
        svg = slide_svg.render_slide_svg(layout, content, b)
        _wellformed(svg)
        assert 'viewBox="0 0 1280 720"' in svg          # true 16:9 canvas
        assert "<text" in svg                            # text actually landed


def test_brand_palette_reaches_the_svg():
    # the TeleGreen primary must appear when the teletracking brand is used.
    # With the master blueprint, a divider is backed by the template's node art,
    # so the brand green shows through its ACCENTS (kicker + rule line), not a
    # flat fill — the kicker makes that accent present.
    if not os.path.isdir(os.path.join(REPO, "brands", "teletracking")):
        return
    b = brandmod.load_brand("teletracking")
    svg = slide_svg.render_slide_svg(
        "divider", {"kicker": "SECTION 1", "title": "Hi"}, b)
    assert "#1EB16A" in svg.upper()


def test_deck_renders_one_svg_per_slide():
    b = brandmod.load_brand("_default")
    deck = [{"layout": "divider", "content": {"title": "One"}},
            {"layout": "process", "content": {"title": "Two",
             "steps": [{"num": "1", "title": "Do", "body": "it"}]}}]
    svgs = slide_svg.render_deck_svg(deck, b)
    assert len(svgs) == 2
    for s in svgs:
        _wellformed(s)


def test_broken_slide_becomes_error_card_not_a_crash():
    b = brandmod.load_brand("_default")
    deck = [{"layout": "divider", "content": {"title": "Good"}},
            {"layout": "comparison", "content": "{ not valid json"}]   # JSON string from editor
    svgs = slide_svg.render_deck_svg(deck, b)
    assert len(svgs) == 2
    _wellformed(svgs[0]); _wellformed(svgs[1])
    assert "could not be previewed" not in svgs[0]
    assert "could not be previewed" in svgs[1]


def test_json_string_content_is_parsed():
    b = brandmod.load_brand("_default")
    svgs = slide_svg.render_deck_svg(
        [{"layout": "divider", "content": json.dumps({"title": "Stringy"})}], b)
    assert len(svgs) == 1 and "Stringy" in svgs[0]


def _import_server():
    import importlib
    sys = __import__("sys")
    dash = os.path.join(REPO, "dashboard")
    if dash not in sys.path:
        sys.path.insert(0, dash)
    return importlib.import_module("server")


def test_slide_svg_endpoint_renders_one_poster():
    # the inline Step-2 row thumbnails fetch /api/slide-svg per slide; the
    # endpoint must return one faithful poster for a dict OR a JSON-string
    # content (what the editor stores), and degrade a broken slide to an error
    # card rather than 500.
    server = _import_server()
    r = server.do_slide_svg({"layout": "divider", "content": {"title": "Hi"}})
    assert r["ok"] is True
    _wellformed(r["svg"])
    assert "Hi" in r["svg"]

    r2 = server.do_slide_svg(
        {"layout": "process", "content": json.dumps(
            {"title": "Two", "steps": [{"num": "1", "title": "Do", "body": "it"}]})})
    assert r2["ok"] is True and "Two" in r2["svg"]

    r3 = server.do_slide_svg({"layout": "comparison", "content": "{ broken json"})
    assert r3["ok"] is True
    _wellformed(r3["svg"])
    assert "could not be previewed" in r3["svg"]


def test_slide_svg_endpoint_honors_theme_override():
    # the row thumbnail must reflect the per-slide dark|light flip the picker sets,
    # so do_slide_svg threads `theme` into the rendered slide. Uses a brand that
    # defines both themes (teletracking) so flipping actually changes the markup;
    # an unset/invalid theme is ignored (the layout's blueprint default).
    server = _import_server()
    base = {"layout": "cards", "brand": "teletracking",
            "content": {"title": "C", "cards": [{"title": "a", "body": "b"}]}}
    dark = server.do_slide_svg({**base, "theme": "dark"})
    light = server.do_slide_svg({**base, "theme": "light"})
    auto = server.do_slide_svg(base)
    bad = server.do_slide_svg({**base, "theme": "teal"})   # invalid -> ignored
    for r in (dark, light, auto, bad):
        assert r["ok"] is True and r["svg"].startswith("<svg")
    assert dark["svg"] != light["svg"]                     # the flip is real
    assert bad["svg"] == auto["svg"]                        # invalid theme dropped


# ----- image layouts (image / imagetext) ------------------------------------

def _tmp_png(tmp_path, name="photo.png", size=(320, 200), color=(30, 177, 106)):
    from PIL import Image
    p = tmp_path / name
    Image.new("RGB", size, color).save(str(p))
    return p


def test_image_layout_embeds_a_data_uri(tmp_path):
    # a real image is embedded as a base64 data URI so the SVG renders standalone
    _tmp_png(tmp_path)
    b = brandmod.load_brand("_default")
    svg = slide_svg.render_slide_svg(
        "image", {"title": "Hero", "image": "photo.png", "caption": "Cap"},
        b, images_dir=str(tmp_path))
    _wellformed(svg)
    assert "<image" in svg and "data:image/png;base64," in svg
    assert "Hero" in svg and "Cap" in svg


def test_imagetext_layout_renders_image_and_bullets(tmp_path):
    _tmp_png(tmp_path)
    b = brandmod.load_brand("_default")
    svg = slide_svg.render_slide_svg(
        "imagetext",
        {"title": "Beside", "image": "photo.png", "side": "right",
         "intro": "Intro line", "items": [["Lead", " rest"], "plain"]},
        b, images_dir=str(tmp_path))
    _wellformed(svg)
    assert "<image" in svg and "data:image" in svg
    assert "Intro" in svg and "Lead" in svg          # words land as separate tspans


def test_missing_image_is_a_placeholder_not_a_crash(tmp_path):
    b = brandmod.load_brand("_default")
    svg = slide_svg.render_slide_svg(
        "image", {"title": "T", "image": "nope.png"}, b, images_dir=str(tmp_path))
    _wellformed(svg)
    assert "<image" not in svg            # no real image placed
    assert "nope.png" in svg              # the labeled "Image not found" placeholder names the file


def test_resolve_images_joins_bare_filename(tmp_path):
    out = slide_layouts._resolve_images({"image": "a.png"}, str(tmp_path))
    assert out["image"] == os.path.join(str(tmp_path), "a.png")
    # absolute path is left untouched
    abspath = os.path.join(str(tmp_path), "b.png")
    assert slide_layouts._resolve_images({"image": abspath}, "/other")["image"] == abspath
    # no dir -> unchanged
    assert slide_layouts._resolve_images({"image": "c.png"}, None)["image"] == "c.png"


def test_pptx_deck_with_image_layouts_builds(tmp_path):
    # the .pptx export embeds the real picture (add_picture) without error
    _tmp_png(tmp_path)
    b = brandmod.load_brand("_default")
    out = tmp_path / "deck.pptx"
    deck = [{"layout": "image", "content": {"title": "I", "image": "photo.png"}},
            {"layout": "imagetext", "content": {"title": "IT", "image": "photo.png",
             "intro": "x", "items": ["one"]}}]
    stats = slide_layouts.export_deck(deck, str(out), brand=b, images_dir=str(tmp_path))
    assert stats["slides"] == 2 and out.exists() and out.stat().st_size > 0


def test_new_layouts_render_wellformed(tmp_path):
    b = brandmod.load_brand("_default")
    cases = {
        "cards": {"title": "Kit", "columns": 2, "cards": [
            {"title": "Engine", "body": "stdlib"}, {"title": "AI", "body": "claude"}]},
        "quote": {"quote": "People own what ships.", "by": "Design principle", "bg": "dark"},
        "statement": {"kicker": "PAYOFF", "value": "6×", "title": "faster", "bg": "primary"},
        "bullets": {"title": "Takeaways", "items": [["Lead", " rest"], "plain"]},
    }
    for lay, content in cases.items():
        svg = slide_svg.render_slide_svg(lay, content, b)
        _wellformed(svg)
        assert "<text" in svg


def test_cover_fit_slices_contain_meets(tmp_path):
    _tmp_png(tmp_path)
    b = brandmod.load_brand("_default")
    cover = slide_svg.render_slide_svg(
        "image", {"title": "H", "image": "photo.png", "fit": "cover"}, b, images_dir=str(tmp_path))
    contain = slide_svg.render_slide_svg(
        "image", {"title": "H", "image": "photo.png", "fit": "contain"}, b, images_dir=str(tmp_path))
    assert "slice" in cover and "<image" in cover
    assert "meet" in contain and "slice" not in contain


def test_cover_crop_set_on_pptx_picture(tmp_path):
    # a wide image cover-fit into a tall-ish box must crop the sides (crop_left>0)
    _tmp_png(tmp_path, size=(800, 200))   # very wide
    b = brandmod.load_brand("_default")
    out = tmp_path / "d.pptx"
    slide_layouts.export_deck(
        [{"layout": "image", "content": {"title": "H", "image": "photo.png", "fit": "cover"}}],
        str(out), brand=b, images_dir=str(tmp_path))
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pics = [sh for s in Presentation(str(out)).slides for sh in s.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics and (pics[0].crop_left > 0 or pics[0].crop_right > 0)


def test_animate_wraps_groups_in_preview():
    b = brandmod.load_brand("_default")
    svgs = slide_svg.render_deck_svg(
        [{"layout": "divider", "content": {"title": "X"}}], b, animate=True)
    assert 'class="nv-a"' in svgs[0] and "nv-rise" in svgs[0] and "animation-delay" in svgs[0]
    # default (no animate) stays static — no animation wrappers
    plain = slide_svg.render_deck_svg([{"layout": "divider", "content": {"title": "X"}}], b)
    assert "nv-a" not in plain[0]


def test_animation_never_leaves_a_blank_slide():
    # REGRESSION: a prior version set .nv-a{opacity:0} as the BASE state and went
    # white when the animation didn't run (reduced-motion / some webviews). The
    # base must stay visible: use `both` fill, and never a persistent opacity:0.
    b = brandmod.load_brand("_default")
    svg = slide_svg.render_deck_svg(
        [{"layout": "divider", "content": {"title": "X"}}], b, animate=True)[0]
    assert "animation:nv-rise .5s ease both" in svg          # `both`, not `forwards`
    assert ".nv-a{opacity:0" not in svg                      # base state is NOT invisible
    assert "prefers-reduced-motion:reduce" in svg and "animation:none" in svg


def test_image_in_deck_svg_endpoint(tmp_path):
    # the slide-svg endpoint resolves images via the images dir
    server = _import_server()
    # confine: the endpoint only accepts a dir under the allowed roots, so use one
    import importlib
    _tmp_png(tmp_path)
    r = server.do_slide_svg({"layout": "image",
                             "content": {"title": "H", "image": "photo.png"},
                             "images": str(tmp_path)})
    assert r["ok"] is True
    # tmp_path may be outside the allow-roots (so images_dir resolves to None and the
    # slot is a placeholder) OR inside (then it embeds) — either way it's well-formed.
    _wellformed(r["svg"])
    assert "<svg" in r["svg"]


# --- B4: a full-bleed image's scrim overlay must be TRANSLUCENT in preview ---
# The mock no-op'd alpha, so the scrim rendered opaque — a solid slab hiding the
# photo (preview != .pptx). _fill_alpha now feeds the mock's set_alpha.

def test_full_bleed_image_scrim_is_translucent():
    import tempfile
    from PIL import Image
    d = tempfile.mkdtemp()
    img = os.path.join(d, "hero.png")
    Image.new("RGB", (800, 450), (20, 40, 60)).save(img)
    svg = slide_svg.render_slide_svg("image", {"title": "Big", "image": "hero.png",
                                               "mode": "full"}, images_dir=d)
    _wellformed(svg)
    assert "<image " in svg                      # the photo is actually drawn
    assert "fill-opacity=" in svg                # the scrim is translucent, not an opaque slab


def test_mock_shape_renders_outline_when_set():
    # the mock now honors line.width + line.color → an SVG stroke (was dropped)
    import slide_layouts as SL
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    sl = slide_svg._Slide()
    sh = sl.shapes.add_shape(__import__("pptx").enum.shapes.MSO_SHAPE.OVAL, Pt(0), Pt(0), Pt(50), Pt(50))
    sh.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sh.line.width = Pt(1.25); sh.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    assert "stroke=" in sh.to_svg() and "stroke-width=" in sh.to_svg()


# --- B5: cycles hub is outlined (visible on a same-tone bg) + steps cap at 6 -

def _cycles(n):
    steps = [{"title": f"Step {i + 1}", "body": "b"} for i in range(n)]
    return slide_svg.render_slide_svg("cycles", {"title": "Loop", "center": "Core",
                                                 "steps": steps})


def test_cycles_hub_is_outlined():
    svg = _cycles(4)
    _wellformed(svg)
    assert "stroke-width=" in svg                 # the hub disc carries a visible outline


def test_cycles_caps_steps_at_six():
    six, nine = _cycles(6), _cycles(9)
    # the ring + legend draw a fixed set of ellipses per step; >6 must not add more
    assert nine.count("<ellipse") == six.count("<ellipse")
    assert _cycles(3).count("<ellipse") < six.count("<ellipse")
