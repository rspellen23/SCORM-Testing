"""Q2 — auto speaker notes per slide (2026-06-29).

A deck slide may carry an optional "notes" string (sibling to "layout"/"theme").
It rides into the PowerPoint NOTES page, never onto the slide itself, and a
notes-free deck creates NO notes slide (so existing decks render unchanged).
A one-click action (`authoring.generate_notes` / `do_deck_notes`) drafts a notes
paragraph per slide for an existing deck.

Covers: the notes-write path (export_deck), the byte-identical no-notes guarantee,
notes riding only the FIRST paginated slide, the prompt injection (build_deck_prompt
+ build_deck_notes_prompt), generate_notes pad/truncate via a stub provider, lint
validation, and the server do_deck normalization + do_deck_notes endpoint.
"""
import os
import sys
import tempfile

import authoring
import brand as brandmod
import slide_layouts as SL

from pptx import Presentation

_DASH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "dashboard")
if _DASH not in sys.path:
    sys.path.insert(0, _DASH)
import server  # noqa: E402


def _build(slides):
    fd, tmp = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    SL.export_deck(slides, tmp, brand=brandmod.load_brand("teletracking"))
    return tmp


def _notes_texts(path):
    prs = Presentation(path)
    out = []
    for s in prs.slides:
        out.append(s.notes_slide.notes_text_frame.text if s.has_notes_slide else None)
    return out


# ---- export writes notes onto the notes page ----

def test_export_deck_writes_notes_to_notes_page():
    slides = [
        {"layout": "divider", "content": {"title": "Hello"}, "notes": "Open warmly; set the frame."},
        {"layout": "bullets", "content": {"title": "Points", "items": ["a", "b"]}, "notes": "Walk each point."},
    ]
    path = _build(slides)
    try:
        texts = _notes_texts(path)
        assert texts[0] == "Open warmly; set the frame."
        assert texts[1] == "Walk each point."
    finally:
        os.unlink(path)


def test_no_notes_creates_no_notes_slide():
    """A notes-free deck must NOT create notes slides (byte-identical guarantee)."""
    slides = [
        {"layout": "divider", "content": {"title": "Hello"}},
        {"layout": "bullets", "content": {"title": "Points", "items": ["a", "b"]}, "notes": "  "},
    ]
    path = _build(slides)
    try:
        assert _notes_texts(path) == [None, None]   # blank/whitespace notes => no notes slide
    finally:
        os.unlink(path)


# ---- pagination: notes ride only the FIRST page ----

def test_cont_carries_notes_only_on_first_page():
    spec = {"layout": "cards", "notes": "Spoken intro for this group."}
    content = {"title": "T", "cards": [{"title": "c"}]}
    first = SL._cont(spec, content, "cards", content["cards"], is_cont=False)
    cont = SL._cont(spec, content, "cards", content["cards"], is_cont=True)
    assert first.get("notes") == "Spoken intro for this group."
    assert "notes" not in cont


def test_paginated_deck_has_notes_on_one_slide_only():
    cards = [{"title": f"c{i}", "body": "x"} for i in range(14)]   # overflows cards capacity
    slides = [{"layout": "cards", "content": {"title": "Many", "cards": cards},
               "notes": "Cover the whole set."}]
    path = _build(slides)
    try:
        texts = _notes_texts(path)
        assert len(texts) > 1                                # it actually paginated
        assert texts[0] == "Cover the whole set."
        assert all(t is None for t in texts[1:])             # continuation pages carry no notes
    finally:
        os.unlink(path)


# ---- prompt injection ----

def test_deck_prompt_mentions_optional_notes():
    p = authoring.build_deck_prompt("T", "f", "aud", None, "SOURCE")
    assert '"notes"' in p and "notes page" in p


def test_deck_notes_prompt_covers_all_slides():
    slides = [{"layout": "divider", "content": {"title": "A"}},
              {"layout": "bullets", "content": {"title": "B", "items": ["x"]}},
              {"layout": "quote", "content": {"quote": "q"}}]
    p = authoring.build_deck_notes_prompt(slides, title="T", audience="ops")
    assert '"notes"' in p
    assert "EXACTLY 3" in p          # array length pinned to slide count
    assert "divider" in p and "quote" in p   # the model sees each slide


# ---- generate_notes pad / truncate via a stub provider ----

def test_generate_notes_pads_and_truncates(monkeypatch):
    slides = [{"layout": "divider", "content": {}}, {"layout": "bullets", "content": {}},
              {"layout": "quote", "content": {}}]

    # model returns too FEW notes -> pad with "" to slide count
    monkeypatch.setattr(authoring, "run_cli",
                        lambda *a, **k: (True, '{"notes": ["one"]}', ""))
    r = authoring.generate_notes("claude", slides)
    assert r["ok"] and r["notes"] == ["one", "", ""]

    # model returns too MANY -> truncate to slide count
    monkeypatch.setattr(authoring, "run_cli",
                        lambda *a, **k: (True, '{"notes": ["a", "b", "c", "d", "e"]}', ""))
    r = authoring.generate_notes("claude", slides)
    assert r["notes"] == ["a", "b", "c"]


def test_generate_notes_handles_bad_payload(monkeypatch):
    slides = [{"layout": "divider", "content": {}}]
    monkeypatch.setattr(authoring, "run_cli", lambda *a, **k: (True, "not json", ""))
    assert authoring.generate_notes("claude", slides)["ok"] is False
    monkeypatch.setattr(authoring, "run_cli", lambda *a, **k: (False, "", "provider down"))
    assert authoring.generate_notes("claude", slides)["ok"] is False
    assert authoring.generate_notes("claude", [])["ok"] is False


# ---- lint validation ----

def test_lint_deck_rejects_non_string_notes():
    ok, _, errs = authoring.lint_deck([{"layout": "divider", "content": {}, "notes": ["a"]}])
    assert not ok and any("notes must be text" in e for e in errs)


def test_lint_deck_allows_string_or_absent_notes():
    ok, _, _ = authoring.lint_deck([{"layout": "divider", "content": {}, "notes": "hi"},
                                    {"layout": "quote", "content": {"quote": "q"}}])
    assert ok


# ---- server normalization + endpoint ----

def test_do_deck_carries_notes_and_drops_empty():
    norm = {}

    def _cap(args):
        # intercept run_cli to read the temp deck JSON the server wrote
        import json
        ci = args.index("--content")
        with open(args[ci + 1], encoding="utf-8") as fh:
            norm["slides"] = json.load(fh)["slides"]
        return True, "ok"

    orig = server.run_cli
    server.run_cli = _cap
    try:
        server.do_deck({"slides": [
            {"layout": "divider", "content": {"title": "A"}, "notes": "say this"},
            {"layout": "quote", "content": {"quote": "q"}, "notes": "   "},   # blank -> dropped
            {"layout": "bullets", "content": {"items": ["x"]}},
        ], "out": tempfile.mkdtemp(), "name": "t"})
    finally:
        server.run_cli = orig
    s = norm["slides"]
    assert s[0].get("notes") == "say this"
    assert "notes" not in s[1] and "notes" not in s[2]


def test_do_deck_notes_endpoint(monkeypatch):
    monkeypatch.setattr(authoring, "run_cli",
                        lambda *a, **k: (True, '{"notes": ["spoken intro", ""]}', ""))
    r = server.do_deck_notes({"provider": "claude", "slides": [
        {"layout": "divider", "content": '{"title": "A"}'},   # content as JSON string from editor
        {"layout": "quote", "content": {"quote": "q"}},
    ]})
    assert r["ok"] and r["notes"] == ["spoken intro", ""]
    assert server.do_deck_notes({"slides": []})["ok"] is False
