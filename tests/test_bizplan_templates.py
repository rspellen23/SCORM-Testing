"""Business-plan template family (region-specs in templates/layouts/bizplan-*.json).

Our own original, brand-agnostic pitch-deck slide types (Step 3): the problem,
the solution (3 pillars), a 4-phase roadmap, and the ask. Proves each loads,
registers as a generic ("default") pickable template, renders the same spec
under two brands, repeats cards over their array, and renders in the SVG preview.
"""
import pytest

import slide_layouts
import slide_svg
import brand as brandmod

pptx = pytest.importorskip("pptx")  # export extra

BIZPLAN = ["bizplan-problem", "bizplan-solution", "bizplan-roadmap", "bizplan-ask"]


def _pptx_texts(path):
    from pptx import Presentation
    p = Presentation(path)
    return " ".join(sh.text_frame.text for s in p.slides
                    for sh in s.shapes if sh.has_text_frame)


def test_all_bizplan_specs_listed_and_generic():
    names = slide_layouts.list_template_layouts()
    info = {r["name"]: r for r in slide_layouts.template_layout_info()}
    for n in BIZPLAN:
        assert n in names, n
        spec = slide_layouts.load_template_spec(n)
        assert spec["name"] == n and spec["regions"]
        assert info[n]["category"] == "default" and info[n]["brand"] is None


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_problem_renders_bulleted_points(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"problem_{brand_name}.pptx")
    content = {"template": "bizplan-problem", "title": "The problem",
               "subtitle": "Bed assignment is slow and manual.",
               "points": ["Phone-tag handoffs", "No live visibility",
                          "Patients wait in the ED"]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("The problem", "manual", "Phone-tag handoffs", "Patients wait"):
        assert w in txt


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_solution_repeats_three_cards(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"solution_{brand_name}.pptx")
    content = {"template": "bizplan-solution", "title": "Our solution",
               "subtitle": "One platform, three moves.",
               "cards": [{"title": "Sense", "body": "See demand"},
                         {"title": "Plan", "body": "Match supply"},
                         {"title": "Act", "body": "Move the patient"}]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Our solution", "Sense", "Plan", "Act", "Move the patient"):
        assert w in txt


def test_roadmap_repeats_four_phases(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "roadmap.pptx")
    content = {"template": "bizplan-roadmap", "title": "Roadmap",
               "phases": [{"title": "Now", "body": "Pilot"},
                          {"title": "Q3", "body": "Roll out"},
                          {"title": "Q4", "body": "Scale"},
                          {"title": "2027", "body": "Expand"}]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Roadmap", "Pilot", "Roll out", "Scale", "Expand"):
        assert w in txt


def test_ask_renders_headline_figure(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "ask.pptx")
    content = {"template": "bizplan-ask", "kicker": "The ask",
               "value": "$2.5M", "title": "Seed round",
               "subtitle": "18-month runway to GA."}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("The ask", "$2.5M", "Seed round", "runway"):
        assert w in txt


@pytest.mark.parametrize("name,content", [
    ("bizplan-problem", {"title": "Problem", "subtitle": "It's slow",
                         "points": ["A", "B"]}),
    ("bizplan-solution", {"title": "Solution", "subtitle": "Three pillars",
                          "cards": [{"title": "One", "body": "x"},
                                    {"title": "Two", "body": "y"},
                                    {"title": "Three", "body": "z"}]}),
    ("bizplan-roadmap", {"title": "Roadmap",
                         "phases": [{"title": "Now", "body": "go"}]}),
    ("bizplan-ask", {"kicker": "Ask", "value": "$1M", "title": "Round",
                     "subtitle": "Runway"}),
])
def test_each_bizplan_spec_renders_in_svg_preview(name, content):
    br = brandmod.load_brand("teletracking")
    content = dict(content, template=name)
    svg = slide_svg.render_slide_svg("template", content, br)
    assert svg.startswith("<svg")
    assert content["title"].split()[0] in svg
