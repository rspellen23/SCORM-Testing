"""Template ingestion, step 2c: AI-assisted role tagging -> region-spec.

The AI only CLASSIFIES each shape into a role; geometry, binds, and spec assembly are
deterministic Python. These tests stub the CLI runner (no live `claude`), so they
verify the deterministic half end-to-end and prove the assembled spec actually renders
through slide_layouts._render_template.
"""
import json

import pytest

pptx = pytest.importorskip("pptx")  # the export extra
import pptx_ingest


# A small fake extracted slide (the shape shape that inspect_pptx emits).
def _shapes():
    return [
        {"kind": "TEXT_BOX (17)", "rect": [0.06, 0.09, 0.70, 0.12],
         "placeholder": None, "text": "How it works", "fill": None},
        {"kind": "AUTO_SHAPE (1)", "rect": [0.06, 0.30, 0.27, 0.50],
         "placeholder": None, "text": "Step one", "fill": "#EEEEEE"},
        {"kind": "AUTO_SHAPE (1)", "rect": [0.36, 0.30, 0.27, 0.50],
         "placeholder": None, "text": "Step two", "fill": "#EEEEEE"},
        {"kind": "AUTO_SHAPE (1)", "rect": [0.66, 0.30, 0.27, 0.50],
         "placeholder": None, "text": "Step three", "fill": "#EEEEEE"},
        {"kind": "PICTURE (13)", "rect": [0.80, 0.06, 0.14, 0.14],
         "placeholder": None, "text": "", "fill": None},
        {"kind": "AUTO_SHAPE (1)", "rect": None,            # inherited geometry
         "placeholder": "BODY (2)", "text": "decorative", "fill": None},
    ]


def _profile():
    return {"name": "demo", "source": "demo.pptx",
            "slides": [{"index": 0, "texts": ["How it works"], "shapes": _shapes()}]}


def test_prompt_lists_shapes_and_role_vocab():
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes()) if s["rect"]]
    p = pptx_ingest.build_roletag_prompt(shapes)
    for role in ("title", "cards", "image", "skip"):
        assert role in p
    assert "#0:" in p and "How it works" in p
    assert "JSON array" in p


def test_extract_json_array_handles_fences_and_prose():
    assert pptx_ingest._extract_json_array('[{"i": 0, "role": "title"}]') == \
        [{"i": 0, "role": "title"}]
    fenced = 'Sure!\n```json\n[{"i": 1, "role": "cards"}]\n```\nDone.'
    assert pptx_ingest._extract_json_array(fenced) == [{"i": 1, "role": "cards"}]
    assert pptx_ingest._extract_json_array("no json here") is None


def test_parse_roletags_aligns_and_defaults_skip():
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes())]
    reply = json.dumps([{"i": 0, "role": "title"}, {"i": 1, "role": "cards"},
                        {"i": 2, "role": "cards"}, {"i": 3, "role": "cards"},
                        {"i": 4, "role": "image"}, {"i": 99, "role": "bogus"}])
    roles = pptx_ingest.parse_roletags(reply, shapes)
    assert roles[0] == "title" and roles[4] == "image"
    assert roles[1] == roles[2] == roles[3] == "cards"
    assert roles[5] == "skip"                       # rect is None -> always skip


def test_unknown_role_falls_back_to_skip():
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes())]
    roles = pptx_ingest.parse_roletags('[{"i": 0, "role": "wizardry"}]', shapes)
    assert roles[0] == "skip"


def test_assemble_merges_cards_and_keeps_binds_unique():
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes())]
    roles = ["title", "cards", "cards", "cards", "image", "skip"]
    spec = pptx_ingest.assemble_region_spec("howitworks", shapes, roles, source="demo.pptx")
    by_role = {}
    for r in spec["regions"]:
        by_role.setdefault(r["role"], []).append(r)
    assert len(by_role["cards"]) == 1                       # three card shapes -> one region
    cards = by_role["cards"][0]
    assert cards["columns"] == 3 and cards["bind"] == "cards"
    # bounding box spans all three cards: x≈0.06, w≈0.87
    assert cards["rect"][0] == pytest.approx(0.06, abs=0.01)
    assert cards["rect"][2] == pytest.approx(0.87, abs=0.02)
    assert by_role["title"][0]["bind"] == "title"
    assert by_role["image"][0]["fit"] == "contain"
    binds = [r["bind"] for r in spec["regions"] if "bind" in r]
    assert len(binds) == len(set(binds))                   # all binds unique


def test_repeated_singular_role_gets_numbered_bind():
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes()[:3])]
    roles = ["body", "body", "body"]
    spec = pptx_ingest.assemble_region_spec("x", shapes, roles)
    binds = sorted(r["bind"] for r in spec["regions"])
    assert binds == ["body", "body2", "body3"]


def test_assembled_spec_renders_through_template_renderer():
    """The deterministic spec must be renderable by slide_layouts._render_template."""
    sl = pytest.importorskip("slide_layouts")
    import brand
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes())]
    roles = ["title", "cards", "cards", "cards", "image", "skip"]
    spec = pptx_ingest.assemble_region_spec("rt_demo", shapes, roles)
    # write into the live template dir so load_template_spec finds it
    import os
    tdir = sl._TEMPLATE_DIR
    os.makedirs(tdir, exist_ok=True)
    path = os.path.join(tdir, "rt_demo.json")
    try:
        with open(path, "w", encoding="utf-8") as fh:
            json.dump(spec, fh)
        sl._TEMPLATE_CACHE.pop("rt_demo", None)
        br = brand.load_brand("teletracking")
        deck = [{"layout": "template",
                 "content": {"template": "rt_demo", "title": "Hi",
                             "cards": [{"title": "A", "body": "b"},
                                       {"title": "B", "body": "b"},
                                       {"title": "C", "body": "b"}]}}]
        outp = "/tmp/_rt_demo_deck.pptx"
        sl.export_deck(deck, outp, brand=br)
        assert os.path.getsize(outp) > 0
    finally:
        if os.path.exists(path):
            os.remove(path)
        sl._TEMPLATE_CACHE.pop("rt_demo", None)


def test_roletag_slide_with_stub_runner():
    """End-to-end through roletag_slide with an injected runner (no live CLI)."""
    def stub(provider, prompt, timeout=None):
        assert provider == "claude" and "JSON array" in prompt
        return True, json.dumps([{"i": 0, "role": "title"}, {"i": 1, "role": "cards"},
                                 {"i": 2, "role": "cards"}, {"i": 3, "role": "cards"},
                                 {"i": 4, "role": "image"}]), ""
    spec, err = pptx_ingest.roletag_slide(_profile(), index=0, name="hiw", runner=stub)
    assert err == "" and spec["name"] == "hiw"
    roles = {r["role"] for r in spec["regions"]}
    assert {"title", "cards", "image"} <= roles


def test_roletag_slide_runner_failure_returns_error():
    def stub(provider, prompt, timeout=None):
        return False, "", "claude not installed"
    spec, err = pptx_ingest.roletag_slide(_profile(), runner=stub)
    assert spec is None and "claude not installed" in err


def test_roletag_slide_all_skip_is_error():
    def stub(provider, prompt, timeout=None):
        return True, "[]", ""
    spec, err = pptx_ingest.roletag_slide(_profile(), runner=stub)
    assert spec is None and "no usable roles" in err


def test_assemble_tags_category_and_brand():
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes()[:1])]
    spec = pptx_ingest.assemble_region_spec("tt", shapes, ["title"],
                                            category="TeleTracking", brand="teletracking")
    assert spec["category"] == "TeleTracking" and spec["brand"] == "teletracking"
    # a generic spec carries neither key
    generic = pptx_ingest.assemble_region_spec("g", shapes, ["title"])
    assert "category" not in generic and "brand" not in generic


def test_template_info_and_starter():
    sl = pytest.importorskip("slide_layouts")
    shapes = [dict(s, _i=i) for i, s in enumerate(_shapes())]
    roles = ["title", "cards", "cards", "cards", "image", "skip"]
    spec = pptx_ingest.assemble_region_spec("ti_demo", shapes, roles,
                                            category="TeleTracking", brand="teletracking")
    starter = sl.build_template_starter(spec)
    assert starter["template"] == "ti_demo"
    assert starter["title"] == "Title"
    assert isinstance(starter["cards"], list) and len(starter["cards"]) == 3
    assert "image" not in starter                          # image slot stays empty

    import os, json
    tdir = sl._TEMPLATE_DIR
    os.makedirs(tdir, exist_ok=True)
    path = os.path.join(tdir, "ti_demo.json")
    try:
        with open(path, "w", encoding="utf-8") as fh:
            json.dump(spec, fh)
        sl._TEMPLATE_CACHE.pop("ti_demo", None)
        info = {t["name"]: t for t in sl.template_layout_info()}
        assert "ti_demo" in info
        rec = info["ti_demo"]
        assert rec["category"] == "TeleTracking" and rec["brand"] == "teletracking"
        assert rec["starter"]["template"] == "ti_demo"
        # the shipped generic specs default to category "default", brand None
        if "bignumber" in info:
            assert info["bignumber"]["category"] == "default"
            assert info["bignumber"]["brand"] is None
    finally:
        if os.path.exists(path):
            os.remove(path)
        sl._TEMPLATE_CACHE.pop("ti_demo", None)


def test_cli_roletag_tags_client(tmp_path, monkeypatch):
    import cli
    from pptx import Presentation
    from pptx.util import Inches
    import json as _json
    f = tmp_path / "tt.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tb.text_frame.text = "Heading"
    prs.save(str(f))

    def stub(provider, prompt, timeout=None):
        return True, '[{"i": 0, "role": "title"}]', ""
    monkeypatch.setattr("authoring.run_cli", stub)
    outp = str(tmp_path / "tt.json")

    class A:
        pptx = str(f); slide = 0; name = "ttc"; provider = "claude"
        client = "TeleTracking"; brand = "teletracking"; out = outp
    cli.cmd_roletag_pptx(A)
    spec = _json.loads(open(outp, encoding="utf-8").read())
    assert spec["category"] == "TeleTracking" and spec["brand"] == "teletracking"


def test_cli_roletag_writes_spec(tmp_path, monkeypatch):
    import cli
    from pptx import Presentation
    from pptx.util import Inches
    f = tmp_path / "deck.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tb.text_frame.text = "Heading"
    prs.save(str(f))

    def stub(provider, prompt, timeout=None):
        return True, '[{"i": 0, "role": "title"}]', ""
    monkeypatch.setattr("authoring.run_cli", stub)
    outp = str(tmp_path / "spec.json")

    class A:
        pptx = str(f); slide = 0; name = "clitest"; provider = "claude"
        client = None; brand = None; out = outp
    cli.cmd_roletag_pptx(A)
    spec = json.loads(open(outp, encoding="utf-8").read())
    assert spec["name"] == "clitest" and spec["regions"][0]["role"] == "title"
