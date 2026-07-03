"""Per-region reverse-text `color` override on template text roles.

A region-spec text role (title/kicker/value/subtitle/body) may carry an optional
`color` palette token (e.g. "white") so it can sit in reversed white on a
full brand-color `background` fill. Absent, the role keeps its default color and
the render is byte-identical to before. This is the renderer half of the bold
section divider; it lives in the SHARED `_render_template`, so both the .pptx
export and the slide_svg preview honor it from one code path.
"""
import re

import pytest

import slide_layouts
import slide_svg
import brand as brandmod

pytest.importorskip("pptx")  # export extra


def _spec_to_temp(tmp_path, monkeypatch, spec):
    """Point the template loader at a temp dir holding one ad-hoc spec."""
    import json
    d = tmp_path / "layouts"
    d.mkdir()
    (d / (spec["name"] + ".json")).write_text(json.dumps(spec), encoding="utf-8")
    monkeypatch.setattr(slide_layouts, "_TEMPLATE_DIR", str(d))
    monkeypatch.setattr(slide_layouts, "_TEMPLATE_CACHE", {})


def _run_colors(path):
    from pptx import Presentation
    return {str(r.font.color.rgb)
            for sh in Presentation(path).slides[0].shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs for r in para.runs if r.text.strip()}


def test_color_token_reverses_text_on_fill(tmp_path, monkeypatch):
    br = brandmod.load_brand("teletracking")
    pal = slide_layouts._palette_of(br)
    _spec_to_temp(tmp_path, monkeypatch, {
        "name": "rev", "title": "rev", "regions": [
            {"role": "background", "rect": [0, 0, 1, 1], "color": "primary"},
            {"role": "title", "rect": [0.1, 0.4, 0.8, 0.2], "bind": "title", "color": "white"},
            {"role": "body", "rect": [0.1, 0.6, 0.8, 0.2], "bind": "body", "color": "white"},
        ]})
    out = str(tmp_path / "rev.pptx")
    slide_layouts.export_deck(
        [{"layout": "template", "content": {"template": "rev",
          "title": "Reversed", "body": ["Point one", "Point two"]}}], out, brand=br)
    assert _run_colors(out) == {str(pal["white"])}


def test_absent_color_keeps_default_ink(tmp_path, monkeypatch):
    br = brandmod.load_brand("teletracking")
    pal = slide_layouts._palette_of(br)
    _spec_to_temp(tmp_path, monkeypatch, {
        "name": "plain", "title": "plain", "regions": [
            {"role": "title", "rect": [0.1, 0.2, 0.8, 0.2], "bind": "title"},
            {"role": "body", "rect": [0.1, 0.5, 0.8, 0.3], "bind": "body"},
        ]})
    out = str(tmp_path / "plain.pptx")
    slide_layouts.export_deck(
        [{"layout": "template", "content": {"template": "plain",
          "title": "Plain", "body": ["a", "b"]}}], out, brand=br)
    ink = str(pal.get("ink", pal["dark"]))
    assert _run_colors(out) == {ink}


def test_unknown_color_token_falls_back_to_default(tmp_path, monkeypatch):
    br = brandmod.load_brand("teletracking")
    pal = slide_layouts._palette_of(br)
    _spec_to_temp(tmp_path, monkeypatch, {
        "name": "bad", "title": "bad", "regions": [
            {"role": "title", "rect": [0.1, 0.2, 0.8, 0.2], "bind": "title",
             "color": "not-a-token"},
        ]})
    out = str(tmp_path / "bad.pptx")
    slide_layouts.export_deck(
        [{"layout": "template", "content": {"template": "bad", "title": "X"}}],
        out, brand=br)
    assert _run_colors(out) == {str(pal.get("ink", pal["dark"]))}


def test_shipped_section_divider_reverses_in_svg_preview():
    # the upgraded workshop-section: white text on a brand fill, in the preview
    br = brandmod.load_brand("teletracking")
    svg = slide_svg.render_slide_svg(
        "template", {"template": "workshop-section",
                     "kicker": "Section 1", "title": "Getting started"}, br)
    assert svg.startswith("<svg")
    fills = set(re.findall(r'fill="(#[0-9a-fA-F]{6})"', svg))
    assert "#FFFFFF" in fills        # reversed white text
    assert "#1EB16A" in fills        # brand-color fill block
