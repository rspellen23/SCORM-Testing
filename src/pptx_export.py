"""Course IR -> PowerPoint (.pptx).

A SECOND output target alongside SCORM/cmi5. It reads the SAME Course IR and
renders a static, editable deck — one slide per source slide (each `heading`
block starts a new slide; the blocks after it are that slide's body).

This is deliberately a *static approximation*. The IR was built for responsive,
interactive, tracked HTML; PowerPoint is a fixed canvas with no runtime. So:

  rendered (flattened): heading, paragraph, list, note, statement, image,
                        imageText, table, cardGrid, knowledgeCheck (shown as a
                        static Q&A slide with the correct option marked).
  dropped  (logged):    continue gates, video/audio/embed, accordion/process/
                        flashcard/categorize/scenario, section wrappers, modals.

Build notes / author meta never reach the IR, so nothing to strip here.
"""
import os
from html.parser import HTMLParser

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 16:9 widescreen canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
TITLE_H = Inches(1.0)
BODY_TOP = Inches(1.35)
BODY_H = SLIDE_H - BODY_TOP - Inches(0.4)
BODY_W = SLIDE_W - MARGIN - MARGIN

NAVY = RGBColor(0x00, 0x3E, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF4, 0xF5)

# block types this exporter cannot represent on a static slide
# (timeline/comparison are rich HTML-player layouts — for a PowerPoint of those,
#  use the dedicated `./build slide --layout timeline|comparison` generator; here
#  they are logged as dropped rather than silently swallowed by the textflow.)
#  chart blocks likewise: the dedicated `./build slide --layout chart` generator emits a
#  NATIVE, editable PowerPoint chart — far better than a flattened bitmap — so here a chart
#  is logged as dropped. (Future: emit a native chart slide inline during the flatten.)
#  quote + infographic have no faithful static-slide form either (a tinted full-bleed
#  pull-quote / a poster section) → dropped-and-logged, not silently swallowed.
# This set is mirrored by blocks.BLOCKS (pptx="drop"/"structural"); the registry
# drift test (tests/test_block_registry.py) fails if the two disagree.
_DROP = {"continue", "video", "audio", "embed", "accordion", "process",
         "flashcard", "categorize", "scenario", "divider", "transition",
         "sectionStart", "sectionEnd", "button", "timeline", "comparison", "chart",
         "quote", "infographic"}


# ---------------------------------------------------------------- HTML -> runs

class _Frag(HTMLParser):
    """Turn an inline-HTML fragment into a list of paragraphs, each a list of
    runs: {text, bold, italic, code}. <p>/<br>/<li> split paragraphs."""
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.paras = [[]]
        self.b = self.i = self.code = 0

    def _newpara(self):
        if self.paras[-1]:
            self.paras.append([])

    def handle_starttag(self, tag, attrs):
        if tag in ("strong", "b"):
            self.b += 1
        elif tag in ("em", "i"):
            self.i += 1
        elif tag == "code":
            self.code += 1
        elif tag in ("p", "li", "br", "tr"):
            self._newpara()

    def handle_endtag(self, tag):
        if tag in ("strong", "b"):
            self.b = max(0, self.b - 1)
        elif tag in ("em", "i"):
            self.i = max(0, self.i - 1)
        elif tag == "code":
            self.code = max(0, self.code - 1)
        elif tag in ("p", "li"):
            self._newpara()

    def handle_data(self, data):
        if not data:
            return
        text = data if data.strip() else data
        self.paras[-1].append({"text": text, "bold": bool(self.b),
                               "italic": bool(self.i), "code": bool(self.code)})


def html_paras(fragment):
    """-> list of paragraphs (each a list of runs). Empty paragraphs dropped."""
    p = _Frag()
    p.feed(fragment or "")
    out = []
    for para in p.paras:
        merged = [r for r in para if r["text"].strip()]
        if merged:
            out.append(merged)
    return out


def plain(fragment):
    return " ".join("".join(r["text"] for r in para).strip()
                    for para in html_paras(fragment)).strip()


class _Table(HTMLParser):
    """Parse the simple <table><thead><tr><th>..><tbody><tr><td>..> the IR emits."""
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.rows, self.cur, self.buf, self.incell = [], None, [], False

    def handle_starttag(self, tag, attrs):
        if tag == "tr":
            self.cur = []
        elif tag in ("td", "th"):
            self.incell, self.buf = True, []

    def handle_endtag(self, tag):
        if tag in ("td", "th"):
            self.cur.append("".join(self.buf).strip())
            self.incell = False
        elif tag == "tr" and self.cur is not None:
            self.rows.append(self.cur)
            self.cur = None

    def handle_data(self, data):
        if self.incell:
            self.buf.append(data)


def parse_table(fragment):
    t = _Table()
    t.feed(fragment or "")
    return [r for r in t.rows if r]


# ----------------------------------------------------------------- IR -> slides

def group_slides(ir):
    """Split the flat block list into slides at each `heading` (level<=2).
    Returns [{"title": str, "blocks": [...]}], plus a list of dropped types."""
    slides, cur, dropped = [], None, []
    for b in ir.get("blocks", []):
        t = b.get("type")
        if t == "heading" and b.get("level", 2) <= 2:
            cur = {"title": plain(b.get("html")), "blocks": []}
            slides.append(cur)
            continue
        if t in _DROP:
            if t not in ("sectionStart", "sectionEnd", "divider", "transition"):
                dropped.append(t)
            continue
        if cur is None:                       # body before the first heading
            cur = {"title": ir.get("title", ""), "blocks": []}
            slides.append(cur)
        cur["blocks"].append(b)
    return slides, dropped


# ----------------------------------------------------------------- pptx helpers

def _img_fit(path, box_w, box_h):
    """Largest (w,h) in EMU fitting box while preserving the image aspect."""
    try:
        from PIL import Image
        w, h = Image.open(path).size
    except Exception:
        return box_w, box_h
    aspect = w / h if h else 1.0
    if box_w / box_h > aspect:                # box wider than image -> height-bound
        return int(box_h * aspect), int(box_h)
    return int(box_w), int(box_w / aspect)


def _add_runs(para, runs, size, color=GREY, font="Open Sans"):
    for rn in runs:
        r = para.add_run()
        r.text = rn["text"]
        r.font.size = Pt(size)
        r.font.bold = rn["bold"]
        r.font.italic = rn["italic"]
        r.font.color.rgb = color
        r.font.name = "Consolas" if rn["code"] else font


def _new_para(tf, first_holder):
    """Reuse the textframe's initial empty paragraph once, then add new ones."""
    if first_holder[0]:
        first_holder[0] = False
        return tf.paragraphs[0]
    return tf.add_paragraph()


def _render_textflow(tf, blocks, accent, unmatched=None):
    """Render the non-media, non-table blocks into one text frame.

    Any block type this flow doesn't recognize is recorded in `unmatched`
    (a set) rather than silently vanishing — a NEW block type added without a
    flatten disposition surfaces as a warning instead of disappearing."""
    tf.word_wrap = True
    first = [True]
    for b in blocks:
        t = b.get("type")
        if t == "heading":                                  # level-3 subheading
            for para in html_paras(b.get("html")):
                p = _new_para(tf, first); p.space_before = Pt(6)
                _add_runs(p, para, 16, NAVY)
                for r in p.runs:
                    r.font.bold = True
        elif t in ("paragraph", "headingParagraph"):
            if t == "headingParagraph":
                for para in html_paras(b.get("headingHtml")):
                    p = _new_para(tf, first); p.space_before = Pt(6)
                    _add_runs(p, para, 16, NAVY)
                    for r in p.runs:
                        r.font.bold = True
            for para in html_paras(b.get("html")):
                p = _new_para(tf, first); p.space_after = Pt(6)
                _add_runs(p, para, 14)
        elif t == "list":
            for item in b.get("items", []):
                paras = html_paras(item) or [[]]
                p = _new_para(tf, first); p.level = 0
                pre = p.add_run(); pre.text = "•  "
                pre.font.size = Pt(14); pre.font.color.rgb = accent
                pre.font.bold = True; pre.font.name = "Open Sans"
                _add_runs(p, paras[0], 14)
        elif t == "note":
            for para in html_paras(b.get("html")):
                p = _new_para(tf, first); p.space_before = Pt(6)
                tag = p.add_run(); tag.text = "Note  "
                tag.font.size = Pt(13); tag.font.bold = True
                tag.font.color.rgb = accent; tag.font.name = "Open Sans"
                _add_runs(p, para, 13, GREY)
                for r in p.runs:
                    r.font.italic = True
        elif t == "statement":
            for para in html_paras(b.get("html")):
                p = _new_para(tf, first); p.space_before = Pt(10); p.space_after = Pt(10)
                _add_runs(p, para, 20, accent)
                for r in p.runs:
                    r.font.bold = True
        elif t in ("image", "imageText"):
            # the picture is placed separately; here we keep its caption/body text
            body = b.get("html") or b.get("caption")
            for para in html_paras(body):
                p = _new_para(tf, first); p.space_after = Pt(6)
                _add_runs(p, para, 14)
        elif t == "cardGrid":
            for c in b.get("cards", []):
                p = _new_para(tf, first); p.space_before = Pt(4)
                head = p.add_run(); head.text = (c.get("title") or "")
                head.font.size = Pt(14); head.font.bold = True
                head.font.color.rgb = NAVY; head.font.name = "Open Sans"
                if c.get("teaser"):
                    tz = p.add_run(); tz.text = " — " + c["teaser"]
                    tz.font.size = Pt(13); tz.font.color.rgb = GREY
                    tz.font.name = "Open Sans"
        elif t == "knowledgeCheck":
            _render_kc(tf, b, accent, first)
        elif unmatched is not None:
            # not rendered here, not pre-pulled (image/table), not in _DROP -> a
            # type with no flatten disposition. Record it so it's reported, never
            # silently swallowed (the audit's silent-fallthrough finding).
            unmatched.add(t)


def _render_kc(tf, b, accent, first):
    p = _new_para(tf, first); p.space_after = Pt(6)
    lab = p.add_run(); lab.text = "Knowledge Check"
    lab.font.size = Pt(12); lab.font.bold = True
    lab.font.color.rgb = accent; lab.font.name = "Open Sans"
    for para in html_paras(b.get("prompt")):
        p = _new_para(tf, first)
        _add_runs(p, para, 15, NAVY)
        for r in p.runs:
            r.font.bold = True
    for o in b.get("options", []):
        p = _new_para(tf, first); p.level = 0
        mark = p.add_run()
        mark.text = ("✓  " if o.get("correct") else "○  ")
        mark.font.size = Pt(13); mark.font.bold = True
        mark.font.name = "Open Sans"
        mark.font.color.rgb = accent if o.get("correct") else GREY
        _add_runs(p, (html_paras(o.get("html")) or [[]])[0], 13,
                  NAVY if o.get("correct") else GREY)
    if b.get("feedback"):
        p = _new_para(tf, first); p.space_before = Pt(6)
        fb = p.add_run(); fb.text = b["feedback"]
        fb.font.size = Pt(12); fb.font.italic = True
        fb.font.color.rgb = GREY; fb.font.name = "Open Sans"


def _resolve(src, blobs):
    """Map a block 'assets/<file>' src to an on-disk path via the blob map."""
    if not src:
        return None
    if src in blobs and os.path.exists(blobs[src]):
        return blobs[src]
    return None


# --------------------------------------------------------------------- builder

def _title_bar(slide, text, accent):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, TITLE_H)  # 1 = rectangle
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = MARGIN; tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text or ""
    r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Open Sans Condensed"


def _add_table(slide, rows, left, top, width, accent):
    if not rows:
        return top
    ncol = max(len(r) for r in rows)
    nrow = len(rows)
    height = Inches(0.4 * nrow)
    gtbl = slide.shapes.add_table(nrow, ncol, left, top, width, height)
    tbl = gtbl.table
    for ci in range(ncol):
        for ri, row in enumerate(rows):
            cell = tbl.cell(ri, ci)
            cell.text = row[ci] if ci < len(row) else ""
            para = cell.text_frame.paragraphs[0]
            for r in para.runs:
                r.font.size = Pt(12); r.font.name = "Open Sans"
                if ri == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
                else:
                    r.font.color.rgb = GREY
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = accent
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return top + height


def _build_slide(prs, slide, accent, logo_path):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _title_bar(s, slide["title"], accent)
    if logo_path:
        try:
            lw = Inches(1.3)
            from PIL import Image
            w, h = Image.open(logo_path).size
            lh = int(lw * h / w)
            s.shapes.add_picture(logo_path, SLIDE_W - lw - Inches(0.3),
                                 Inches(0.5) - lh // 2 + Emu(int(TITLE_H) // 2) - Emu(int(Inches(0.5))),
                                 width=lw)
        except Exception:
            pass
    return s


def export_pptx(ir, blobs, out_path, brand=None, logo=True,
                transition=None, transition_dir=None, transition_speed="med"):
    """Render the Course IR to a .pptx at out_path. Returns a stats dict.

    If `transition` is set (e.g. "fade"/"push"), it is applied to every slide
    (see src/pptx_transitions.py for supported effects)."""
    accent = RGBColor.from_string((brand.accent if brand else "#3B82F6").lstrip("#"))
    logo_path = brand.asset(brand.get("logo", "Logo.png")) if (brand and logo) else None

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- title slide
    title = prs.slides.add_slide(prs.slide_layouts[6])
    band = title.shapes.add_shape(1, 0, Inches(2.4), SLIDE_W, Inches(2.7))
    band.fill.solid(); band.fill.fore_color.rgb = accent; band.line.fill.background()
    tf = band.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.8)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = ir.get("title", "Course")
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Open Sans Condensed"
    if logo_path:
        try:
            from PIL import Image
            w, h = Image.open(logo_path).size
            lw = Inches(2.0); lh = int(lw * h / w)
            title.shapes.add_picture(logo_path, Inches(0.8), Inches(0.7), width=lw)
        except Exception:
            pass

    slides, dropped = group_slides(ir)
    stats = {"slides": len(slides) + 1, "dropped": {}}
    for d in dropped:
        stats["dropped"][d] = stats["dropped"].get(d, 0) + 1
    unmatched = set()

    for sl in slides:
        s = _build_slide(prs, sl, accent, logo_path)
        blocks = sl["blocks"]
        # pick the one image to place visually (first resolvable image/imageText)
        media = next((b for b in blocks if b.get("type") in ("image", "imageText")
                      and _resolve(b.get("src"), blobs)), None)
        tables = [b for b in blocks if b.get("type") == "table"]
        # text flow = everything except the placed image and tables; an imageText's
        # own copy stays in the flow (rendered as text, its picture placed below).
        text_content = [b for b in blocks if b.get("type") not in ("table", "image")]

        img_path = _resolve(media.get("src"), blobs) if media else None
        img_left = media.get("side") == "left" if media else False

        if img_path and not tables:
            # two columns: image one side, text the other
            col_w = (BODY_W - Inches(0.4)) // 2
            iw, ih = _img_fit(img_path, col_w, BODY_H)
            if img_left:
                s.shapes.add_picture(img_path, MARGIN + (col_w - iw) // 2,
                                     BODY_TOP + (BODY_H - ih) // 2, width=iw)
                tx_left = MARGIN + col_w + Inches(0.4)
            else:
                s.shapes.add_picture(img_path, MARGIN + col_w + Inches(0.4) + (col_w - iw) // 2,
                                     BODY_TOP + (BODY_H - ih) // 2, width=iw)
                tx_left = MARGIN
            box = s.shapes.add_textbox(tx_left, BODY_TOP, col_w, BODY_H)
            _render_textflow(box.text_frame, text_content, accent, unmatched)
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            y = BODY_TOP
            if text_content:
                if not tables:
                    th = BODY_H if not img_path else Inches(2.2)
                else:
                    th = Inches(2.0)
                box = s.shapes.add_textbox(MARGIN, y, BODY_W, th)
                _render_textflow(box.text_frame, text_content, accent, unmatched)
                y = y + th + Inches(0.15)
            if img_path:
                iw, ih = _img_fit(img_path, BODY_W, Inches(2.6))
                s.shapes.add_picture(img_path, MARGIN + (BODY_W - iw) // 2, y, width=iw)
                y = y + ih + Inches(0.15)
            for tb in tables:
                y = _add_table(s, parse_table(tb.get("html")), MARGIN, y, BODY_W, accent)
                y = y + Inches(0.2)

    if unmatched:
        # a type reached the textflow with no rendering branch and no _DROP entry
        # -> it would have vanished silently. Report it as dropped + warn loudly.
        import sys
        for t in unmatched:
            stats["dropped"][t] = stats["dropped"].get(t, 0) + 1
        print(f"[pptx_export] WARNING: block type(s) {sorted(unmatched)} have no "
              f"flatten disposition (not rendered, not in _DROP) — dropped. Add them "
              f"to blocks.BLOCKS / _DROP.", file=sys.stderr)

    if transition:
        import pptx_transitions
        pptx_transitions.apply_all(prs, transition, transition_speed, transition_dir)
        stats["transition"] = transition

    prs.save(out_path)
    return stats
