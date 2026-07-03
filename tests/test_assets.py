"""Vector asset library (src/assets.py) + its render integration.

SVG is the source of truth; it's recolored on-brand and (for the .pptx, which can't
take SVG) rasterized to a cached PNG via resvg-py. Raster assertions are skipped when
resvg-py isn't installed (the feature degrades to a placeholder there).
"""
import os
import pytest

import assets
import slide_layouts

ICON = "icon-stethoscope.svg"
_HAS_RASTER = assets.rasterize_svg('<svg xmlns="http://www.w3.org/2000/svg" '
                                   'viewBox="0 0 10 10"><rect width="10" height="10"/></svg>') is not None


# ---- recolor / normalize -------------------------------------------------

def test_norm_hex_forms():
    assert assets.norm_hex("1EB16A") == "#1EB16A"
    assert assets.norm_hex("#1eb16a") == "#1EB16A"
    assert assets.norm_hex("abc") == "#AABBCC"      # 3-digit expands
    assert assets.norm_hex(None) == "#000000"       # default
    assert assets.norm_hex("not-a-color") == "#000000"


def test_recolor_swaps_currentcolor_and_keeps_none():
    svg = '<svg fill="none" stroke="currentColor"><path fill="none" d="M0 0"/></svg>'
    out = assets.recolor_svg(svg, "#1EB16A")
    assert "currentColor" not in out
    assert "#1EB16A" in out                         # stroke recolored
    assert 'fill="none"' in out                     # transparent fills untouched


def test_recolor_repaints_explicit_fill():
    out = assets.recolor_svg('<svg><path fill="#000000" d="M0 0"/></svg>', "#1EB16A")
    assert 'fill="#1EB16A"' in out and "#000000" not in out


def test_recolor_strips_comments():
    out = assets.recolor_svg("<!-- tags: x -->\n<svg stroke='currentColor'></svg>", "#123456")
    assert "<!--" not in out


# ---- library listing / resolution ---------------------------------------

def test_library_bundled_and_listed():
    icons = assets.list_icons()
    assert ICON in icons and len(icons) >= 10
    assert all(n.endswith(".svg") for n in icons)


def test_icon_path_is_basename_confined():
    assert assets.icon_path(ICON) and os.path.isfile(assets.icon_path(ICON))
    assert assets.icon_path("../../etc/passwd") is None      # escape attempt -> None
    assert assets.icon_path("does-not-exist.svg") is None


def test_icon_preview_svg_recolored():
    svg = assets.icon_preview_svg(ICON, "#1EB16A")
    assert svg.startswith("<svg") and "#1EB16A" in svg and "currentColor" not in svg
    assert assets.icon_preview_svg("missing.svg") == ""


def test_resolve_images_falls_back_to_library():
    # no images_dir: a bare icon name resolves to the shared library
    out = slide_layouts._resolve_images({"image": ICON}, None)
    assert out["image"].endswith(os.path.join("assets", "icons", ICON))
    # unknown name -> left as-is (placeholder handles it downstream)
    assert slide_layouts._resolve_images({"image": "nope.svg"}, None)["image"] == "nope.svg"
    # a real brand image dir still wins when it has the file
    assert slide_layouts._resolve_images({"image": ICON}, assets.ICON_DIR)["image"] \
        == os.path.join(assets.ICON_DIR, ICON)


# ---- rasterize / render (resvg-py required) ------------------------------

@pytest.mark.skipif(not _HAS_RASTER, reason="resvg-py not installed")
def test_svg_asset_to_png_caches():
    p1 = assets.svg_asset_to_png(assets.icon_path(ICON), "#1EB16A")
    p2 = assets.svg_asset_to_png(assets.icon_path(ICON), "#1EB16A")
    assert p1 and os.path.isfile(p1) and p1 == p2            # cache hit
    from PIL import Image
    im = Image.open(p1)
    assert im.size[0] >= 256 and im.mode == "RGBA"
    # a different color is a different cache entry
    assert assets.svg_asset_to_png(assets.icon_path(ICON), "#0B2C37") != p1


@pytest.mark.skipif(not _HAS_RASTER, reason="resvg-py not installed")
def test_icon_places_in_real_pptx():
    import brand as B
    br = B.load_brand("teletracking")
    out = os.path.join(os.path.dirname(__file__), "_assets_tmp.pptx")
    try:
        slide_layouts.export_deck(
            [{"layout": "imagetext",
              "content": {"title": "Care", "subtitle": "Flow", "image": ICON}}],
            out, brand=br)
        from pptx import Presentation
        p = Presentation(out)
        pics = [sh for s in p.slides for sh in s.shapes if sh.shape_type == 13]
        assert pics                                          # a real picture, not a placeholder
        txt = " ".join(sh.text_frame.text for s in p.slides
                       for sh in s.shapes if sh.has_text_frame)
        assert "could not be placed" not in txt and "Image not found" not in txt
    finally:
        if os.path.exists(out):
            os.remove(out)
