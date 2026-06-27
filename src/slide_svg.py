"""HTML/SVG preview of slide-deck posters — faithful to the .pptx output.

The slide layouts in `slide_layouts.py` draw a poster onto a 13.333"x7.5"
PowerPoint slide via a tiny kit of primitives — `_drawkit(s)` returns
(fill, rect, text), and the renderers also call `s.shapes.add_shape(...)`
directly for ovals/chevrons. This module provides a SECOND backend: a mock
"slide" whose `.shapes` records the same geometry and emits SVG instead of
PowerPoint shapes. We then run the EXACT SAME renderer functions against it,
so the preview shares one geometry source with the real export and cannot
drift from it. No new dependencies; pure stdlib + the existing modules.

The `chart` layout is the one exception — it inserts a native, editable
PowerPoint chart object (not drawkit shapes), so for the preview we reuse the
existing `chart_svg` engine (which already emits SVG) nested under the header
band, with the brand palette supplied as CSS custom properties.

Entry points:
  render_slide_svg(layout, content, brand) -> one <svg> string
  render_deck_svg(slides, brand)           -> list of <svg> strings
"""
import re
import os
import base64
import html as _html

import slide_layouts as SL
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----- canvas geometry -------------------------------------------------------
# The PPTX canvas is 13.333"x7.5" expressed in EMU (914400 EMU per inch). The
# SVG mirrors it at 96 px/inch -> a 1280x720 viewBox (true 16:9), so any EMU
# coordinate the renderers compute maps linearly to a px coordinate.
_EMU_PER_IN = 914400.0
_PXPERIN = 96.0
_SCALE = _PXPERIN / _EMU_PER_IN
VW, VH = 1280, 720

FONT_STACK = ("-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,"
              "'Helvetica Neue',Arial,sans-serif")


def _px(emu):
    return emu * _SCALE


def _ptpx(pt):
    return pt * _PXPERIN / 72.0


def _hexs(c):
    """RGBColor (indexable r,g,b) -> #RRGGBB."""
    try:
        return "#%02X%02X%02X" % (c[0], c[1], c[2])
    except Exception:
        return "#000000"


def _esc(s):
    return _html.escape(str(s if s is not None else ""), quote=True)


# ----- mock python-pptx surface ---------------------------------------------
# Only the API the slide renderers actually touch is implemented. Each mock
# object knows how to serialize itself to SVG; the slide keeps an ordered list
# so paint order (z-order) matches creation order, exactly like PowerPoint.

class _Color:
    __slots__ = ("rgb",)
    def __init__(self):
        self.rgb = None


class _Fill:
    __slots__ = ("fore_color",)
    def __init__(self):
        self.fore_color = _Color()
    def solid(self):
        pass
    def background(self):
        pass


class _Line:
    __slots__ = ("fill", "color", "width")
    def __init__(self):
        self.fill = _Fill()
        self.color = _Color()
        self.width = None             # EMU (a Pt(...) length); None = no stroke


class _Shape:
    """A rect / rounded-rect / oval / chevron. Color is set after creation via
    the kit's fill() helper, so we read it lazily at serialization time."""
    def __init__(self, shape_type, x, y, w, h):
        self.shape_type = shape_type
        self.x, self.y, self.w, self.h = x, y, w, h
        self.fill = _Fill()
        self.line = _Line()
        self._alpha = None        # fill translucency 0-100; set via set_alpha()

    def set_alpha(self, pct):
        """Honor _fill_alpha in the preview. Without this the mock no-op'd alpha, so
        a full-bleed image's scrim overlay rendered OPAQUE — a solid slab hiding the
        photo (preview != .pptx). _fill_alpha duck-types this method for the mock."""
        self._alpha = pct

    def _op(self):
        return (f' fill-opacity="{max(0, min(100, self._alpha)) / 100:.3f}"'
                if self._alpha is not None and self._alpha < 100 else "")

    def _stroke(self):
        """An outline, when the renderer set line width + color (e.g. the cycles hub
        outline, timeline card borders). Without this the mock dropped all outlines,
        so a same-tone shape was invisible in the preview though visible in the .pptx."""
        lc, lw = self.line.color.rgb, self.line.width
        if lc is None or not lw:
            return ""
        return f' stroke="{_hexs(lc)}" stroke-width="{max(0.5, _px(lw)):.1f}"'

    def to_svg(self):
        col = _hexs(self.fill.fore_color.rgb) if self.fill.fore_color.rgb is not None else "none"
        op = self._op() + self._stroke()
        x, y, w, h = _px(self.x), _px(self.y), _px(self.w), _px(self.h)
        st = self.shape_type
        if st == MSO_SHAPE.OVAL:
            return (f'<ellipse cx="{x + w / 2:.1f}" cy="{y + h / 2:.1f}" '
                    f'rx="{w / 2:.1f}" ry="{h / 2:.1f}" fill="{col}"{op}/>')
        if st == MSO_SHAPE.CHEVRON:
            # a right-pointing arrow filling the box
            p = (f'{x:.1f},{y:.1f} {x + w * 0.6:.1f},{y:.1f} {x + w:.1f},{y + h / 2:.1f} '
                 f'{x + w * 0.6:.1f},{y + h:.1f} {x:.1f},{y + h:.1f} '
                 f'{x + w * 0.4:.1f},{y + h / 2:.1f}')
            return f'<polygon points="{p}" fill="{col}"{op}/>'
        rx = min(w, h) * 0.12 if st == MSO_SHAPE.ROUNDED_RECTANGLE else 0
        return (f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                f'rx="{rx:.1f}" fill="{col}"{op}/>')


class _Run:
    __slots__ = ("text", "font")
    def __init__(self):
        self.text = ""
        self.font = type("F", (), {"size": None, "bold": False, "color": _Color()})()


class _Para:
    def __init__(self):
        self.alignment = PP_ALIGN.LEFT
        self.space_after = Pt(0)
        self.space_before = Pt(0)
        self.line_spacing = 1.0
        self.runs = []
    def add_run(self):
        r = _Run(); self.runs.append(r); return r


class _TextFrame:
    def __init__(self):
        self.word_wrap = True
        self.vertical_anchor = MSO_ANCHOR.TOP
        self.margin_left = self.margin_right = self.margin_top = self.margin_bottom = Pt(2)
        self.paragraphs = [_Para()]
    def add_paragraph(self):
        p = _Para(); self.paragraphs.append(p); return p


class _Textbox:
    def __init__(self, x, y, w, h):
        self.x, self.y, self.w, self.h = x, y, w, h
        self.text_frame = _TextFrame()

    # --- text layout: greedy word-wrap with estimated glyph widths ----------
    def _tokens(self, para):
        toks = []
        for r in para.runs:
            size = r.font.size.pt if r.font.size is not None else 12.0
            col = _hexs(r.font.color.rgb) if r.font.color.rgb is not None else "#000000"
            bold = bool(r.font.bold)
            words = (r.text or "").split(" ")
            for i, w in enumerate(words):
                if w == "" and i not in (0, len(words) - 1):
                    continue
                toks.append((w, _ptpx(size), col, bold))
        return toks

    @staticmethod
    def _wmeas(word, size_px, bold):
        return len(word) * size_px * (0.56 if bold else 0.52)

    def to_svg(self):
        bx, by = _px(self.x), _px(self.y)
        bw, bh = _px(self.w), _px(self.h)
        pad = 3.0
        inner = max(8.0, bw - 2 * pad)
        tf = self.text_frame
        # build wrapped lines (each = list of tokens), tracking line heights
        lines, line_h = [], []
        for pi, para in enumerate(tf.paragraphs):
            toks = self._tokens(para)
            if not toks:
                continue
            cur, cur_w, maxsz = [], 0.0, 0.0
            space = 0.28
            for (word, sz, col, bold) in toks:
                ww = self._wmeas(word, sz, bold)
                sw = sz * space
                add = ww + (sw if cur else 0)
                if cur and cur_w + add > inner:
                    lines.append((cur, para.alignment)); line_h.append(maxsz * para.line_spacing * 1.22)
                    cur, cur_w, maxsz = [], 0.0, 0.0
                    add = ww
                cur.append((word, sz, col, bold, sz * space))
                cur_w += add
                maxsz = max(maxsz, sz)
            if cur:
                lines.append((cur, para.alignment)); line_h.append(maxsz * para.line_spacing * 1.22)
            # paragraph gap
            if pi < len(tf.paragraphs) - 1 and lines:
                line_h[-1] += para.space_after.pt * _PXPERIN / 72.0
        if not lines:
            return ""
        total = sum(line_h)
        if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
            cy = by + max(0.0, (bh - total) / 2.0)
        elif tf.vertical_anchor == MSO_ANCHOR.BOTTOM:
            cy = by + max(0.0, bh - total)
        else:
            cy = by + pad
        out = []
        for (toks, align), lh in zip(lines, line_h):
            baseline = cy + lh * 0.78
            if align == PP_ALIGN.CENTER:
                anchor, tx = "middle", bx + bw / 2.0
            elif align == PP_ALIGN.RIGHT:
                anchor, tx = "end", bx + bw - pad
            else:
                anchor, tx = "start", bx + pad
            spans = []
            for (word, sz, col, bold, sw) in toks:
                fw = ' font-weight="700"' if bold else ""
                spans.append(f'<tspan font-size="{sz:.1f}" fill="{col}"{fw}>{_esc(word)}</tspan>')
            # join tspans with normal spaces (xml:space preserve keeps them)
            joined = '<tspan> </tspan>'.join(spans)
            out.append(f'<text x="{tx:.1f}" y="{baseline:.1f}" text-anchor="{anchor}" '
                       f'xml:space="preserve">{joined}</text>')
            cy += lh
        return "".join(out)


_MAX_PREVIEW_PX = 900   # downscale embedded images so a preview/thumbnail SVG stays light


def _data_uri(path):
    """Embed an image as a base64 data URI so the SVG renders standalone (in the
    slideshow overlay and the row thumbnails, which are injected as innerHTML and
    can't load file:// refs). Downscaled via Pillow to keep the payload small;
    falls back to the raw bytes if Pillow can't open it."""
    try:
        import io
        from PIL import Image
        with Image.open(path) as im:
            if im.mode not in ("RGB", "RGBA", "L"):
                im = im.convert("RGBA")
            if max(im.size) > _MAX_PREVIEW_PX:
                im = im.copy()
                im.thumbnail((_MAX_PREVIEW_PX, _MAX_PREVIEW_PX))
            buf = io.BytesIO()
            im.save(buf, format="PNG")
            return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")
    except Exception:
        try:
            with open(path, "rb") as fh:
                raw = fh.read()
            ext = os.path.splitext(path)[1].lower().lstrip(".")
            mime = {"jpg": "jpeg", "jpeg": "jpeg", "png": "png", "gif": "gif",
                    "webp": "webp", "svg": "svg+xml"}.get(ext, "png")
            return f"data:image/{mime};base64," + base64.b64encode(raw).decode("ascii")
        except Exception:
            return ""


class _Picture:
    """A placed image. _place_image sizes the box for either a 'contain' fit
    (box already matches the image aspect) or a 'cover' fit (full box + crop_*
    set by _crop_cover). When cropped we slice (cover); otherwise we meet."""
    __slots__ = ("path", "x", "y", "w", "h", "name",
                 "crop_left", "crop_right", "crop_top", "crop_bottom")
    def __init__(self, path, x, y, w, h):
        self.path, self.x, self.y, self.w, self.h = path, x, y, w, h
        self.name = ""
        self.crop_left = self.crop_right = self.crop_top = self.crop_bottom = 0.0
    def to_svg(self):
        x, y, w, h = _px(self.x), _px(self.y), _px(self.w), _px(self.h)
        uri = _data_uri(self.path)
        if not uri:
            return (f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                    f'fill="#EEF2F3"/>')
        cropped = max(self.crop_left, self.crop_right, self.crop_top, self.crop_bottom) > 0
        par = "xMidYMid slice" if cropped else "xMidYMid meet"
        return (f'<image x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                f'preserveAspectRatio="{par}" href="{_esc(uri)}"/>')


class _Shapes:
    def __init__(self, objs):
        self._objs = objs
    def add_shape(self, shape_type, x, y, w, h):
        sh = _Shape(shape_type, x, y, w, h); self._objs.append(sh); return sh
    def add_textbox(self, x, y, w, h):
        tb = _Textbox(x, y, w, h); self._objs.append(tb); return tb
    def add_picture(self, path, x, y, width, height):
        pic = _Picture(path, x, y, width, height); self._objs.append(pic); return pic
    def add_chart(self, *a, **k):  # charts use a dedicated SVG path, never the mock
        raise NotImplementedError("chart layout is rendered via chart_svg, not the mock backend")


class _Slide:
    def __init__(self):
        self._objs = []
        self.shapes = _Shapes(self._objs)
    def body_svg(self, animate=False):
        if not animate:
            return "".join(o.to_svg() for o in self._objs)
        # cascade: wrap each shape in a group with an incremental delay so the
        # slideshow can build the slide in, like the course entrance animations.
        out = []
        for i, o in enumerate(self._objs):
            out.append(f'<g class="nv-a" style="animation-delay:{i * 45}ms">{o.to_svg()}</g>')
        return "".join(out)


# ----- public renderers ------------------------------------------------------

# IMPORTANT: the base state is VISIBLE (no opacity:0 on .nv-a itself). The
# entrance comes from `animation ... both`, whose `from` keyframe (opacity:0) is
# applied only while the animation runs/waits (backwards+forwards fill). If the
# animation never runs — reduced-motion, a webview that won't run @keyframes in
# an innerHTML-inserted SVG, etc. — the element keeps its natural opacity:1 and
# the slide is still shown. (A prior version set opacity:0 as the base and went
# blank when the animation didn't fire.)
_ANIM_CSS = ("@keyframes nv-rise{from{opacity:0;transform:translateY(14px)}"
             "to{opacity:1;transform:none}}"
             ".nv-slide-svg .nv-a{animation:nv-rise .5s ease both;"
             "transform-box:fill-box}"
             "@media (prefers-reduced-motion:reduce){.nv-slide-svg .nv-a"
             "{animation:none}}")


def _wrap(body, animate=False):
    extra = _ANIM_CSS if animate else ""
    return (f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}" '
            f'preserveAspectRatio="xMidYMid meet" class="nv-slide-svg">'
            f'<style>.nv-slide-svg text{{font-family:{FONT_STACK}}}{extra}</style>'
            f'<rect width="{VW}" height="{VH}" fill="#FFFFFF"/>{body}</svg>')


_CHART_VARS = {
    # chart_svg series colors + ink/line, mapped to the brand palette so the
    # nested chart matches the native-PPTX chart's primary/secondary/... series.
    "--brand-accent": "primary", "--brand-accent2": "secondary",
    "--brand-accent-ink": "tertiary", "--brand-heading": "dark",
    "--brand-ink": "dark", "--brand-ink-soft": "grey", "--brand-line": "_line",
}


def _chart_slide_svg(content, pal):
    import chart_svg
    # header band + source line drawn faithfully via the real kit
    s = _Slide()
    fill, rect, text = SL._drawkit(s)
    SL._header(s, content, pal, rect, text)
    if content.get("source"):
        from pptx.util import Inches
        text(Inches(0.6), Inches(6.92), Inches(12.13), Inches(0.4),
             [[("Source: ", 10, pal["dark"], True), (content["source"], 10, pal["grey"], False)]])
    head = s.body_svg()
    # the chart itself: reuse chart_svg, extract its inner <svg>, place + color it
    block = {k: content.get(k) for k in ("chart", "categories", "series", "xLabel", "yLabel")}
    figure = chart_svg.render_chart(block)
    m = re.search(r"<svg\b.*?</svg>", figure, re.S)
    if not m:
        return _wrap(head)
    inner = m.group(0)
    from pptx.util import Inches
    cx, cy = _px(Inches(0.6)), _px(Inches(1.5))
    cw, ch = _px(Inches(12.13)), _px(Inches(5.2))
    light = "#E2E8F0"
    style = ";".join(
        f"{var}:{(light if role == '_line' else _hexs(pal[role]))}"
        for var, role in _CHART_VARS.items())
    inner = inner.replace(
        "<svg", f'<svg x="{cx:.1f}" y="{cy:.1f}" width="{cw:.1f}" height="{ch:.1f}" '
        f'style="{style}"', 1)
    return _wrap(head + inner)


def render_slide_svg(layout, content, brand=None, images_dir=None, animate=False, theme=None):
    """Render ONE slide's poster to an <svg> string, reusing the exact PPTX
    geometry. `brand` is a Brand object or dict (or None for neutral).
    `images_dir` resolves a bare `image` filename to a file on disk. `animate`
    wraps shapes so the slideshow can cascade them in (off for thumbnails).
    `theme` (dark|light) is the cross-cutting per-slide override, so the preview
    matches a themed .pptx exactly (same _slide_pal resolution)."""
    base_pal = SL._palette_of(brand)
    bp = SL._load_blueprint(brand)
    layout = layout if layout in SL.LAYOUTS else "infographic"
    pal, bg = SL._slide_pal(base_pal, brand, bp, layout, theme)
    if layout == "chart":
        return _chart_slide_svg(content or {}, pal)
    s = _Slide()
    fill, rect, text = SL._drawkit(s)
    if bg:
        SL._place_background(s, bg)            # match the .pptx background substrate
    content = SL._resolve_images(content or {}, images_dir)
    SL.RENDERERS[layout](s, content, pal, fill, rect, text)
    return _wrap(s.body_svg(animate=animate), animate=animate)


def render_deck_svg(slides, brand=None, images_dir=None, animate=False):
    """Render an ordered list of {layout, content} slides to a list of <svg>
    strings. `content` may be a dict or a JSON string (from the editor). A slide
    that fails to parse or render becomes an error placeholder rather than
    killing the whole deck preview."""
    import json
    out = []
    for i, spec in enumerate(slides or []):
        spec = spec or {}
        layout = spec.get("layout", "infographic")
        try:
            content = spec.get("content") or {}
            if isinstance(content, str):
                content = json.loads(content) if content.strip() else {}
            out.append(render_slide_svg(layout, content, brand, images_dir, animate,
                                        spec.get("theme")))
        except Exception as e:
            out.append(_wrap(
                f'<rect width="{VW}" height="{VH}" fill="#FBE9E7"/>'
                f'<text x="{VW/2}" y="{VH/2}" text-anchor="middle" font-size="22" '
                f'fill="#B71C1C">Slide {i + 1} ({_esc(layout)}) could not be previewed: '
                f'{_esc(str(e)[:120])}</text>'))
    return out
