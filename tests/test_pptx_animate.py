"""pptx_animate — per-element entrance animations injected as OOXML.

These confirm the timing tree is injected (one per shape, staggered), that the
.pptx reopens (the hand-authored XML parses/loads), and that 'none' removes it.
Real PowerPoint rendering is a manual check; here we guard well-formedness and
that a build never crashes from animation.
"""
import os
import tempfile

import slide_layouts as SL
import brand as brandmod
import pptx_animate
from pptx import Presentation
from pptx.oxml.ns import qn

DECK = [{"layout": "divider", "content": {"title": "T"}},
        {"layout": "cards", "content": {"title": "C", "cards": [
            {"title": "a", "body": "b"}, {"title": "c", "body": "d"}]}}]


def _build(animate, transition=None):
    out = tempfile.mktemp(suffix=".pptx")
    SL.export_deck(DECK, out, brand=brandmod.load_brand("_default"),
                   transition=transition, animate=animate)
    return out


def test_animation_injects_timing_and_reopens():
    for eff in ("rise", "fade", "flyleft", "flyright"):
        out = _build(eff)
        prs = Presentation(out)                       # reopen validates the XML
        timed = [s for s in prs.slides if s._element.find(qn("p:timing")) is not None]
        assert len(timed) == len(DECK), f"{eff}: every slide should carry <timing>"
        os.unlink(out)


def test_animation_targets_each_shape():
    out = _build("rise")
    prs = Presentation(out)
    s0 = list(prs.slides)[0]
    timing = s0._element.find(qn("p:timing"))
    spids = {t.get("spid") for t in timing.iter(qn("p:spTgt"))}
    shape_ids = {str(sh.shape_id) for sh in s0.shapes}
    assert shape_ids and shape_ids <= spids        # every shape is an animation target
    os.unlink(out)


def test_none_leaves_no_timing():
    out = _build("none")
    prs = Presentation(out)
    assert all(s._element.find(qn("p:timing")) is None for s in prs.slides)
    os.unlink(out)


def test_animation_coexists_with_transition():
    # transition (<p:transition>) and animation (<p:timing>) on the same slide
    out = _build("rise", transition="fade")
    s0 = list(Presentation(out).slides)[0]._element
    assert s0.find(qn("p:transition")) is not None
    assert s0.find(qn("p:timing")) is not None
    os.unlink(out)


def test_apply_unknown_effect_raises():
    import pytest
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(ValueError):
        pptx_animate.apply(s, "explode")


def test_clustering_groups_overlapping_shapes():
    # Two clusters: a card-like trio (panel + accent strip inside it + a label
    # inside it) and, far away, a separate box. Inner pieces overlap the panel →
    # one cluster; the distant box → its own cluster.
    from pptx import Presentation as P
    from pptx.util import Inches as I
    from pptx.enum.shapes import MSO_SHAPE
    prs = P()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(1), I(1), I(3), I(2))
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(1), I(1), I(3), I(0.1))
    label = s.shapes.add_textbox(I(1.2), I(1.4), I(2.6), I(0.5))
    far = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(8), I(5), I(2), I(1))
    clusters = pptx_animate.cluster_shapes([panel, accent, label, far])
    assert len(clusters) == 2
    big = max(clusters, key=len)
    assert {sh.shape_id for sh in big} == {panel.shape_id, accent.shape_id, label.shape_id}
    assert [sh.shape_id for c in clusters for sh in c].count(far.shape_id) == 1


def test_real_powerpoint_build_structure():
    # The generator emits the structure real PowerPoint uses (reverse-engineered
    # from James's reference deck): one mainSeq -> ONE container -> N timed
    # sub-steps -> grouped effects. Assert that exact nesting + auto-start.
    out = _build("rise")
    s1 = list(Presentation(out).slides)[1]            # the cards slide
    timing = s1._element.find(qn("p:timing"))
    mainseq = next(c for c in timing.iter(qn("p:cTn")) if c.get("nodeType") == "mainSeq")
    # mainSeq holds exactly one container <p:par>
    containers = mainseq.find(qn("p:childTnLst")).findall(qn("p:par"))
    assert len(containers) == 1
    container_ctn = containers[0].find(qn("p:cTn"))
    # auto-start: container waits delay=0 (not "indefinite" = on-click)
    assert container_ctn.find(qn("p:stCondLst") + "/" + qn("p:cond")).get("delay") == "0"
    substeps = container_ctn.find(qn("p:childTnLst")).findall(qn("p:par"))
    effects = [c for c in timing.iter(qn("p:cTn")) if c.get("presetClass") == "entr"]
    assert len(effects) == len(list(s1.shapes))       # one entrance effect per shape
    assert 2 <= len(substeps) < len(effects)          # grouped into timed sub-steps
    # sub-step delays are absolute and strictly increasing (timed cascade)
    delays = [int(p.find(qn("p:cTn") + "/" + qn("p:stCondLst") + "/" + qn("p:cond")).get("delay"))
              for p in substeps]
    assert delays == sorted(delays) and delays[0] == 0 and delays[-1] > 0
    for p in substeps:                                 # each sub-step nests >=1 effect
        assert [c for c in p.iter(qn("p:cTn")) if c.get("presetClass") == "entr"]
    os.unlink(out)
