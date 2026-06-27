"""Standalone INFOGRAPHIC slides (.pptx) from a content file — brand-driven.

A reusable slide-template target, separate from course IR. Feed a JSON content
file describing the slide and it renders one editable, on-brand PowerPoint slide.
Colors come from the active brand profile (palette / defaultAccent), so the same
template renders client-branded or neutral depending on --brand.

Layouts:
  infographic  — header band (title + subtitle) | left "challenge" column
                 (heading + intro + bulleted items + callout box) | right
                 "framework" column (heading + sublabel + numbered cards) |
                 a goals row (label chip + cards) | footer band.
                 Every section is optional; omit a key to skip it.

Content schema (all sections optional):
{
  "title": "...", "subtitle": "...",
  "left":  {"heading": "...", "intro": "...",
            "items": [["bold lead", " rest of line"], "plain line", ...],
            "callout": "..."},
  "right": {"heading": "...", "sublabel": "...",
            "cards": [{"num": "1", "title": "...", "body": "...",
                       "accent": "primary|secondary|tertiary|dark"}, ...]},
  "goals": {"label": ["OUR", "GOALS"],
            "items": [{"title": "...", "body": "...", "accent": "primary"}, ...]},
  "footer": "..."
}
"""
import json
import os
import layouts
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0x44, 0x4B, 0x4F)
CARD  = RGBColor(0xF4, 0xF6, 0xF7)

W, H = Inches(13.333), Inches(7.5)


def _hex(s, fallback):
    s = (s or "").lstrip("#")
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return fallback


def _lighten(c, amt=0.88):
    """Blend a color toward white (amt = fraction of white)."""
    return RGBColor(int(c[0] + (255 - c[0]) * amt),
                    int(c[1] + (255 - c[1]) * amt),
                    int(c[2] + (255 - c[2]) * amt))


def palette(brand):
    """Resolve 4 brand roles (+ derived tint) with graceful fallbacks."""
    p = (brand.get("palette", {}) if brand else {}) or {}
    accent = (brand.get("defaultAccent") if brand else None) or "#3B82F6"
    primary   = _hex(p.get("green")  or p.get("accent")  or accent, RGBColor(0x1E, 0xB1, 0x6A))
    secondary = _hex(p.get("teal")   or p.get("accent2") or accent, RGBColor(0x06, 0x96, 0x96))
    tertiary  = _hex(p.get("blue")   or p.get("accent2") or accent, RGBColor(0x53, 0x9B, 0xD2))
    dark      = _hex(p.get("deepNavy") or p.get("navy") or p.get("heading") or "#1F2937",
                     RGBColor(0x0B, 0x2C, 0x37))
    return {"primary": primary, "secondary": secondary, "tertiary": tertiary,
            "dark": dark, "tint": _lighten(primary), "white": WHITE,
            "grey": GREY, "card": CARD}


def _accent(pal, name, default="primary"):
    return pal.get(name, pal.get(default, pal["primary"]))


LAYOUTS = ("infographic", "process", "comparison", "timeline", "divider", "chart",
           "image", "imagetext", "cards", "quote", "statement", "bullets",
           "agenda", "closing", "sectionheader", "cycles")


def _resolve_images(content, images_dir):
    """Return content with the `image` field resolved to a usable path: an
    absolute path is kept as-is; a bare filename is joined to `images_dir`.
    A blank/missing image is left alone. Shallow-copies only when it rewrites,
    so callers can pass the original dict safely."""
    if not isinstance(content, dict):
        return content
    img = content.get("image")
    if not img or not isinstance(img, str) or os.path.isabs(img) or not images_dir:
        return content
    c = dict(content)
    c["image"] = os.path.join(images_dir, img)
    return c


def _img_size(path):
    """(width, height) in px via Pillow, or (0, 0) if unreadable. Pillow is a
    python-pptx dependency, so it's present wherever this export module runs."""
    try:
        from PIL import Image
        with Image.open(path) as im:
            return im.size
    except Exception:
        return (0, 0)


def _crop_cover(pic, iw, ih, bw, bh):
    """Crop a box-filling picture down to the box's aspect ratio (CSS `cover`):
    the image fills the whole box edge-to-edge, trimming the overflow centered.
    Works on the real python-pptx Picture and on slide_svg's mock alike (both
    expose crop_left/right/top/bottom)."""
    if not (iw and ih and bw and bh):
        return
    img_ar, box_ar = iw / ih, bw / bh
    if img_ar > box_ar:                          # image wider -> trim sides
        frac = max(0.0, (1 - box_ar / img_ar) / 2)
        pic.crop_left = frac; pic.crop_right = frac
    elif img_ar < box_ar:                        # image taller -> trim top/bottom
        frac = max(0.0, (1 - img_ar / box_ar) / 2)
        pic.crop_top = frac; pic.crop_bottom = frac


def _place_image(s, path, x, y, w, h, pal, rect, text, fit=None):
    """Place an image in the box per the brand's design.image tokens. Default is
    template-faithful: a frameless, edge-to-edge COVER picture (fills the box,
    overflow cropped centered) — no floating rounded card. fit="contain" keeps the
    whole image (aspect-preserving, centered) for screenshots/diagrams. The
    `framePt` token draws a thin mat behind the picture; `shadow` lifts it off the
    page. A missing/unreadable file becomes a labeled placeholder, never a crash.
    The SAME code drives the SVG preview (slide_svg reuses this), so preview ==.pptx."""
    tok = (pal.get("design") or {}).get("image") or {}
    fit = fit or tok.get("fit", "cover")
    framePt = tok.get("framePt", 0) or 0
    frameColor = _tok_color(tok.get("frameColor", "card"), pal, pal.get("card"))
    if not path or not os.path.isfile(path):
        rect(x, y, w, h, pal["card"], rounded=True)
        text(x, y, w, h,
             [[("Image not found", 13, pal["cardInk"], True)],
              [(os.path.basename(path) if path else "(no image set)", 11, pal["muted"], False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None
    iw, ih = _img_size(path)

    def _mat(bx, by, bw, bh):
        if framePt and frameColor is not None:
            m = Pt(framePt)
            rect(int(bx - m), int(by - m), int(bw + 2 * m), int(bh + 2 * m), frameColor)

    def _shadow(pic):
        if tok.get("shadow"):
            try:
                _add_shadow(pic)
            except Exception:
                pass

    if fit == "cover":
        try:
            _mat(x, y, w, h)
            pic = s.shapes.add_picture(path, x, y, w, h)   # fill the box
            _crop_cover(pic, iw, ih, int(w), int(h))        # then trim to box aspect
            _shadow(pic)
            return True
        except Exception:
            pass                                            # fall through to contain
    # contain: scale whole image into the box, centered
    pad = Inches(tok.get("padIn", 0.0))
    bw, bh = int(w - 2 * pad), int(h - 2 * pad)
    if iw and ih:
        scale = min(bw / iw, bh / ih)
        dw, dh = int(iw * scale), int(ih * scale)
    else:
        dw, dh = bw, bh
    px = int(x) + int((int(w) - dw) / 2)
    py = int(y) + int((int(h) - dh) / 2)
    _mat(px, py, dw, dh)
    try:
        pic = s.shapes.add_picture(path, px, py, dw, dh)
        _shadow(pic)
        return True
    except Exception:
        text(x, y, w, h, [[("Image could not be placed", 12, pal["cardInk"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None


def role(shape, name):
    """Tag a shape with its animation ROLE via its name (visible in PowerPoint's
    Selection Pane and persisted in the .pptx). Scheme `<tier>.<part>[.<unit>]`
    where tier is background|focal|content (see the animation library): background
    loads first, then focal containers, then content. Returns the shape."""
    if name and shape is not None:
        shape.name = name
    return shape


def _drawkit(s):
    """Return (fill, rect, text) drawing helpers bound to one slide `s`."""
    def fill(sh, c):
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

    def rect(x, y, w, h, color, rounded=False, role=None):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
        fill(shp, color)
        if role:
            shp.name = role
        return shp

    def text(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=2, line=1.0, role=None):
        tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        for m in ("left", "right", "top", "bottom"):
            setattr(tf, f"margin_{m}", Pt(2))
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(space_after)
            p.space_before = Pt(0); p.line_spacing = line
            for (txt, size, color, bold) in para:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        if role:
            tb.name = role
        return tb
    return fill, rect, text


def _render_infographic(s, content, pal, fill, rect, text):
    """INFOGRAPHIC layout — header band | left challenge column | right numbered
    cards | goals row | footer. Every section is optional."""
    cycle = ["primary", "secondary", "tertiary", "dark"]
    # ---- header band (shared) ----
    _header(s, content, pal, rect, text)

    # ---- left column ----
    left = content.get("left") or {}
    LX, LW = Inches(0.55), Inches(6.35)
    if left.get("heading"):
        text(LX, Inches(1.45), LW, Inches(0.35), [[(left["heading"], 16, pal["primary"], True)]])
    if left.get("intro"):
        text(LX, Inches(1.88), LW, Inches(0.3), [[(left["intro"], 11.5, pal["muted"], False)]])
    cy, row_h = Inches(2.28), Inches(0.46)
    for item in left.get("items", []):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, LX, cy + Emu(int(Inches(0.06))),
                                 Inches(0.13), Inches(0.13))
        fill(dot, pal["primary"])
        if isinstance(item, (list, tuple)) and len(item) == 2:
            runs = [[(item[0], 11.5, pal["ink"], True), (item[1], 11.5, pal["muted"], False)]]
        else:
            runs = [[(str(item), 11.5, pal["muted"], False)]]
        text(LX + Inches(0.28), cy, LW - Inches(0.28), row_h, runs)
        cy = cy + row_h
    if left.get("callout"):
        rect(LX, Inches(4.7), LW, Inches(1.0), pal["tint"], rounded=True)
        rect(LX, Inches(4.7), Inches(0.09), Inches(1.0), pal["primary"])
        text(LX + Inches(0.28), Inches(4.78), LW - Inches(0.45), Inches(0.85),
             [[(left["callout"], 12, pal["dark"], True)]], anchor=MSO_ANCHOR.MIDDLE, line=1.05)

    # ---- right column ----
    right = content.get("right") or {}
    RX, RW = Inches(7.15), Inches(5.65)
    if right.get("heading"):
        text(RX, Inches(1.45), RW, Inches(0.35), [[(right["heading"], 16, pal["secondary"], True)]])
    if right.get("sublabel"):
        text(RX, Inches(1.84), RW, Inches(0.3), [[(right["sublabel"], 10.5, pal["muted"], True)]])
    cards = right.get("cards", [])
    yy, ch, gap = Inches(2.24), Inches(0.78), Inches(0.12)
    cycle = ["primary", "secondary", "tertiary", "dark"]
    for i, c in enumerate(cards):
        accent = _accent(pal, c.get("accent", cycle[i % 4]))
        rect(RX, yy, RW, ch, pal["card"], rounded=True)
        rect(RX, yy, Inches(0.62), ch, accent, rounded=True)
        text(RX, yy, Inches(0.62), ch, [[(str(c.get("num", i + 1)), 24, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(RX + Inches(0.78), yy + Inches(0.07), RW - Inches(0.9), ch,
             [[(c.get("title", ""), 12.5, pal["cardInk"], True)],
              [(c.get("body", ""), 10.5, pal["muted"], False)]], line=1.0, space_after=1)
        yy = yy + ch + gap

    # ---- goals row ----
    goals = content.get("goals") or {}
    GY = Inches(5.95)
    label = goals.get("label") or []
    if label:
        rect(Inches(0.55), GY, Inches(1.7), Inches(0.9), pal["dark"], rounded=True)
        runs = [[(label[0], 13, pal["white"], True)]]
        if len(label) > 1:
            runs.append([(label[1], 13, pal["primary"], True)])
        text(Inches(0.55), GY, Inches(1.7), Inches(0.9), runs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    gitems = goals.get("items", [])
    gx = Inches(2.45) if label else Inches(0.55)
    gw = Inches(3.42); ggap = Inches(0.13)
    for i, g in enumerate(gitems):
        accent = _accent(pal, g.get("accent", cycle[i % 4]))
        rect(gx, GY, gw, Inches(0.9), pal["card"], rounded=True)
        rect(gx, GY, gw, Inches(0.08), accent, rounded=True)
        text(gx + Inches(0.2), GY + Inches(0.13), gw - Inches(0.35), Inches(0.8),
             [[(g.get("title", ""), 12, pal["cardInk"], True)],
              [(g.get("body", ""), 10, pal["muted"], False)]], line=1.0, space_after=1)
        gx = gx + gw + ggap

    # ---- footer ----
    footer = content.get("footer")
    if footer:
        rect(0, Inches(7.06), W, Inches(0.44), pal["primary"])
        text(0, Inches(7.06), _footer_tw(pal), Inches(0.44), [[(footer, 12, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _place_band_logo(s, pal)
    return {"layout": "infographic", "cards": len(cards), "goals": len(gitems)}


def _brand_logo(brand):
    """Absolute path to the brand's logo asset (for the header mark), or None.
    Resolves via Brand.asset() so it falls back to brands/_default; a plain dict
    or None brand yields None. Never raises."""
    try:
        if hasattr(brand, "asset"):
            return brand.asset(brand.get("logo") or "logo.svg")
    except Exception:
        pass
    return None


def _palette_of(brand):
    pal = palette(brand if isinstance(brand, dict) else getattr(brand, "data", {}) or {})
    pal["logo"] = _brand_logo(brand)
    # AGNOSTIC SEMANTIC DEFAULTS — these reproduce today's "previous templates"
    # look (dark header band, navy ink). A brand blueprint overrides them per
    # layout in _slide_pal(); a brand with no blueprint keeps exactly this.
    pal.setdefault("ink", pal["dark"])         # primary body/heading text
    pal.setdefault("muted", pal["grey"])       # secondary text
    pal.setdefault("cardInk", pal["dark"])     # text on a card surface
    pal.setdefault("title_ink", pal["white"])  # header title (white on the band)
    pal.setdefault("rule_color", pal["primary"])
    pal.setdefault("rail", pal["dark"])        # timeline axis rail
    pal.setdefault("node", pal["primary"])     # timeline node fill
    pal.setdefault("header_mode", "band")      # band (default) | rule | none
    pal.setdefault("bg_image", None)
    return pal


def _load_blueprint(brand):
    """The brand's master blueprint (brands/<brand>/blueprint.json) or None.
    Resolved through Brand.asset so it falls back to brands/_default. A brand
    with no blueprint => None => the engine's neutral, client-agnostic look."""
    try:
        if hasattr(brand, "asset"):
            p = brand.asset("blueprint.json")
            if p and os.path.isfile(p):
                with open(p, encoding="utf-8") as fh:
                    return json.load(fh)
    except Exception:
        pass
    return None


def _bp_color(val, pal, bpal, fallback):
    """Resolve a blueprint color: a literal #hex, a base-palette role
    (primary/secondary/tertiary/dark/tint/...), or a brand-palette key
    (green/teal/deepNavy/accentGreen/...). Falls back gracefully."""
    if not val:
        return fallback
    if isinstance(val, str) and val.startswith("#"):
        return _hex(val, fallback)
    if val in pal and isinstance(pal[val], RGBColor):
        return pal[val]
    if isinstance(bpal, dict) and val in bpal:
        return _hex(bpal[val], fallback)
    return fallback


def _slide_pal(base_pal, brand, bp, layout, theme_override=None):
    """Return (pal, bg_image_path) for ONE layout. With no blueprint the pal is
    unchanged (agnostic defaults) and bg is None. With a blueprint, the layout's
    theme overlays semantic colors + header mode and selects the full-bleed
    background art + a theme-appropriate logo. Engine stays brand-agnostic — all
    specifics live in the brand's blueprint, not here.

    `theme_override` is the cross-cutting per-slide `theme: dark|light` flag: it
    swaps ONLY the color theme (semantic colors + background art + theme logo),
    keeping the layout's header mode and structure. It is honored only when the
    brand actually defines that theme; otherwise the layout's default theme is
    kept. A brand with NO blueprint ignores it (there are no themes to switch)."""
    pal = dict(base_pal)
    pal["design"] = (bp or {}).get("design") or {}
    if not bp:
        return pal, None
    bpal = {}
    if hasattr(brand, "get"):
        bpal = brand.get("palette", {}) or {}
    elif isinstance(brand, dict):
        bpal = brand.get("palette", {}) or {}
    lay = dict(bp.get("defaults", {}))
    lay.update(bp.get("layouts", {}).get(layout, {}))
    theme_name = lay.get("theme", "light")
    if isinstance(theme_override, str) and theme_override.lower() in bp.get("themes", {}):
        theme_name = theme_override.lower()      # cross-cutting per-slide override
    theme = bp.get("themes", {}).get(theme_name, {})
    pal["theme"] = theme_name
    pal["header_mode"] = lay.get("header", "rule")
    pal["ink"] = _bp_color(theme.get("ink"), pal, bpal, base_pal["dark"])
    pal["muted"] = _bp_color(theme.get("muted"), pal, bpal, base_pal["grey"])
    pal["card"] = _bp_color(theme.get("card"), pal, bpal, base_pal["card"])
    pal["cardInk"] = _bp_color(theme.get("cardInk"), pal, bpal, pal["ink"])
    pal["title_ink"] = _bp_color(theme.get("title"), pal, bpal, pal["ink"])
    pal["rule_color"] = _bp_color(theme.get("rule"), pal, bpal, base_pal["primary"])
    pal["rail"] = _bp_color(theme.get("rail"), pal, bpal, base_pal["grey"])
    pal["node"] = _bp_color(theme.get("node"), pal, bpal, base_pal["primary"])
    # full-bleed background art for this theme
    bg_key = theme.get("bg", theme_name)
    bg_rel = (bp.get("backgrounds", {}) or {}).get(bg_key)
    bg_path = brand.asset(bg_rel) if (bg_rel and hasattr(brand, "asset")) else None
    pal["bg_image"] = bg_path
    # theme-appropriate wordmark (white on dark, color on light) overrides base
    logo_rel = (bp.get("logos", {}) or {}).get(theme.get("logo"))
    logo_path = brand.asset(logo_rel) if (logo_rel and hasattr(brand, "asset")) else None
    if logo_path:
        pal["logo"] = logo_path
    # white-on-dark wordmark for sitting INSIDE the footer band (the band is the
    # brand primary/green; brand book: white logo on green/navy backgrounds).
    white_rel = (bp.get("logos", {}) or {}).get("dark")
    pal["footer_logo"] = brand.asset(white_rel) if (white_rel and hasattr(brand, "asset")) else None
    return pal, bg_path


def _place_background(s, path):
    """Lay the brand's full-bleed background art behind everything (cover-fit).
    Added before any other shape so it sits at the bottom of the z-order."""
    if not path or not os.path.isfile(path):
        return
    try:
        iw, ih = _img_size(path)
        pic = s.shapes.add_picture(path, 0, 0, W, H)
        _crop_cover(pic, iw, ih, int(W), int(H))
        pic.name = "bg.image"
    except Exception:
        pass


# ----------------------------------------------- DESIGN TOKENS (design language)
# The template's elements decomposed into tokens (radius, outline pt, fill alpha,
# shadow, type scale, spacing) the generative renderer draws from — so slides are
# template-faithful AND adapt to content. Tokens live in the brand blueprint's
# `design` block (pal["design"]); absent that, the engine's built-in defaults
# below reproduce the prior agnostic look.

def _dtok(pal, group, key, default):
    """Read a design token pal['design'][group][key], else `default`."""
    return ((pal.get("design") or {}).get(group) or {}).get(key, default)


def _tok_color(name, pal, default=None):
    """Resolve a token color: '#hex', a semantic pal key (card/ink/node/...), or
    'none' (-> None, meaning transparent)."""
    if not name or name == "none":
        return None
    if isinstance(name, str) and name.startswith("#"):
        return _hex(name, default or pal["dark"])
    v = pal.get(name)
    return v if isinstance(v, RGBColor) else (default)


def _set_radius(shp, pct):
    try:
        shp.adjustments[0] = max(0.0, min(0.5, pct / 100.0))
    except Exception:
        pass


def _fill_alpha(shp, pct):
    """Make a shape's solid fill translucent (pct 0-100). No-op off real pptx."""
    if pct is None or pct >= 100:
        return
    if hasattr(shp, "set_alpha"):       # SVG-preview mock shape — honors alpha as fill-opacity
        shp.set_alpha(pct)
        return
    try:
        spPr = shp._element.spPr
        sf = spPr.find(qn("a:solidFill"))
        if sf is None:
            return
        clr = sf.find(qn("a:srgbClr"))
        if clr is None:
            clr = sf.find(qn("a:schemeClr"))
        if clr is None:
            return
        from lxml import etree
        etree.SubElement(clr, qn("a:alpha")).set("val", str(int(max(0, min(100, pct)) * 1000)))
    except Exception:
        pass


def _add_shadow(shp):
    """Add a soft outer shadow (matches the template's card shadow). Best-effort."""
    try:
        from lxml import etree
        spPr = shp._element.spPr
        eff = spPr.find(qn("a:effectLst"))
        if eff is None:
            eff = etree.SubElement(spPr, qn("a:effectLst"))
        sh = etree.SubElement(eff, qn("a:outerShdw"))
        sh.set("blurRad", "50800"); sh.set("dist", "25400")
        sh.set("dir", "5400000"); sh.set("rotWithShape", "0")
        c = etree.SubElement(sh, qn("a:srgbClr")); c.set("val", "0B2C37")
        etree.SubElement(c, qn("a:alpha")).set("val", "26000")
    except Exception:
        pass


def _fit_pt(text_s, w_emu, h_emu, base_pt, min_pt):
    """Rough shrink-to-fit: scale a font down toward min_pt so `text_s` fits a box
    (w_emu x h_emu). Approximate (no real metrics) but keeps dense content from
    overflowing readable bounds. Pagination handles count overflow; this handles
    a single long passage."""
    import math
    if not text_s:
        return base_pt
    w_in = max(0.1, w_emu / 914400.0)
    h_in = max(0.1, h_emu / 914400.0)
    pt = float(base_pt)
    chars = len(text_s)
    while pt > min_pt:
        cpl = max(1, int((w_in * 72) / (0.50 * pt)))      # chars per line
        lines = math.ceil(chars / cpl)
        if lines * (pt * 1.2 / 72) <= h_in:
            break
        pt -= 0.5
    return round(max(pt, min_pt), 1)


def _box_capacity_chars(w_emu, h_emu, pt):
    """Approx how many characters fit a w_emu x h_emu box at `pt`, using the SAME
    rough metric as `_fit_pt` (≈0.50*pt avg char width, 1.2 line height). Used to
    decide when text overflows even at the shrink-to-fit floor."""
    import math
    w_in = max(0.1, w_emu / 914400.0)
    h_in = max(0.1, h_emu / 914400.0)
    cpl = max(1, int((w_in * 72) / (0.50 * pt)))          # chars per line
    line_in = pt * 1.2 / 72
    max_lines = max(1, math.floor(h_in / line_in))
    return cpl * max_lines


def _fit_body(text_s, w_emu, h_emu, base_pt, min_pt):
    """Shrink-to-fit returning (pt, text). Scales the font toward min_pt via `_fit_pt`;
    if the text STILL overflows the box at min_pt, hard-truncates it with an ellipsis
    so it can't silently bleed past the bounds (the B6 floor-overflow failure mode).
    The common case — text that fits — is returned unchanged."""
    pt = _fit_pt(text_s, w_emu, h_emu, base_pt, min_pt)
    if not text_s:
        return pt, text_s
    cap = _box_capacity_chars(w_emu, h_emu, pt)
    if len(text_s) <= cap:
        return pt, text_s
    if cap <= 1:
        return pt, "…"
    return pt, text_s[:cap - 1].rstrip() + "…"


def _box(s, x, y, w, h, pal, tok, accent=None):
    """Draw a card per design tokens: rounded rect with radius%, outline pt+color,
    fill color+alpha, optional shadow. `tok` is a card-token dict; `accent`, if
    given, overrides the outline color. Returns the shape."""
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_radius(shp, tok.get("radiusPct", 8.33))
    fillc = _tok_color(tok.get("fillColor", "card"), pal, pal.get("card"))
    if fillc is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fillc
        _fill_alpha(shp, tok.get("fillAlpha", 100))
    opt = tok.get("outlinePt", 0)
    try:
        if opt:
            shp.line.width = Pt(opt)
            oc = accent or _tok_color(tok.get("outlineColor", "card"), pal, fillc or pal["dark"])
            if oc is not None:
                shp.line.color.rgb = oc
        else:
            shp.line.fill.background()
    except Exception:
        pass                                  # SVG mock line has no width/color
    if tok.get("shadow") and fillc is not None:
        _add_shadow(shp)
    return shp


def _apply_tx(s, transition, transition_dir, transition_speed):
    if transition and transition != "none":
        import pptx_transitions
        pptx_transitions.apply(s, transition, transition_speed, transition_dir)


def _apply_anim(s, animate, animate_speed="med", layout=None):
    """Inject entrance animations, but NEVER let an animation problem break a
    build — a malformed/unsupported case is swallowed so the deck still saves.
    `layout` lets pptx_animate replay a verified library sequence for that layout."""
    if animate and animate != "none":
        try:
            import pptx_animate
            pptx_animate.apply(s, animate, animate_speed, layout=layout)
        except Exception:
            pass


def _new_presentation():
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    return prs


def _add_slide(prs, content, base_pal, layout, images_dir=None, brand=None, bp=None,
               theme=None):
    """Render ONE slide of the given layout into `prs`; returns (slide, stats).
    Resolves the layout's blueprint (theme + full-bleed background) so the slide
    is drawn on the brand's master template substrate, not a blank canvas.
    `theme` is the optional cross-cutting dark|light override for this slide."""
    if layout not in LAYOUTS:
        raise ValueError(f"unknown slide layout: {layout!r} (have: {', '.join(LAYOUTS)})")
    pal, bg = _slide_pal(base_pal, brand, bp, layout, theme)
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _place_background(s, bg)                          # behind everything
    fill, rect, text = _drawkit(s)
    content = _resolve_images(content or {}, images_dir)
    stats = RENDERERS[layout](s, content, pal, fill, rect, text)
    return s, stats


def export_slide(content, out, brand=None, layout="infographic",
                 transition=None, transition_dir=None, transition_speed="med",
                 images_dir=None, animate=None, animate_speed="med", theme=None):
    """Render ONE standalone, on-brand slide to a .pptx. `theme` (dark|light)
    optionally overrides the layout's default theme."""
    pal = _palette_of(brand)
    bp = _load_blueprint(brand)
    prs = _new_presentation()
    s, stats = _add_slide(prs, content, pal, layout, images_dir, brand, bp, theme)
    _apply_tx(s, transition, transition_dir, transition_speed)
    _apply_anim(s, animate, animate_speed, layout=layout)
    prs.save(out)
    return stats


# ============================ NATIVE TEMPLATE RENDERING ======================
# When the brand blueprint carries a `native` mapping + a real PowerPoint
# template, the engine builds ON that template and pours content into the named
# layouts' placeholders — so the template's actual elements, styles, and image
# slots come through (not a generative approximation). Anything the template
# can't host falls back to the generative renderer on the template background.
# Engine stays brand-agnostic: a brand without a `native` blueprint never enters
# this path.

def _native_template(brand, bp):
    """Absolute path to the brand's template .pptx (from blueprint.native), or None."""
    rel = ((bp or {}).get("native") or {}).get("template")
    p = brand.asset(rel) if (rel and hasattr(brand, "asset")) else None
    return p if (p and os.path.isfile(p)) else None


def _item_count(layout, content):
    """How many repeating items a layout's content has — drives native layout
    selection (e.g. 4 cards -> the 4-block template layout)."""
    if layout == "comparison":
        return len(layouts.normalize_comparison(content).get("panels") or [])
    if layout == "timeline":
        return len(layouts.normalize_timeline(content).get("milestones") or [])
    if layout == "cards":
        return len(content.get("cards") or [])
    if layout == "bullets":
        return len(content.get("items") or [])
    return 0


def _native_spec(bp, layout, count):
    """Pick the template-layout candidate for (layout, count): a candidate whose
    `count` matches wins; a candidate with no `count` matches anything; else None
    (=> generative fallback). This is the "native for fixed, generative for
    variable" fit policy, expressed in the blueprint."""
    cands = (((bp or {}).get("native") or {}).get("layouts") or {}).get(layout)
    if not cands:
        return None
    nocount = None
    for c in cands:
        if "count" not in c:
            nocount = c
        elif c.get("count") == count:
            return c
    return nocount


def _strip_slides(prs):
    """Remove a template's example slides (keeping its masters/layouts) so the
    deck starts empty. Drops the relationship too, avoiding duplicate parts."""
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rId = sid.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        lst.remove(sid)


def _layout_by_name(prs, name):
    for m in prs.slide_masters:
        for lay in m.slide_layouts:
            if lay.name == name:
                return lay
    return None


def _set_paras(tf, paras):
    """Fill a placeholder text frame with paragraphs, INHERITING the template's
    own styling (color/size/bullets/caps come from the layout). `paras` is a list
    of paragraphs, each a list of (text, bold) runs. Bold is only set when True so
    the placeholder's default weight is otherwise preserved."""
    tf.clear()
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        for (txt, bold) in para:
            r = p.add_run()
            r.text = txt
            if bold:
                r.font.bold = True


def _norm_item(it):
    """A bullet item -> one paragraph's runs. [lead, rest] -> bold lead + rest."""
    if isinstance(it, (list, tuple)) and len(it) == 2:
        return [(layouts.strip_tags(str(it[0])), True),
                (layouts.strip_tags(" " + str(it[1])), False)]
    return [(layouts.strip_tags(str(it)), False)]


def _fill_native(slide, spec, content):
    """Pour `content` into a native template slide's placeholders per `spec`."""
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    def setp(idx, paras):
        if idx is not None and idx in phs and paras:
            _set_paras(phs[idx].text_frame, paras)

    if content.get("title"):
        setp(spec.get("title"), [[(layouts.strip_tags(content["title"]), False)]])
    if content.get("intro"):
        setp(spec.get("intro"), [[(layouts.strip_tags(content["intro"]), False)]])
    kind = spec.get("kind")
    slots = spec.get("slots") or []
    if kind == "text":
        if content.get("subtitle"):
            setp(spec.get("subtitle"), [[(layouts.strip_tags(content["subtitle"]), False)]])
    elif kind == "bullets":
        paras = [_norm_item(it) for it in (content.get("items") or [])]
        if slots:
            setp(slots[0], paras)
    elif kind == "panels":
        panels = (layouts.normalize_comparison(content).get("panels") or [])
        for slot, panel in zip(slots, panels):
            paras = []
            if panel.get("heading"):
                paras.append([(layouts.strip_tags(panel["heading"]), True)])
            if panel.get("sublabel"):
                paras.append([(layouts.strip_tags(panel["sublabel"]), False)])
            for it in (panel.get("items") or []):
                paras.append(_norm_item(it))
            if panel.get("callout"):
                paras.append([(layouts.strip_tags(panel["callout"]), True)])
            setp(slot, paras)
    elif kind == "cards":
        for slot, c in zip(slots, (content.get("cards") or [])):
            paras = []
            if c.get("title"):
                paras.append([(layouts.strip_tags(c["title"]), True)])
            if c.get("body"):
                paras.append([(layouts.strip_tags(c["body"]), False)])
            setp(slot, paras)
    elif kind == "milestones":
        ms = (layouts.normalize_timeline(content).get("milestones") or [])
        for slot, m in zip(slots, ms):
            paras = []
            if m.get("phase"):
                paras.append([(layouts.strip_tags(m["phase"]), True)])
            if m.get("title"):
                paras.append([(layouts.strip_tags(m["title"]), True)])
            if m.get("body"):
                paras.append([(layouts.strip_tags(m["body"]), False)])
            setp(slot, paras)


def _export_deck_native(slides, out, brand, bp, base_pal, images_dir,
                        transition, transition_dir, transition_speed,
                        animate, animate_speed):
    """Assemble a deck ON the brand's real template: each slide uses a native
    template layout when one fits, else the generative renderer on a template
    'Blank' layout (which carries the brand background)."""
    prs = Presentation(_native_template(brand, bp))
    _strip_slides(prs)
    blank_dark = _layout_by_name(prs, "Blank Dark")
    blank_light = _layout_by_name(prs, "Blank Light")
    used = []
    for spec0 in slides:
        spec0 = spec0 or {}
        layout = spec0.get("layout", "infographic")
        if layout not in LAYOUTS:
            layout = "infographic"
        content = _resolve_images(spec0.get("content") or {}, images_dir)
        nat = _native_spec(bp, layout, _item_count(layout, content))
        native_layout = _layout_by_name(prs, nat["layout"]) if nat else None
        if nat and native_layout is not None:
            s = prs.slides.add_slide(native_layout)
            _fill_native(s, nat, content)
            used.append("native:" + nat["layout"])
        else:                                    # generative fallback on a blank
            pal, _ = _slide_pal(base_pal, brand, bp, layout, spec0.get("theme"))
            theme = pal.get("theme", "light")
            blank = (blank_light if (theme == "light" and blank_light is not None)
                     else (blank_dark if blank_dark is not None else prs.slide_layouts[6]))
            s = prs.slides.add_slide(blank)
            pal["bg_image"] = True               # blank layout carries the bg art
            pal["logo"] = None                   # ...and its own wordmark; don't double it
            fill, rect, text = _drawkit(s)
            RENDERERS[layout](s, content, pal, fill, rect, text)
            used.append("generative:" + layout)
        _apply_tx(s, transition, transition_dir, transition_speed)
        _apply_anim(s, animate, animate_speed, layout=layout)
    prs.save(out)
    return {"layout": "deck", "slides": len(used), "render": used}


def _cont(spec, content, key, chunk, is_cont):
    """One paginated slide: same layout, content with `key` replaced by `chunk`;
    continuation slides get a ' (cont.)' title."""
    c = dict(content)
    c[key] = chunk
    if is_cont and c.get("title"):
        c = dict(c, title=str(c["title"]) + " (cont.)")
    out = {"layout": spec.get("layout"), "content": c}
    if spec.get("theme"):                         # carry the cross-cutting theme flag
        out["theme"] = spec["theme"]
    return out


def _paginate(spec, design):
    """Split one slide spec into N slides when its content exceeds the layout's
    readable capacity (overflow -> continuation slide of the same layout). Capacity
    comes from design tokens; absent design, nothing is split."""
    spec = spec or {}
    layout = spec.get("layout", "infographic")
    content = spec.get("content")
    if not isinstance(content, dict):
        return [spec]
    cap = (design or {}).get("capacity") or {}
    if layout == "cards":
        items = content.get("cards") or []
        n = int(cap.get("cards", 0) or 0)
        if n and len(items) > n:
            return [_cont(spec, content, "cards", items[i:i + n], i > 0)
                    for i in range(0, len(items), n)]
    elif layout == "timeline":
        norm = layouts.normalize_timeline(content)
        items = norm.get("milestones") or []
        n = int(cap.get("timeline", 0) or 0)
        if n and len(items) > n:
            return [_cont(spec, norm, "milestones", items[i:i + n], i > 0)
                    for i in range(0, len(items), n)]
    elif layout == "bullets":
        items = content.get("items") or []
        cols = max(1, min(2, int(content.get("columns") or (2 if len(items) > 5 else 1))))
        n = int(cap.get("bulletsPerColumn", 0) or 0) * cols
        if n and len(items) > n:
            return [_cont(spec, content, "items", items[i:i + n], i > 0)
                    for i in range(0, len(items), n)]
    # NOTE: `cycles` is deliberately NOT paginated. A cycle is a single
    # visualization (a loop); splitting it across slides ruins the graphic. When a
    # sequence is too long for a clean loop, the model should pick `process` or
    # `timeline` instead (see LAYOUT_PURPOSE) — we never break the wheel in half.
    return [spec]


def export_deck(slides, out, brand=None,
                transition=None, transition_dir=None, transition_speed="med",
                images_dir=None, animate=None, animate_speed="med"):
    """Assemble an ORDERED list of slides into ONE multi-slide .pptx deck.

    `slides` is a list of {"layout": <name>, "content": {...}}. The same
    transition (if any) is applied to every slide, like a real deck; `animate`
    (none|fade|rise|flyleft|flyright) cascades each slide's elements in on entry.
    `images_dir` resolves bare `image` filenames in a slide's content."""
    if not slides:
        raise ValueError("deck has no slides")
    pal = _palette_of(brand)
    bp = _load_blueprint(brand)
    design = (bp or {}).get("design")
    # PRIMARY path = token-driven generative (template look via design tokens +
    # background, adaptive + paginated). The native fixed-layout build is an
    # opt-in fallback, used only when a brand has a `native` map but NO `design`.
    if (not design) and _native_template(brand, bp) and ((bp.get("native") or {}).get("layouts")):
        return _export_deck_native(slides, out, brand, bp, pal, images_dir,
                                   transition, transition_dir, transition_speed,
                                   animate, animate_speed)
    expanded = []
    for spec in slides:
        expanded.extend(_paginate(spec, design or {}))
    prs = _new_presentation()
    used = []
    for spec in expanded:
        layout = (spec or {}).get("layout", "infographic")
        s, _ = _add_slide(prs, (spec or {}).get("content"), pal, layout, images_dir, brand, bp,
                          (spec or {}).get("theme"))
        _apply_tx(s, transition, transition_dir, transition_speed)
        _apply_anim(s, animate, animate_speed, layout=layout)
        used.append(layout)
    prs.save(out)
    return {"layout": "deck", "slides": len(used), "layouts": used}


def _render_process(s, content, pal, fill, rect, text):
    """PROCESS layout — a horizontal numbered-step poster.

    Header band (title + subtitle) | optional intro line | a row of 3-6
    connected step cards (top accent strip · numbered circle · title · body)
    joined by chevrons | optional footer band. Colors come from the brand.

    Content schema:
    {
      "title": "...", "subtitle": "...", "intro": "...",
      "steps": [{"num": "1", "title": "...", "body": "...",
                 "accent": "primary|secondary|tertiary|dark"}, ...],  # 3-6
      "footer": "..."
    }
    """
    cycle = ["primary", "secondary", "tertiary", "dark"]
    intro = content.get("intro", "")
    steps = (content.get("steps") or [])[:6]   # cap at 6 to stay on-slide
    footer = content.get("footer")

    # ---- header band (shared) ----
    _header(s, content, pal, rect, text)

    # ---- optional intro ----
    card_top = Inches(1.6)
    if intro:
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.4),
             [[(intro, 12.5, pal["muted"], False)]])
        card_top = Inches(2.05)

    # ---- step cards row ----
    n = max(1, len(steps))
    margin, gap = Inches(0.55), Inches(0.34)
    avail = W - 2 * margin
    cw = int((avail - (n - 1) * gap) / n)
    card_bottom = Inches(6.85) if not footer else Inches(6.55)
    card_h = int(card_bottom - card_top)
    num_d = Inches(0.92)
    x = margin
    for i, st in enumerate(steps):
        accent = _accent(pal, st.get("accent", cycle[i % 4]))
        rect(x, card_top, cw, card_h, pal["card"], rounded=True)
        rect(x, card_top, cw, Inches(0.1), accent, rounded=True)  # top accent strip
        cx = x + int((cw - num_d) / 2)
        ny = card_top + Inches(0.4)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, ny, num_d, num_d)
        fill(circ, accent)
        text(cx, ny, num_d, num_d, [[(str(st.get("num", i + 1)), 30, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(x + Inches(0.18), card_top + Inches(1.55), cw - Inches(0.36), Inches(0.8),
             [[(st.get("title", ""), 15, pal["cardInk"], True)]], align=PP_ALIGN.CENTER)
        text(x + Inches(0.18), card_top + Inches(2.35), cw - Inches(0.36),
             card_h - int(Inches(2.55)),
             [[(st.get("body", ""), 11.5, pal["muted"], False)]],
             align=PP_ALIGN.CENTER, line=1.05)
        if i < n - 1:  # connector chevron in the gap
            chd = Inches(0.3)
            chx = x + cw + int((gap - chd) / 2)
            chy = card_top + int(card_h / 2) - int(chd / 2)
            chev = s.shapes.add_shape(MSO_SHAPE.CHEVRON, chx, chy, chd, chd)
            fill(chev, pal["primary"])
        x = x + cw + gap

    # ---- footer band ----
    if footer:
        rect(0, Inches(7.06), W, Inches(0.44), pal["primary"])
        text(0, Inches(7.06), _footer_tw(pal), Inches(0.44), [[(footer, 12, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _place_band_logo(s, pal)

    return {"layout": "process", "steps": len(steps)}


def _place_header_logo(s, path):
    """Draw the brand logo (a white wordmark, for dark bands) in the top-right of
    the header band. Returns the placed width (EMU), or None if it couldn't be
    placed — missing file, a non-raster asset (add_picture needs a raster, so an
    SVG logo is skipped), or an unreadable image. Never raises; tagged `bg.logo`
    so the animation engine treats it as a background mark."""
    if not path or not os.path.isfile(path):
        return None
    if os.path.splitext(path)[1].lower() not in (".png", ".jpg", ".jpeg"):
        return None
    iw, ih = _img_size(path)
    if not (iw and ih):
        return None
    h = Inches(0.4)
    w = int(h * iw / ih)
    x = int(W - Inches(0.55) - w)
    y = int((Inches(1.18) - h) / 2)
    try:
        pic = s.shapes.add_picture(path, x, y, w, h)
        pic.name = "bg.logo"
        return w
    except Exception:
        return None


def _header(s, content, pal, rect, text):
    """Shared slide header. Mode comes from the brand blueprint (pal["header_mode"]):
      band — a solid dark band + accent rule + white title (the agnostic default,
             i.e. the "previous templates" look)
      rule — title in theme ink over the brand background + a thin accent rule
             under it + the wordmark bottom-right (the official template look)
      none — caller draws its own full-bleed chrome (divider/closing/quote)
    Returns the y-offset (EMU) where body content can begin (stable across modes
    so existing body geometry is unaffected)."""
    mode = pal.get("header_mode", "band")
    if mode == "none":
        return Inches(1.55)
    title_w = Inches(12.2)
    if mode == "band":
        rect(0, 0, W, Inches(1.18), pal["dark"], role="bg.band")
        rect(0, Inches(1.18), W, Inches(0.06), pal["primary"], role="bg.rule")
        title_color, sub_color = pal["white"], pal["primary"]
        logo_w = _place_header_logo(s, pal.get("logo"))
        if logo_w:                   # keep the title clear of the logo
            title_w = max(Inches(4.0), int(W - Inches(0.55) - logo_w - Inches(0.35)))
    else:                            # rule (template look) — no band
        title_color = pal.get("title_ink", pal["dark"])
        sub_color = pal.get("rule_color", pal["primary"])
    if content.get("title"):
        text(Inches(0.55), Inches(0.14), title_w, Inches(0.6),
             [[(content["title"], _dtok(pal, "type", "titleSizePt", 28), title_color, True)]],
             role="content.title")
    if content.get("subtitle"):
        text(Inches(0.57), Inches(0.74), Inches(12.2), Inches(0.35),
             [[(content["subtitle"], 14, sub_color, False)]], role="content.subtitle")
    if mode == "rule":
        rect(Inches(0.55), Inches(1.16), Inches(12.23), Inches(0.03),
             pal.get("rule_color", pal["primary"]), role="bg.rule")
        # If this slide carries a footer band, the wordmark lives INSIDE it
        # (white, centered) — but it's placed by the footer step so it lands ABOVE
        # the band (see _place_band_logo). Here we only place the bottom-right
        # clear-space logo for the no-band case (theme-colored, on the bg).
        if not (content.get("footer") and pal.get("footer_logo")):
            _place_footer_logo(s, pal.get("logo"))
    return Inches(1.55)


def _place_footer_logo(s, path, in_band=False):
    """Place the brand wordmark bottom-right. By default at the official template's
    clear-space position (~11.03in, 6.78in, 1.49in wide). With in_band=True the
    wordmark is vertically centered INSIDE the footer band (y 7.06..7.50, the
    green primary bar) — the white-on-dark variant reads against the band.
    Raster-only; never raises."""
    if not path or not os.path.isfile(path):
        return
    if os.path.splitext(path)[1].lower() not in (".png", ".jpg", ".jpeg"):
        return
    iw, ih = _img_size(path)
    if not (iw and ih):
        return
    w = Inches(1.49)
    h = int(w * ih / iw)
    if in_band:
        y = int(Inches(7.06) + (Inches(0.44) - h) / 2)   # centered in the band
    else:
        y = int(Inches(6.78))
    try:
        pic = s.shapes.add_picture(path, int(Inches(11.03)), y, int(w), h)
        pic.name = "bg.logo"
    except Exception:
        pass


def _footer_tw(pal):
    """Width for footer-band text. Narrowed to clear the in-band wordmark when one
    is present (the template look), so the text never runs under the logo."""
    return Inches(10.6) if (pal.get("footer_logo") and pal.get("header_mode") == "rule") else W


def _place_band_logo(s, pal):
    """Place the white wordmark INTO the footer band. Called right AFTER the band
    is drawn so the logo sits in the top z-layer (PPTX z-order = creation order);
    placing it in _header instead would leave it behind the later-drawn band.
    No-op unless the template-look in-band logo applies (brand has a footer_logo
    and is in rule mode)."""
    if pal.get("footer_logo") and pal.get("header_mode") == "rule":
        _place_footer_logo(s, pal["footer_logo"], in_band=True)


def _footer(s, footer, pal, rect, text):
    if footer:
        rect(0, Inches(7.06), W, Inches(0.44), pal["primary"], role="bg.footer")
        text(0, Inches(7.06), _footer_tw(pal), Inches(0.44), [[(footer, 12, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, role="content.footer")
        _place_band_logo(s, pal)


def _render_comparison(s, content, pal, fill, rect, text):
    """COMPARISON layout — 2 or 3 side-by-side panels.

    Header band | optional intro | 2-3 panels, each: accent header bar (heading)
    · optional sublabel · bulleted items · optional callout box | optional footer.
    Best for old-vs-new, option A/B/C. Colors come from the brand.

    Content schema:
    {
      "title","subtitle","intro",
      "columns": [                                   # 2 or 3
        {"heading":"...","sublabel":"...","accent":"primary|secondary|tertiary|dark",
         "items": [["bold lead"," rest"], "plain line", ...],   # up to 6
         "callout":"..."}, ...
      ],
      "footer":"..."
    }
    """
    # same accent order as the course block (render.py _ACCENT_CYCLE) so a panel
    # gets the same color on either surface.
    ct = (pal.get("design") or {}).get("card", {})
    cap = ((pal.get("design") or {}).get("capacity") or {}).get("comparisonItems", 6)
    cycle = ["primary", "secondary", "tertiary", "dark"]
    content = layouts.normalize_comparison(content)   # accept course `panels` too
    top = _header(s, content, pal, rect, text)
    if content.get("intro"):
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.4),
             [[(content["intro"], _dtok(pal, "type", "introSizePt", 12.5), pal["muted"], False)]])
        top = Inches(2.0)
    cols = (content.get("panels") or [])[:3]
    n = max(1, len(cols))
    footer = content.get("footer")
    margin, gap = Inches(0.55), Inches(0.4)
    avail = W - 2 * margin
    pw = int((avail - (n - 1) * gap) / n)
    p_bottom = Inches(6.85) if not footer else Inches(6.55)
    p_h = int(p_bottom - top)
    x = margin
    for i, col in enumerate(cols):
        accent = _accent(pal, col.get("accent", cycle[i % 4]))
        _box(s, x, top, pw, p_h, pal, ct)                            # panel (card token)
        rect(x, top, pw, Inches(0.62), accent, rounded=True)         # accent header bar
        # tags stripped: a course panel may carry inline HTML (<strong> etc.) the
        # slide canvas can't render.
        text(x + Inches(0.1), top, pw - Inches(0.2), Inches(0.62),
             [[(layouts.strip_tags(col.get("heading", "")), 15, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        iy = top + Inches(0.78)
        if col.get("sublabel"):
            text(x + Inches(0.22), iy, pw - Inches(0.44), Inches(0.3),
                 [[(layouts.strip_tags(col["sublabel"]), 10.5, pal["muted"], True)]])
            iy = iy + Inches(0.36)
        for item in (col.get("items") or [])[:cap]:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22),
                                     iy + Inches(0.06), Inches(0.12), Inches(0.12))
            fill(dot, accent)
            if isinstance(item, (list, tuple)) and len(item) == 2:
                runs = [[(layouts.strip_tags(item[0]), 11, pal["cardInk"], True),
                         (layouts.strip_tags(item[1]), 11, pal["muted"], False)]]
            else:
                runs = [[(layouts.strip_tags(str(item)), 11, pal["muted"], False)]]
            text(x + Inches(0.46), iy, pw - Inches(0.6), Inches(0.5), runs, line=1.02)
            iy = iy + Inches(0.46)
        if col.get("callout"):
            cy = top + p_h - int(Inches(0.95))
            rect(x + Inches(0.18), cy, pw - Inches(0.36), Inches(0.8), pal["tint"], rounded=True)
            text(x + Inches(0.34), cy, pw - Inches(0.68), Inches(0.8),
                 [[(layouts.strip_tags(col["callout"]), 11, pal["dark"], True)]],
                 anchor=MSO_ANCHOR.MIDDLE, line=1.03)
        x = x + pw + gap
    _footer(s, footer, pal, rect, text)
    return {"layout": "comparison", "columns": len(cols)}


def _render_timeline(s, content, pal, fill, rect, text):
    """TIMELINE / ROADMAP layout — a horizontal axis with 3-6 milestones that
    alternate above and below the line.

    Header band | optional intro | axis line + numbered nodes; each milestone
    card shows a phase chip · title · body | optional footer.

    Content schema:
    {
      "title","subtitle","intro",
      "milestones": [                                # 3-6
        {"phase":"NOW / Q1","title":"...","body":"...",
         "accent":"primary|secondary|tertiary|dark"}, ...
      ],
      "footer":"..."
    }
    """
    tl = (pal.get("design") or {}).get("timeline", {})
    content = layouts.normalize_timeline(content)   # accept course `html` as well as slide `body`
    _header(s, content, pal, rect, text)
    if content.get("intro"):
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.4),
             [[(content["intro"], _dtok(pal, "type", "introSizePt", 12.5), pal["muted"], False)]])
    ms = content.get("milestones") or []            # count overflow already paginated
    n = max(1, len(ms))
    footer = content.get("footer")
    axis_y = Inches(4.05)
    cardw = Inches(2.2)
    half = int(cardw / 2)
    pad = Inches(1.55)                       # keep extreme cards on-slide
    first, last = pad, W - pad
    span = int(last - first)
    node_d = Inches(tl.get("nodeDiaIn", 0.34))
    node_col = _tok_color(tl.get("nodeColor", "node"), pal, pal["node"])
    rail_col = _tok_color(tl.get("axisColor", "rail"), pal, pal["rail"])
    num_col = _tok_color(tl.get("numberColor", "dark"), pal, pal["dark"])
    phase_col = _tok_color(tl.get("phaseColor", "node"), pal, pal["node"])
    title_col = _tok_color(tl.get("titleColor", "ink"), pal, pal["ink"])
    body_col = _tok_color(tl.get("bodyColor", "muted"), pal, pal["muted"])
    minb = tl.get("minBodyPt", 8.5)
    # design-token milestone card: border-only (mint outline), template radius
    tlcard = {"radiusPct": tl.get("cardRadiusPct", 9.6), "outlinePt": tl.get("cardOutlinePt", 3),
              "outlineColor": tl.get("cardOutlineColor", "node"),
              "fillColor": tl.get("cardFillColor", "none"),
              "fillAlpha": tl.get("cardFillAlpha", 0), "shadow": False}
    axis_pt = tl.get("axisPt", 1.5)
    rect(first, axis_y - Pt(axis_pt / 2), span, Pt(axis_pt), rail_col, role="bg.rail")
    for i, m in enumerate(ms):
        u = i + 1                                # 1-based unit index for role names
        cx = first if n == 1 else first + int(span * i / (n - 1))
        above = (i % 2 == 0)
        if above:
            rect(cx - Inches(0.012), Inches(3.55), Inches(0.024), int(axis_y - Inches(3.55)),
                 node_col, role=f"bg.connector.{u}")
        else:
            rect(cx - Inches(0.012), axis_y, Inches(0.024), Inches(0.5), node_col,
                 role=f"bg.connector.{u}")
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - int(node_d / 2),
                                  axis_y - int(node_d / 2), node_d, node_d)
        fill(node, node_col); role(node, f"focal.node.{u}")
        text(cx - int(node_d / 2), axis_y - int(node_d / 2), node_d, node_d,
             [[(str(i + 1), 13, num_col, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, role=f"content.number.{u}")
        # card (clamped within margins)
        cardx = max(int(Inches(0.4)), min(cx - half, int(W - Inches(0.4) - cardw)))
        cardy = Inches(1.7) if above else Inches(4.55)
        card_h = Inches(1.85)
        role(_box(s, cardx, cardy, cardw, card_h, pal, tlcard), f"focal.card.{u}")
        cpad = Inches(0.16)
        yy = cardy + Inches(0.12)
        if m.get("phase"):
            text(cardx + cpad, yy, cardw - 2 * cpad, Inches(0.28),
                 [[(layouts.strip_tags(m["phase"]), 10, phase_col, True)]], role=f"content.phase.{u}")
            yy = yy + Inches(0.3)
        text(cardx + cpad, yy, cardw - 2 * cpad, Inches(0.38),
             [[(layouts.strip_tags(m.get("title", "")), 12.5, title_col, True)]], role=f"content.mtitle.{u}")
        yy = yy + Inches(0.4)
        body = layouts.strip_tags(m.get("body", ""))
        bh = int(cardy + card_h - Inches(0.12) - yy)
        bpt, body = _fit_body(body, cardw - 2 * cpad, bh, 10, minb)
        text(cardx + cpad, yy, cardw - 2 * cpad, bh,
             [[(body, bpt, body_col, False)]], line=1.0, role=f"content.mbody.{u}")
    _footer(s, footer, pal, rect, text)
    return {"layout": "timeline", "milestones": len(ms)}


def _render_divider(s, content, pal, fill, rect, text):
    """DIVIDER / SECTION-TITLE layout — a full-bleed branded title screen.

    Full-bleed background (brand role) · optional kicker label + accent rule ·
    large centered title · optional subtitle · optional footer band. Use as a
    section break or a course/deck title slide.

    Content schema:
    {
      "kicker":"SECTION 2", "title":"...", "subtitle":"...",
      "bg":"dark|primary|secondary|tertiary",   # background role (default dark)
      "footer":"..."
    }
    """
    if not pal.get("bg_image"):                       # brand art is the bg if present
        bg = _accent(pal, content.get("bg", "dark"), "dark")
        rect(0, 0, W, H, bg)                          # full bleed
    ink = pal.get("title_ink", pal["white"])
    cx, cw = Inches(1.0), Inches(11.33)
    if content.get("kicker"):
        text(cx, Inches(2.55), cw, Inches(0.4),
             [[(content["kicker"], _dtok(pal, "type", "kickerSizePt", 16), pal["rule_color"], True)]],
             align=PP_ALIGN.CENTER)
        rect(int(W / 2) - Inches(0.6), Inches(3.05), Inches(1.2), Inches(0.05), pal["rule_color"])
    text(cx, Inches(3.25), cw, Inches(1.4),
         [[(content.get("title", ""), _dtok(pal, "type", "dividerTitlePt", 44), ink, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if content.get("subtitle"):
        text(cx, Inches(4.75), cw, Inches(0.8),
             [[(content["subtitle"], _dtok(pal, "type", "dividerSubPt", 18),
                pal.get("muted", pal["tint"]), False)]],
             align=PP_ALIGN.CENTER, line=1.1)
    _footer(s, content.get("footer"), pal, rect, text)
    return {"layout": "divider"}


def _render_chart(s, content, pal, fill, rect, text):
    """CHART layout — header band (title + subtitle) | a NATIVE, editable PowerPoint
    chart (column / stacked / line / pie) brand-colored | optional source footnote.

    Unlike the other layouts (which draw shapes), this inserts a real PowerPoint
    chart object, so the figures stay editable in PowerPoint. Colors come from the
    brand palette (series = primary/secondary/tertiary/dark, cycling).

    Content schema:
    {
      "title": "...", "subtitle": "...",
      "chart": "bar|line|pie|stackedBar|groupedBar",
      "categories": ["Q1", "Q2", ...],
      "series": [{"name": "Admits", "data": [120, 145, ...]}, ...],
      "xLabel": "...", "yLabel": "...", "source": "..."
    }
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

    # ---- header band (shared) ----
    _header(s, content, pal, rect, text)

    ctype = content.get("chart", "bar")
    cats = content.get("categories") or []
    series = [sr for sr in (content.get("series") or []) if isinstance(sr, dict)]
    if not series:
        text(Inches(0.6), Inches(3.4), Inches(12.13), Inches(0.6),
             [[("No chart data provided.", 14, pal["muted"], False)]], align=PP_ALIGN.CENTER)
        return {"layout": "chart", "chart": ctype, "series": 0}

    cd = CategoryChartData()
    cd.categories = cats or [""]
    for i, sr in enumerate(series):
        data = list(sr.get("data") or [])
        if cats:                                  # align each series to the category count
            data = (data + [None] * len(cats))[:len(cats)]
        cd.add_series(sr.get("name") or f"Series {i + 1}", data)

    xl = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
          "groupedBar": XL_CHART_TYPE.COLUMN_CLUSTERED,
          "stackedBar": XL_CHART_TYPE.COLUMN_STACKED,
          "line": XL_CHART_TYPE.LINE_MARKERS,
          "pie": XL_CHART_TYPE.PIE}.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)

    gf = s.shapes.add_chart(xl, Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.2), cd)
    chart = gf.chart
    chart.has_title = False
    order = [pal["primary"], pal["secondary"], pal["tertiary"], pal["dark"]]
    plot = chart.plots[0]

    if ctype == "pie":
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        for i, pt in enumerate(plot.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = order[i % len(order)]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = True
        dl.show_value = False
        dl.number_format = "0%"
        dl.number_format_is_linked = False
        try:
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except (ValueError, KeyError):
            pass
    else:
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        for i, ser in enumerate(plot.series):
            c = order[i % len(order)]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = c
            try:                                  # line charts: also color the line
                ser.format.line.color.rgb = c
            except Exception:
                pass
        # axis titles
        try:
            if content.get("yLabel"):
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = content["yLabel"]
            if content.get("xLabel"):
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = content["xLabel"]
        except Exception:
            pass

    # ---- source footnote (the provenance line; carried through from the block) ----
    if content.get("source"):
        text(Inches(0.6), Inches(6.92), Inches(12.13), Inches(0.4),
             [[("Source: ", 10, pal["ink"], True), (content["source"], 10, pal["muted"], False)]])
    return {"layout": "chart", "chart": ctype, "series": len(series)}


def _render_image(s, content, pal, fill, rect, text):
    """IMAGE layout — header band (title + subtitle) | one large framed image |
    optional caption | optional footer. For a hero shot, screenshot, or diagram.

    Content schema:
    { "title","subtitle", "image": "file.png", "caption":"...", "footer":"..." }
    """
    mode = (content.get("mode") or "hero").lower()
    tok = (pal.get("design") or {}).get("image") or {}
    footer = content.get("footer")
    caption = content.get("caption")
    cap_pt = tok.get("captionSizePt", 12)
    cap_color = _tok_color(tok.get("captionColor", "muted"), pal, pal["grey"])

    # ---- full-bleed: image behind everything, title overlaid on a scrim ----
    if mode == "full" and content.get("image"):
        _place_image(s, content.get("image"), 0, 0, W, H, pal, rect, text, fit="cover")
        scrim = rect(0, 0, W, H, _tok_color(tok.get("overlayColor", "dark"), pal, pal["dark"]),
                     role="bg.scrim")
        try:
            _fill_alpha(scrim, tok.get("overlayAlpha", 42))
        except Exception:
            pass
        if content.get("title"):
            text(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.7),
                 [[(content["title"], _dtok(pal, "type", "dividerTitlePt", 48), pal["white"], True)]],
                 role="content.title")
        if content.get("subtitle"):
            text(Inches(0.85), Inches(4.45), Inches(11.6), Inches(0.7),
                 [[(content["subtitle"], _dtok(pal, "type", "dividerSubPt", 18), pal["white"], False)]],
                 role="content.subtitle")
        _footer(s, footer, pal, rect, text)
        return {"layout": "image", "image": True, "mode": "full"}

    _header(s, content, pal, rect, text)

    # ---- banner: a wide image band under the header, content/caption below ----
    if mode == "banner" and content.get("image"):
        top = Inches(1.5)
        band_h = Inches(tok.get("bannerIn", 2.7))
        _place_image(s, content.get("image"), 0, top, W, int(band_h), pal, rect, text, fit="cover")
        cy = top + band_h + Inches(0.24)
        if content.get("intro"):
            text(Inches(0.8), cy, Inches(11.73), Inches(1.6),
                 [[(content["intro"], tok.get("introSizePt", 16), pal["ink"], False)]], line=1.18)
            cy = cy + Inches(1.7)
        if caption:
            text(Inches(0.8), cy, Inches(11.73), Inches(0.4),
                 [[(caption, cap_pt, cap_color, False)]], align=PP_ALIGN.CENTER)
        _footer(s, footer, pal, rect, text)
        return {"layout": "image", "image": True, "mode": "banner"}

    # ---- hero (default): one large edge-to-edge image filling the body ----
    top = Inches(1.5)
    bottom = Inches(6.55) if footer else Inches(7.3)
    cap_h = Inches(0.5) if caption else Inches(0)
    _place_image(s, content.get("image"), 0, top, W,
                 int(bottom - top - cap_h), pal, rect, text,
                 fit=content.get("fit", "cover"))
    if caption:
        text(Inches(0.8), bottom - cap_h + Inches(0.06), Inches(11.73), Inches(0.4),
             [[(caption, cap_pt, cap_color, False)]], align=PP_ALIGN.CENTER)
    _footer(s, footer, pal, rect, text)
    return {"layout": "image", "image": bool(content.get("image")), "mode": "hero"}


def _render_imagetext(s, content, pal, fill, rect, text):
    """IMAGE + TEXT layout — header band | a framed image on one side, a text
    column (intro + bulleted items + optional callout) on the other. Mirrors the
    course `imageText` block. `side` puts the image left (default) or right.

    Content schema:
    { "title","subtitle", "image":"file.png", "side":"left|right",
      "intro":"...", "items": [["bold lead"," rest"], "plain line", ...],
      "callout":"...", "footer":"..." }
    """
    _header(s, content, pal, rect, text)
    footer = content.get("footer")
    tok = (pal.get("design") or {}).get("image") or {}
    top = Inches(1.55)
    bottom = Inches(6.55) if footer else Inches(7.1)        # text-column bottom
    img_bottom = Inches(7.06) if footer else H              # image runs to the slide edge
    img_w = Inches(tok.get("insetIn", 6.0))
    gap = Inches(0.5)
    if (content.get("side") or "left").lower() == "right":
        img_x = int(W - img_w)
        text_x = Inches(0.7)
        colw = int(img_x - text_x - gap)
    else:
        img_x = 0
        text_x = int(img_w + gap)
        colw = int(W - text_x - Inches(0.7))
    _place_image(s, content.get("image"), img_x, top, int(img_w), int(img_bottom - top),
                 pal, rect, text, fit=content.get("fit", "cover"))

    # ---- text column ----
    tx, tw = text_x, colw
    intro_pt = tok.get("introSizePt", 16)
    body_pt = tok.get("bodySizePt", 14)
    callout_pt = tok.get("calloutSizePt", 14)
    cy = top + Inches(0.05)
    if content.get("intro"):
        text(tx, cy, tw, Inches(1.0), [[(content["intro"], intro_pt, pal["ink"], False)]], line=1.14)
        cy = cy + Inches(1.05)
    row_h = Inches(0.56)
    for item in (content.get("items") or [])[:7]:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, tx, cy + Inches(0.09),
                                 Inches(0.14), Inches(0.14))
        fill(dot, pal["primary"])
        if isinstance(item, (list, tuple)) and len(item) == 2:
            runs = [[(layouts.strip_tags(item[0]), body_pt, pal["ink"], True),
                     (layouts.strip_tags(item[1]), body_pt, pal["muted"], False)]]
        else:
            runs = [[(layouts.strip_tags(str(item)), body_pt, pal["muted"], False)]]
        text(tx + Inches(0.32), cy, tw - Inches(0.32), row_h, runs, line=1.06)
        cy = cy + row_h
    if content.get("callout"):
        cy = max(cy + Inches(0.08), bottom - Inches(1.05))
        rect(tx, cy, tw, Inches(1.0), pal["tint"], rounded=True)
        rect(tx, cy, Inches(0.09), Inches(1.0), pal["primary"])
        text(tx + Inches(0.26), cy, tw - Inches(0.4), Inches(1.0),
             [[(content["callout"], callout_pt, pal["dark"], True)]],
             anchor=MSO_ANCHOR.MIDDLE, line=1.06)
    _footer(s, footer, pal, rect, text)
    return {"layout": "imagetext", "image": bool(content.get("image")),
            "items": len(content.get("items") or [])}


def _render_cards(s, content, pal, fill, rect, text):
    """CARDS layout — header band | optional intro | a grid of equal cards
    (2–4 columns), each with a top accent strip, a title, and a body line.
    For a set of parallel items that aren't a sequence (a process) or a contrast
    (a comparison): components, options, features, the kit.

    Content schema:
    { "title","subtitle","intro", "columns": 2|3|4,
      "cards": [{"title":"...","body":"...","accent":"primary|..."}, ...],
      "footer":"..." }
    """
    ct = (pal.get("design") or {}).get("card", {})
    cycle = ["primary", "secondary", "tertiary", "dark"]
    top = _header(s, content, pal, rect, text)
    if content.get("intro"):
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.4),
             [[(content["intro"], _dtok(pal, "type", "introSizePt", 12.5), pal["muted"], False)]])
        top = Inches(2.0)
    cards = content.get("cards") or []          # count overflow already paginated
    n = max(1, len(cards))
    cols = int(content.get("columns") or (2 if n <= 4 else (3 if n <= 6 else 4)))
    cols = max(1, min(4, cols))
    rows = (n + cols - 1) // cols
    footer = content.get("footer")
    margin, gap = Inches(0.55), Inches(ct.get("gapIn", 0.34))
    cw = int((W - 2 * margin - (cols - 1) * gap) / cols)
    bottom = Inches(6.85) if not footer else Inches(6.55)
    ch = int((int(bottom - top) - (rows - 1) * gap) / max(1, rows))
    pad = Inches(ct.get("padIn", 0.24))
    tsize, bsize = ct.get("titleSizePt", 15), ct.get("bodySizePt", 11.5)
    minb = ct.get("minBodyPt", 9)
    tcol = _tok_color(ct.get("titleColor", "cardInk"), pal, pal["cardInk"])
    bcol = _tok_color(ct.get("bodyColor", "muted"), pal, pal["muted"])
    strip = ct.get("accentStripPt", 0)
    for i, c in enumerate(cards):
        r, q = divmod(i, cols)
        x = margin + q * (cw + gap)
        y = top + r * (ch + gap)
        accent = _accent(pal, c.get("accent", cycle[i % 4]))
        _box(s, x, y, cw, ch, pal, ct)
        if strip:
            rect(x, y, cw, Pt(strip), accent, rounded=True)
        title_h = Inches(0.52)
        text(x + pad, y + pad, cw - 2 * pad, title_h,
             [[(layouts.strip_tags(c.get("title", "")), tsize, tcol, True)]])
        body = layouts.strip_tags(c.get("body", ""))
        by = y + pad + title_h
        bh = int(y + ch - pad - by)
        bpt, body = _fit_body(body, cw - 2 * pad, bh, bsize, minb)
        text(x + pad, by, cw - 2 * pad, bh, [[(body, bpt, bcol, False)]], line=1.08)
    _footer(s, footer, pal, rect, text)
    return {"layout": "cards", "cards": len(cards), "columns": cols}


def _render_quote(s, content, pal, fill, rect, text):
    """QUOTE layout — a full-bleed branded pull-quote: a big quotation mark, the
    quote in large type, and an attribution rule + byline.

    Content schema: { "quote":"...", "by":"...", "bg":"dark|primary|secondary|tertiary" }
    """
    if not pal.get("bg_image"):                       # brand art is the bg if present
        bg = _accent(pal, content.get("bg", "dark"), "dark")
        rect(0, 0, W, H, bg)
    text(Inches(0.9), Inches(1.05), Inches(3.0), Inches(2.0),
         [[("“", 120, pal["primary"], True)]])               # oversized open-quote
    quote = content.get("quote") or content.get("title") or ""
    text(Inches(1.1), Inches(2.5), Inches(11.13), Inches(2.6),
         [[(quote, 30, pal.get("title_ink", pal["white"]), True)]],
         anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    by = content.get("by")
    if by:
        rect(Inches(1.15), Inches(5.35), Inches(0.7), Inches(0.06), pal["primary"])
        text(Inches(1.1), Inches(5.5), Inches(11.13), Inches(0.5),
             [[("— " + by, 15, pal.get("muted", pal["tint"]), False)]])
    return {"layout": "quote"}


def _render_statement(s, content, pal, fill, rect, text):
    """STATEMENT / BIG-STAT layout — a full-bleed branded slide built around ONE
    line: an optional kicker, an optional huge `value` (a metric), a large
    statement, and an optional supporting subtitle. For an impact line or a KPI.

    Content schema:
    { "kicker":"...", "value":"6×", "title":"the statement",
      "subtitle":"supporting line", "bg":"primary|dark|secondary|tertiary|white" }
    """
    bgname = content.get("bg", "primary")
    if pal.get("bg_image"):
        # themed brand: the full-bleed art IS the theme's background, so the text
        # follows the ACTIVE theme (a light-theme override then reads correctly).
        ink = pal.get("title_ink", pal["white"])
        accent_ink = pal.get("rule_color", pal["primary"])
        sub_col = pal.get("muted", pal["tint"])
    else:                                             # agnostic: bg comes from content
        light = bgname in ("white", "light", "tint")
        bg = pal["white"] if bgname == "white" else (pal["tint"] if bgname in ("light", "tint")
                                                     else _accent(pal, bgname, "primary"))
        rect(0, 0, W, H, bg)
        ink = pal["dark"] if light else pal["white"]
        accent_ink = pal["primary"] if light else pal["white"]
        sub_col = pal["grey"] if light else pal["tint"]
    cx, cw = Inches(1.2), Inches(10.93)
    if content.get("kicker"):
        text(cx, Inches(1.7), cw, Inches(0.5),
             [[(content["kicker"], 16, accent_ink, True)]], align=PP_ALIGN.CENTER)
    has_val = bool(content.get("value"))
    if has_val:
        text(cx, Inches(2.15), cw, Inches(1.8),
             [[(str(content["value"]), 96, ink, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    title = content.get("title") or content.get("statement") or ""
    ty = Inches(4.05) if has_val else Inches(2.9)
    th = Inches(1.6) if has_val else Inches(1.9)
    text(cx, ty, cw, th, [[(title, 30 if has_val else 36, ink, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.12)
    if content.get("subtitle"):
        text(cx, ty + th, cw, Inches(0.9),
             [[(content["subtitle"], 16, sub_col, False)]],
             align=PP_ALIGN.CENTER, line=1.1)
    return {"layout": "statement"}


def _render_bullets(s, content, pal, fill, rect, text):
    """BULLETS / AGENDA layout — header band | optional lead | a clean list of
    accent-dotted bullets in one or two columns. For an agenda, takeaways, or a
    simple list that doesn't need cards.

    Content schema:
    { "title","subtitle","intro", "columns": 1|2,
      "items": [["bold lead"," rest"], "plain line", ...], "footer":"..." }
    """
    bsz = _dtok(pal, "type", "bodySizePt", 14)
    top = _header(s, content, pal, rect, text)
    if content.get("intro"):
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.5),
             [[(content["intro"], _dtok(pal, "type", "introSizePt", 13), pal["muted"], False)]])
        top = Inches(2.05)
    items = content.get("items") or []
    footer = content.get("footer")
    cols = max(1, min(2, int(content.get("columns") or (2 if len(items) > 5 else 1))))
    margin, gap = Inches(0.7), Inches(0.6)
    colw = int((W - 2 * margin - (cols - 1) * gap) / cols)
    bottom = Inches(6.85) if not footer else Inches(6.55)
    per = (len(items) + cols - 1) // cols if items else 0
    row_h = Inches(0.62)
    for i, item in enumerate(items):
        c = i // per if per else 0
        r = i % per if per else i
        x = margin + c * (colw + gap)
        y = top + Inches(0.1) + r * row_h
        if y > bottom - Inches(0.3):
            continue
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.09),
                                 Inches(0.16), Inches(0.16))
        fill(dot, _accent(pal, ["primary", "secondary", "tertiary", "dark"][i % 4]))
        if isinstance(item, (list, tuple)) and len(item) == 2:
            runs = [[(layouts.strip_tags(item[0]), bsz, pal["ink"], True),
                     (layouts.strip_tags(item[1]), bsz, pal["muted"], False)]]
        else:
            runs = [[(layouts.strip_tags(str(item)), bsz, pal["ink"], False)]]
        text(x + Inches(0.34), y, colw - Inches(0.34), row_h, runs, line=1.05)
    _footer(s, footer, pal, rect, text)
    return {"layout": "bullets", "items": len(items), "columns": cols}


def _agenda_items(content):
    """Normalize agenda entries to [{'title','desc'}]. Accepts a dict
    {'title','desc'}, a ['title','desc'] pair, or a plain 'title' string."""
    out = []
    for it in (content.get("items") or []):
        if isinstance(it, dict):
            out.append({"title": str(it.get("title", "")), "desc": str(it.get("desc", "") or "")})
        elif isinstance(it, (list, tuple)):
            out.append({"title": str(it[0]) if it else "", "desc": str(it[1]) if len(it) > 1 else ""})
        else:
            out.append({"title": str(it), "desc": ""})
    return out


def _render_agenda(s, content, pal, fill, rect, text):
    """AGENDA / TABLE-OF-CONTENTS layout — header band | a numbered list of
    sections, each a large accent numeral + section title + optional one-line
    description. Mirrors the official 'Agenda' template layout. For an agenda,
    a contents page, or an ordered set of topics.

    Content schema:
    { "title":"Agenda","subtitle":"...","intro":"...",
      "items": [ {"title":"Section","desc":"optional one-liner"}, "Plain title", ... ],
      "columns": 1|2, "footer":"..." }
    """
    cycle = ["primary", "secondary", "tertiary", "dark"]
    if "title" not in content:
        content = dict(content, title="Agenda")
    top = _header(s, content, pal, rect, text)
    if content.get("intro"):
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.5),
             [[(content["intro"], 13, pal["muted"], False)]])
        top = Inches(2.05)
    items = _agenda_items(content)
    footer = content.get("footer")
    cols = max(1, min(2, int(content.get("columns") or (2 if len(items) > 5 else 1))))
    margin, gap = Inches(0.7), Inches(0.7)
    colw = int((W - 2 * margin - (cols - 1) * gap) / cols)
    bottom = Inches(6.85) if not footer else Inches(6.55)
    per = (len(items) + cols - 1) // cols if items else 0
    avail = int(bottom - (top + Inches(0.1)))
    row_h = min(Inches(1.15), int(avail / max(1, per))) if per else Inches(1.0)
    num_w = Inches(0.9)
    for i, it in enumerate(items):
        c = i // per if per else 0
        r = i % per if per else i
        x = margin + c * (colw + gap)
        y = top + Inches(0.1) + r * row_h
        if y + Inches(0.2) > bottom:
            continue
        accent = _accent(pal, cycle[i % 4])
        box = role(rect(x, y, num_w, Inches(0.8), accent, rounded=True), "focal.num")
        text(x, y, num_w, Inches(0.8), [[(f"{i + 1:02d}", 26, pal["white"], True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, role="content.num")
        tx = x + num_w + Inches(0.28)
        tw = colw - num_w - Inches(0.28)
        text(tx, y + Inches(0.02), tw, Inches(0.45),
             [[(layouts.strip_tags(it["title"]), 17, pal["ink"], True)]],
             anchor=MSO_ANCHOR.MIDDLE, role="content.title")
        if it["desc"]:
            text(tx, y + Inches(0.44), tw, Inches(0.4),
                 [[(layouts.strip_tags(it["desc"]), 12.5, pal["muted"], False)]],
                 line=1.04, role="content.body")
    _footer(s, footer, pal, rect, text)
    return {"layout": "agenda", "items": len(items), "columns": cols}


def _place_logo_centered(s, path, y, h):
    """Place a wordmark logo horizontally centered at vertical position `y`,
    height `h` (EMU). Returns True if placed. Raster-only (add_picture); never
    raises. Used by the full-bleed closing slide."""
    if not path or not os.path.isfile(path):
        return False
    if os.path.splitext(path)[1].lower() not in (".png", ".jpg", ".jpeg"):
        return False
    iw, ih = _img_size(path)
    if not (iw and ih):
        return False
    w = int(h * iw / ih)
    try:
        pic = s.shapes.add_picture(path, int((W - w) / 2), int(y), w, int(h))
        pic.name = "bg.logo"
        return True
    except Exception:
        return False


def _render_closing(s, content, pal, fill, rect, text):
    """CLOSING / THANK-YOU layout — a full-bleed branded closing slide: optional
    kicker, a large headline (default 'Thank you'), an optional supporting line,
    optional contact lines, and the brand wordmark centered near the bottom.
    Mirrors the official 'Closing Slide' template layout.

    Content schema:
    { "kicker":"...", "title":"Thank you", "subtitle":"...",
      "contact": ["Name, Role", "name@example.com", "example.com"],
      "bg":"dark|primary|secondary|tertiary|white" }
    """
    bgname = content.get("bg", "dark")
    if pal.get("bg_image"):
        # themed brand: text follows the active theme over the full-bleed art.
        ink = pal.get("title_ink", pal["white"])
        accent_ink = pal.get("rule_color", pal["primary"])
        sub_col = pal.get("muted", pal["tint"])
        contact_col = pal.get("muted", pal["white"])
    else:                                             # agnostic: bg comes from content
        light = bgname in ("white", "light", "tint")
        bg = (pal["white"] if bgname == "white"
              else pal["tint"] if bgname in ("light", "tint")
              else _accent(pal, bgname, "dark"))
        rect(0, 0, W, H, bg, role="bg.band")
        ink = pal["dark"] if light else pal["white"]
        accent_ink = pal["primary"] if light else pal["white"]
        sub_col = pal["grey"] if light else pal["tint"]
        contact_col = pal["grey"] if light else pal["white"]
    cx, cw = Inches(1.2), Inches(10.93)
    if content.get("kicker"):
        text(cx, Inches(2.0), cw, Inches(0.5),
             [[(content["kicker"], 16, accent_ink, True)]],
             align=PP_ALIGN.CENTER, role="content.kicker")
    title = content.get("title") or "Thank you"
    text(cx, Inches(2.5), cw, Inches(1.5), [[(title, 54, ink, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, role="content.title")
    if content.get("subtitle"):
        text(cx, Inches(4.0), cw, Inches(0.8),
             [[(content["subtitle"], 18, sub_col, False)]],
             align=PP_ALIGN.CENTER, line=1.12, role="content.subtitle")
    contact = content.get("contact") or []
    if contact:
        runs = [[(layouts.strip_tags(str(line)), 13, contact_col, i == 0)]
                for i, line in enumerate(contact)]
        text(cx, Inches(4.85), cw, Inches(1.1), runs, align=PP_ALIGN.CENTER,
             line=1.25, role="content.contact")
    # brand wordmark, centered near the bottom (white mark reads on the dark bg)
    _place_logo_centered(s, pal.get("logo"), Inches(6.5), Inches(0.42))
    return {"layout": "closing"}


def _render_sectionheader(s, content, pal, fill, rect, text):
    """SECTION HEADER layout — a full-bleed branded SECTION BREAK with a large
    section number. Mirrors the official 'Section Header' template: a left-aligned
    number + kicker lockup, a large title, an accent rule, and a supporting line.
    Distinct from `divider` (a centered title screen) — this is a numbered break
    INSIDE a deck, used to open a major section.

    Content schema:
    { "number":"02", "kicker":"SECTION", "title":"...", "subtitle":"...",
      "bg":"dark|primary|secondary|tertiary", "footer":"..." }
    """
    sh_t = (pal.get("design") or {}).get("sectionheader", {})
    if not pal.get("bg_image"):                       # brand art is the bg if present
        rect(0, 0, W, H, _accent(pal, content.get("bg", "dark"), "dark"), role="bg.band")
    ink = pal.get("title_ink", pal["white"])
    LX, CW = Inches(0.92), Inches(11.5)
    num = str(content.get("number", "")).strip()
    kicker = content.get("kicker")
    num_col = _tok_color(sh_t.get("numberColor", "rule_color"), pal, pal.get("rule_color", pal["primary"]))
    kick_col = _tok_color(sh_t.get("kickerColor", "rule_color"), pal, pal.get("rule_color", pal["primary"]))
    sub_col = _tok_color(sh_t.get("subColor", "muted"), pal, pal.get("tint", pal["muted"]))
    y = Inches(2.0)
    if num or kicker:
        if num:
            nshp = text(LX - Inches(0.08), Inches(1.5), Inches(4.0), Inches(1.95),
                        [[(num, sh_t.get("numberSizePt", 96), num_col, True)]],
                        anchor=MSO_ANCHOR.MIDDLE, role="bg.number")
            _ = nshp
        if kicker:
            kx = LX + (Inches(2.45) if num else Emu(0))
            text(kx, Inches(2.02), CW, Inches(0.45),
                 [[((kicker.upper() if sh_t.get("kickerCaps", True) else kicker),
                    sh_t.get("kickerSizePt", 15), kick_col, True)]], role="content.kicker")
        y = Inches(3.55)
    text(LX, y, CW, Inches(1.6),
         [[(content.get("title", ""), sh_t.get("titleSizePt", 40), ink, True)]],
         line=1.04, role="content.title")
    rule_col = _tok_color(sh_t.get("ruleColor", "rule_color"), pal, pal.get("rule_color", pal["primary"]))
    if rule_col is not None:
        rect(LX + Inches(0.02), y + Inches(1.66), Inches(1.4), Pt(sh_t.get("rulePt", 3)),
             rule_col, role="bg.rule")
    if content.get("subtitle"):
        text(LX, y + Inches(1.92), CW, Inches(1.1),
             [[(content["subtitle"], sh_t.get("subSizePt", 18), sub_col, False)]],
             line=1.18, role="content.subtitle")
    _footer(s, content.get("footer"), pal, rect, text)
    return {"layout": "sectionheader"}


def _render_cycles(s, content, pal, fill, rect, text):
    """CYCLES layout — a CIRCULAR / cyclical process: a ring of numbered nodes on
    the left (an optional center label), beside a numbered legend column on the
    right (one row per step). Mirrors the official 'Cycles' template (left ring +
    right text column). Use for a repeating loop (PDCA, a feedback cycle, a
    continuous-improvement wheel) — NOT a one-way `process` or `timeline`.

    Adaptive: the nodes place evenly around the ring for any 3–6 steps. A cycle is
    ONE visualization and is never split across slides (per the no-split rule), so
    steps beyond 6 are capped here — a longer sequence should use `process`/`timeline`,
    which the generator is steered to pick.

    Content schema:
    { "title","subtitle","intro", "center":"optional hub label",
      "steps": [ {"title":"...","body":"...","accent":"primary|..."}, ... ],
      "footer":"..." }
    """
    import math
    cy_t = (pal.get("design") or {}).get("cycles", {})
    cyc = ["primary", "secondary", "tertiary", "dark"]
    _header(s, content, pal, rect, text)
    top = Inches(1.7)
    if content.get("intro"):
        text(Inches(0.55), Inches(1.45), Inches(12.23), Inches(0.4),
             [[(content["intro"], _dtok(pal, "type", "introSizePt", 12.5), pal["muted"], False)]])
        top = Inches(2.05)
    steps = (content.get("steps") or [])[:6]          # cap: a cycle stays whole on one slide (never split)
    n = max(1, len(steps))
    footer = content.get("footer")
    bottom = Inches(6.85) if not footer else Inches(6.5)
    white = _tok_color(cy_t.get("numberColor", "#FFFFFF"), pal, pal["white"])

    # ---- LEFT: ring of numbered nodes (the cycle) ----
    region_w = Inches(5.7)
    ccx = int(Inches(0.4) + region_w / 2)
    ccy = int((top + bottom) / 2)
    R = Inches(cy_t.get("ringRadiusIn", 1.95))
    node_d = Inches(cy_t.get("nodeDiaIn", 0.92))
    dot_d = Inches(cy_t.get("dotDiaIn", 0.16))
    dot_col = _tok_color(cy_t.get("dotColor", "rule_color"), pal, pal.get("rule_color", pal["primary"]))
    hub_d = Inches(cy_t.get("hubDiaIn", 1.95))
    hub_col = _tok_color(cy_t.get("hubColor", "card"), pal, pal.get("card", pal["white"]))
    if hub_col is not None:
        hub = s.shapes.add_shape(MSO_SHAPE.OVAL, ccx - int(hub_d / 2), ccy - int(hub_d / 2), hub_d, hub_d)
        fill(hub, hub_col); _fill_alpha(hub, cy_t.get("hubAlpha", 100)); role(hub, "bg.hub")
        # Outline the hub so it stays visible when its fill matches the background
        # (a light-theme `card` hub on a light bg was an invisible disc otherwise).
        hub_oc = _tok_color(cy_t.get("hubOutlineColor", "rule_color"), pal,
                            pal.get("rule_color", pal.get("muted", pal["dark"])))
        if hub_oc is not None:
            hub.line.width = Pt(cy_t.get("hubOutlinePt", 1.25))
            hub.line.color.rgb = hub_oc
    if content.get("center"):
        text(ccx - int(hub_d / 2), ccy - int(hub_d / 2), hub_d, hub_d,
             [[(layouts.strip_tags(content["center"]), cy_t.get("centerSizePt", 13),
                _tok_color(cy_t.get("centerColor", "ink"), pal, pal["ink"]), True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.0, role="content.center")
    # connector dots at the midpoints between consecutive nodes -> reads as a loop
    if dot_col is not None and n > 1:
        for i in range(n):
            a = -math.pi / 2 + 2 * math.pi * (i + 0.5) / n
            dx = ccx + int(R * math.cos(a)); dy = ccy + int(R * math.sin(a))
            d = s.shapes.add_shape(MSO_SHAPE.OVAL, dx - int(dot_d / 2), dy - int(dot_d / 2), dot_d, dot_d)
            fill(d, dot_col); role(d, f"bg.dot.{i + 1}")
    for i, st in enumerate(steps):
        u = i + 1
        a = -math.pi / 2 + 2 * math.pi * i / n
        nx = ccx + int(R * math.cos(a)); ny = ccy + int(R * math.sin(a))
        acc = _accent(pal, st.get("accent", cyc[i % 4]))
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, nx - int(node_d / 2), ny - int(node_d / 2), node_d, node_d)
        fill(node, acc); role(node, f"focal.node.{u}")
        text(nx - int(node_d / 2), ny - int(node_d / 2), node_d, node_d,
             [[(str(u), cy_t.get("nodeNumPt", 19), white, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, role=f"content.number.{u}")

    # ---- RIGHT: numbered legend column (one row per step) ----
    LXr, LWr = Inches(6.85), Inches(6.0)
    row_h = int((bottom - top) / n)
    chip_d = Inches(cy_t.get("chipDiaIn", 0.46))
    tcol = _tok_color(cy_t.get("legendTitleColor", "ink"), pal, pal["ink"])
    bcol = _tok_color(cy_t.get("legendBodyColor", "muted"), pal, pal["muted"])
    rule_col = _tok_color(cy_t.get("ruleColor", "rail"), pal, pal.get("rail", pal["muted"]))
    minb = cy_t.get("minBodyPt", 9)
    for i, st in enumerate(steps):
        u = i + 1
        ry = top + i * row_h
        acc = _accent(pal, st.get("accent", cyc[i % 4]))
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, LXr, ry + Inches(0.04), chip_d, chip_d)
        fill(chip, acc); role(chip, f"focal.chip.{u}")
        text(LXr, ry + Inches(0.04), chip_d, chip_d,
             [[(str(u), 14, white, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             role=f"content.chipnum.{u}")
        tx = LXr + chip_d + Inches(0.22)
        tw = int(LXr + LWr - tx)
        text(tx, ry, tw, Inches(0.4),
             [[(layouts.strip_tags(st.get("title", "")), cy_t.get("legendTitlePt", 14.5), tcol, True)]],
             role=f"content.title.{u}")
        body = layouts.strip_tags(st.get("body", ""))
        by = ry + Inches(0.4)
        bh = max(int(Inches(0.2)), int(ry + row_h - Inches(0.12) - by))
        bpt, body = _fit_body(body, tw, bh, cy_t.get("legendBodyPt", 11.5), minb)
        text(tx, by, tw, bh, [[(body, bpt, bcol, False)]], line=1.05, role=f"content.body.{u}")
        if i < n - 1 and rule_col is not None:
            rect(tx, ry + row_h - Inches(0.04), tw, Pt(cy_t.get("rulePt", 1)), rule_col,
                 role=f"bg.rule.{u}")
    _footer(s, footer, pal, rect, text)
    return {"layout": "cycles", "steps": len(steps)}


RENDERERS = {
    "infographic": _render_infographic,
    "process": _render_process,
    "comparison": _render_comparison,
    "timeline": _render_timeline,
    "divider": _render_divider,
    "chart": _render_chart,
    "image": _render_image,
    "imagetext": _render_imagetext,
    "cards": _render_cards,
    "quote": _render_quote,
    "statement": _render_statement,
    "bullets": _render_bullets,
    "agenda": _render_agenda,
    "closing": _render_closing,
    "sectionheader": _render_sectionheader,
    "cycles": _render_cycles,
}


def export_slide_file(content_path, out, brand=None, layout="infographic", **kw):
    with open(content_path, encoding="utf-8") as fh:
        content = json.load(fh)
    return export_slide(content, out, brand=brand, layout=layout, **kw)


def export_deck_file(content_path, out, brand=None, **kw):
    """Build a deck from a JSON file shaped {"slides": [{"layout","content"}, ...]}
    (a bare list of slide specs is also accepted). An `images` key on the deck
    object supplies a default images folder if `images_dir` wasn't passed."""
    with open(content_path, encoding="utf-8") as fh:
        data = json.load(fh)
    slides = data.get("slides") if isinstance(data, dict) else data
    if isinstance(data, dict):
        if data.get("images") and not kw.get("images_dir"):
            kw["images_dir"] = data["images"]
        if data.get("animate") and not kw.get("animate"):
            kw["animate"] = data["animate"]
    return export_deck(slides or [], out, brand=brand, **kw)
