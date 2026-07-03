"""Infographic template family (region-specs in templates/layouts/infographic-*.json).

Our own original, brand-agnostic data-viz slide types (Step 3): a four-up stat
grid, a two-column comparison, a four-step process, and a single stat callout.
Proves each loads, registers as a generic ("default") pickable template, renders
the same spec under two brands, repeats cards over their array, and renders in
the SVG preview.
"""
import pytest

import slide_layouts
import slide_svg
import brand as brandmod

pptx = pytest.importorskip("pptx")  # export extra

INFOGRAPHIC = ["infographic-statgrid", "infographic-comparison",
               "infographic-process", "infographic-callout"]


def _pptx_texts(path):
    from pptx import Presentation
    p = Presentation(path)
    return " ".join(sh.text_frame.text for s in p.slides
                    for sh in s.shapes if sh.has_text_frame)


def test_all_infographic_specs_listed_and_generic():
    names = slide_layouts.list_template_layouts()
    info = {r["name"]: r for r in slide_layouts.template_layout_info()}
    for n in INFOGRAPHIC:
        assert n in names, n
        spec = slide_layouts.load_template_spec(n)
        assert spec["name"] == n and spec["regions"]
        assert info[n]["category"] == "default" and info[n]["brand"] is None


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_statgrid_repeats_four_stats(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"statgrid_{brand_name}.pptx")
    content = {"template": "infographic-statgrid", "title": "Impact",
               "stats": [{"title": "6x", "body": "Faster"},
                         {"title": "40%", "body": "Fewer calls"},
                         {"title": "12k", "body": "Beds"},
                         {"title": "98%", "body": "Uptime"}]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Impact", "6x", "40%", "Beds", "Uptime"):
        assert w in txt


def test_comparison_repeats_two_columns(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "comparison.pptx")
    content = {"template": "infographic-comparison", "kicker": "Before / after",
               "title": "Manual vs. Nova",
               "columns": [{"title": "Manual", "body": "Phone tag"},
                           {"title": "Nova", "body": "One queue"}]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("Before / after", "Manual", "Nova", "Phone tag", "One queue"):
        assert w in txt


@pytest.mark.parametrize("brand_name", ["teletracking", "_default"])
def test_process_repeats_four_steps(tmp_path, brand_name):
    br = brandmod.load_brand(brand_name)
    out = str(tmp_path / f"process_{brand_name}.pptx")
    content = {"template": "infographic-process", "title": "How it works",
               "subtitle": "Four steps, start to finish.",
               "steps": [{"title": "Request", "body": "Open"},
                         {"title": "Match", "body": "Assign"},
                         {"title": "Move", "body": "Transport"},
                         {"title": "Close", "body": "Confirm"}]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("How it works", "Request", "Match", "Transport", "Confirm"):
        assert w in txt


def test_callout_renders_figure_and_points(tmp_path):
    br = brandmod.load_brand("teletracking")
    out = str(tmp_path / "callout.pptx")
    content = {"template": "infographic-callout", "value": "73%",
               "title": "Adoption in 90 days",
               "points": ["Across 12 units", "No new headcount"]}
    slide_layouts.export_deck([{"layout": "template", "content": content}], out, brand=br)
    txt = _pptx_texts(out)
    for w in ("73%", "Adoption in 90 days", "Across 12 units", "No new headcount"):
        assert w in txt


@pytest.mark.parametrize("name,content", [
    ("infographic-statgrid", {"title": "Stats",
                              "stats": [{"title": "1", "body": "a"},
                                        {"title": "2", "body": "b"}]}),
    ("infographic-comparison", {"kicker": "vs", "title": "A vs B",
                                "columns": [{"title": "A", "body": "x"},
                                            {"title": "B", "body": "y"}]}),
    ("infographic-process", {"title": "Process", "subtitle": "Steps",
                             "steps": [{"title": "One", "body": "go"}]}),
    ("infographic-callout", {"value": "50%", "title": "Half",
                             "points": ["one", "two"]}),
])
def test_each_infographic_spec_renders_in_svg_preview(name, content):
    br = brandmod.load_brand("teletracking")
    content = dict(content, template=name)
    svg = slide_svg.render_slide_svg("template", content, br)
    assert svg.startswith("<svg")
    assert content["title"].split()[0] in svg
