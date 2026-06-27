"""Per-element ENTRANCE animations for slides, injected as OOXML.

The course player has simple entrance effects (fade / float-in, slide-in from a
side, cascaded across blocks). This brings the same feel to the .pptx deck:
shapes animate in automatically on slide entry — no click needed. python-pptx has
no animation API, so (like pptx_transitions) we append a `<p:timing>` tree to the
slide part directly.

GROUPED FLOW (2026-06-25): the old build animated every shape strictly one after
another (~0.65s each), so a busy slide took ~8–9s to assemble and each piece read
as disconnected. Now shapes are clustered by SPATIAL PROXIMITY — a card's panel +
its accent strip + title + body, or a bullet's dot + its text line, physically
overlap (or sit adjacent), so they fall into one cluster and enter TOGETHER. Only
clusters are staggered, top-to-bottom, by a short offset that OVERLAPS the fade
(stagger < dur) — a quick, subtle cascade (header → content → footer). A typical
slide now lands in ~1–1.5s. Grouping needs no renderer changes; it reads geometry.

Design choices for robustness (this XML is hand-authored and can't be validated
against real PowerPoint in this environment — James is the oracle for a real
PowerPoint check, per the project gotcha log):
- Canonical build: tmRoot → interactive seq (mainSeq). PowerPoint populates the
  Animation pane from the mainSeq; a bare timeline is silently ignored.
- A cluster LEADER node starts the cluster; its members are `withEffect` (begin
  WITH the leader → simultaneous). Each later leader is `withEffect` with a
  `delay` = the inter-cluster stagger, so leaders chain at a fixed offset while
  members ride their leader. (withEffect begins with the previous sibling; the
  delay is measured from that begin — so stagger accumulates leader-to-leader.)
- Only safe entrance effects: a fade, plus an optional small motion (float up /
  fly in from a side) layered on the fade. Motion offsets are kept subtle.
- Every shape is targeted by its spid (shape id). The injection is wrapped by the
  caller in try/except per slide, and a malformed build fails loudly at
  parse_xml time (never a silent corrupt save).

Usage:
    import pptx_animate
    pptx_animate.apply(slide, "rise")              # one slide
    pptx_animate.apply_all(prs, "fade", speed="med")
"""
import json
import os
import re

from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

EFFECTS = ("none", "fade", "rise", "flyleft", "flyright")
# FADE duration (ms) — how fast one group's shapes appear (500 = James's reference
# default). Sub-steps are placed at absolute times `i * (dur + gap)`.
SPEEDS = {"slow": 700, "med": 500, "fast": 350}
# GAP (ms) added on top of the fade before the next group starts.
STAGGERS = {"slow": 250, "med": 100, "fast": 40}
# A shape this thin in either dimension, OR this small in area, is treated as a
# DECORATION/CONNECTOR (a dot, number bubble, accent rule, chevron, baseline). Such
# shapes never bind two groups together — they ATTACH to their nearest content
# shape — so a connector between two cards can't fuse the cards into one group.
_SMALL_MIN = Inches(0.22)
_SMALL_AREA = Inches(0.36) * Inches(0.36)
_BIG = 1 << 62
_SLIDE_W = Inches(13.333)               # 16:9 deck width; left half vs right half
_SLIDE_MID = _SLIDE_W / 2
_SLIDE_MIDY = Inches(7.5) / 2           # above/below the slide's vertical middle
_TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates", "anim")


def _remove_timing(slide):
    sld = slide._element
    for ex in sld.findall(qn("p:timing")):
        sld.remove(ex)


def _box(shape):
    """(left, top, right, bottom) in EMU, or None if the shape isn't positioned."""
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return None
    if None in (l, t, w, h):
        return None
    return (l, t, l + w, t + h)


def _is_small(b):
    """A decoration/connector — thin in a dimension or tiny in area."""
    w, h = b[2] - b[0], b[3] - b[1]
    return min(w, h) <= _SMALL_MIN or w * h <= _SMALL_AREA


def _strong_overlap(a, b):
    """True only when the boxes genuinely OVERLAP and the intersection covers >=40%
    of the smaller box — i.e. one sits inside/over the other (a card's title over its
    panel), NOT merely abut (stacked bullet rows touch but don't overlap)."""
    vo = min(a[3], b[3]) - max(a[1], b[1])
    ho = min(a[2], b[2]) - max(a[0], b[0])
    if vo <= 0 or ho <= 0:
        return False
    inter = vo * ho
    smaller = min((a[3] - a[1]) * (a[2] - a[0]), (b[3] - b[1]) * (b[2] - b[0]))
    return smaller > 0 and inter * 5 >= smaller * 2


def _rect_gap2(a, b):
    """Squared edge-to-edge gap between two boxes (0 if they touch/overlap)."""
    dx = max(0, a[0] - b[2], b[0] - a[2])
    dy = max(0, a[1] - b[3], b[1] - a[3])
    return dx * dx + dy * dy


def cluster_shapes(shapes):
    """Group shapes into logical clusters that should enter together.

    Two-phase, geometry-only (no renderer changes):
      1) UNION content shapes that strongly overlap — a card's panel/accent/title/
         body, a header band with its title+subtitle. Merely touching (stacked
         bullet rows, a header band and the row beneath it) does NOT union, so
         lists stay one-cluster-per-row and the header stays separate.
      2) ATTACH each leftover small decoration/connector (dot, number bubble, rule,
         chevron, timeline baseline) to its NEAREST content shape's cluster. A
         connector binds to one group instead of bridging two.

    Returns clusters as lists of shapes, ordered top-to-bottom then left-to-right;
    shapes inside a cluster keep document (creation) order. Unpositioned shapes each
    form their own singleton so they're never dropped."""
    n = len(shapes)
    boxes = [_box(s) for s in shapes]
    parent = list(range(n))

    def find(i):
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    def union(i, j):
        ri, rj = find(i), find(j)
        if ri != rj:
            parent[max(ri, rj)] = min(ri, rj)

    # Phase 1 — union by strong overlap (containment), content shapes included.
    for i in range(n):
        if boxes[i] is None:
            continue
        for j in range(i + 1, n):
            if boxes[j] is not None and _strong_overlap(boxes[i], boxes[j]):
                union(i, j)

    # Phase 2 — any cluster made up ENTIRELY of decorations/connectors (a lone dot,
    # a ring+digit number bubble, a baseline) attaches to the nearest CONTENT
    # shape's cluster, so it rides a real group instead of animating on its own or
    # bridging two groups. Targets are computed from the phase-1 grouping, then
    # applied, so attach order doesn't matter.
    content = [i for i in range(n) if boxes[i] is not None and not _is_small(boxes[i])]
    content_roots = {find(c) for c in content}
    phase1 = {}
    for i in range(n):
        if boxes[i] is not None:
            phase1.setdefault(find(i), []).append(i)
    attach = []
    for root, members in phase1.items():
        if root in content_roots:
            continue                       # cluster already anchored by content
        best, best_d = None, None
        for c in content:
            if find(c) == root:
                continue
            d = min(_rect_gap2(boxes[m], boxes[c]) for m in members)
            if best_d is None or d < best_d:
                best, best_d = c, d
        if best is not None:
            attach.append((members[0], best))
    for m, c in attach:
        union(m, c)

    groups = {}
    for idx in range(n):
        groups.setdefault(find(idx), []).append(idx)

    def sort_key(members):
        ys = [boxes[m][1] for m in members if boxes[m] is not None]
        xs = [boxes[m][0] for m in members if boxes[m] is not None]
        return (min(ys) if ys else _BIG, min(xs) if xs else _BIG)

    ordered = sorted(groups.values(), key=sort_key)
    return [[shapes[m] for m in sorted(members)] for members in ordered]


# Effects are built with a UNIFORM motion path per group, NOT PowerPoint's per-shape
# Fly In (which moves each shape by its own width and so de-synchronises a group).
# Every shape in a sub-step gets the IDENTICAL translation, so the group enters as
# one rigid unit — the look James gets by grouping shapes in PowerPoint. Direction →
# a relative motion path (slide-fraction offsets); position decides left vs right.
_FLY_PATH = {           # Fly In: slide in from off the matching edge
    "left":  ("2", "8", "M -0.55 0 L 0 0"),
    "right": ("2", "2", "M 0.55 0 L 0 0"),
    "top":   ("2", "1", "M 0 -0.55 L 0 0"),
    "bottom":("2", "4", "M 0 0.55 L 0 0"),
}
_FLOAT_PATH = {         # Float In: fade + a short drift into place
    "up":    ("42", "0", "M 0 0.08 L 0 0"),
    "down":  ("42", "0", "M 0 -0.08 L 0 0"),
    "left":  ("42", "0", "M 0.08 0 L 0 0"),
    "right": ("42", "0", "M -0.08 0 L 0 0"),
}


def _set(spid, nid):
    return (f'<p:set><p:cBhvr><p:cTn id="{nid()}" dur="1" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>')


def _animeffect(spid, dur, nid, filt):
    return (f'<p:animEffect transition="in" filter="{filt}"><p:cBhvr>'
            f'<p:cTn id="{nid()}" dur="{dur}"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'</p:cBhvr></p:animEffect>')


def _anim(spid, dur, nid, attr, v0, v1):
    def val(v):
        return f'<p:fltVal val="0"/>' if v == "0" else f'<p:strVal val="{v}"/>'
    return (f'<p:anim calcmode="lin" valueType="num"><p:cBhvr>'
            f'<p:cTn id="{nid()}" dur="{dur}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>{attr}</p:attrName></p:attrNameLst></p:cBhvr>'
            f'<p:tavLst><p:tav tm="0"><p:val>{val(v0)}</p:val></p:tav>'
            f'<p:tav tm="100000"><p:val>{val(v1)}</p:val></p:tav></p:tavLst></p:anim>')


def _animmotion(spid, dur, nid, path):
    """A UNIFORM relative translation — identical for every shape in the group, so
    the whole group moves as one rigid unit (no per-shape de-sync)."""
    return (f'<p:animMotion origin="layout" path="{path}" pathEditMode="relative">'
            f'<p:cBhvr><p:cTn id="{nid()}" dur="{dur}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>ppt_x</p:attrName>'
            f'<p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr></p:animMotion>')


def _effect_body(spid, kind, direction, dur, nid):
    """Return (presetID, presetSubtype, behaviors_xml) for one shape's entrance."""
    if kind == "fade":
        return "10", "0", _set(spid, nid) + _animeffect(spid, dur, nid, "fade")
    if kind == "floatin":
        preset, sub, path = _FLOAT_PATH.get(direction or "up", _FLOAT_PATH["up"])
        return preset, sub, (_set(spid, nid) + _animeffect(spid, dur, nid, "fade")
                             + _animmotion(spid, dur, nid, path))
    if kind == "flyin":
        preset, sub, path = _FLY_PATH.get(direction or "left", _FLY_PATH["left"])
        return preset, sub, _set(spid, nid) + _animmotion(spid, dur, nid, path)
    if kind == "stretch":
        if direction in ("bottom", "top"):
            # vertical grow from the bottom (sub4) or top (sub1) edge — height 0->full
            # while the center shifts so the anchored edge stays put (James's connector).
            frm_y = "#ppt_y+#ppt_h/2" if direction == "bottom" else "#ppt_y-#ppt_h/2"
            body = (_set(spid, nid)
                    + _anim(spid, dur, nid, "ppt_x", "#ppt_x", "#ppt_x")
                    + _anim(spid, dur, nid, "ppt_y", frm_y, "#ppt_y")
                    + _anim(spid, dur, nid, "ppt_w", "#ppt_w", "#ppt_w")
                    + _anim(spid, dur, nid, "ppt_h", "0", "#ppt_h"))
            return "17", ("4" if direction == "bottom" else "1"), body
        # horizontal stretch from center (sub10): width 0->full, height held (the rail)
        body = (_set(spid, nid)
                + _anim(spid, dur, nid, "ppt_w", "0", "#ppt_w")
                + _anim(spid, dur, nid, "ppt_h", "#ppt_h", "#ppt_h"))
        return "17", "10", body
    if kind == "blinds":
        return "3", "10", _set(spid, nid) + _animeffect(spid, dur, nid, "blinds(horizontal)")
    if kind == "wipe":
        d = direction if direction in ("up", "down", "left", "right") else "down"
        sub = {"down": "4", "up": "1", "left": "8", "right": "2"}[d]
        return "22", sub, _set(spid, nid) + _animeffect(spid, dur, nid, f"wipe({d})")
    return "10", "0", _set(spid, nid) + _animeffect(spid, dur, nid, "fade")


def _effect_par(spid, kind, direction, dur, nid, node_type):
    """One shape's entrance EFFECT (innermost par). `node_type` is relative to its
    siblings inside the sub-step: leader 'afterEffect'/'clickEffect', rest
    'withEffect' (so the whole group enters together)."""
    eff_id = nid()
    preset, sub, body = _effect_body(spid, kind, direction, dur, nid)
    return (
        f'<p:par><p:cTn id="{eff_id}" presetID="{preset}" presetClass="entr" '
        f'presetSubtype="{sub}" fill="hold" grpId="0" nodeType="{node_type}">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>{body}</p:childTnLst></p:cTn></p:par>')


def _substep_par(spids, kind, direction, dur, abs_delay, nid, first_node):
    """One BUILD SUB-STEP: a wrapper par at absolute time `abs_delay` (ms) whose
    children are the effects for `spids`, all entering TOGETHER (leader = first_node
    — 'afterEffect' for auto, 'clickEffect' for on-click — the rest 'withEffect').
    This container→sub-step→effect nesting is the structure real PowerPoint uses;
    a flat list of effects under the mainSeq plays one-at-a-time. See
    reference_pptx_entrance_animation_structure.md."""
    sid = nid()
    effs = "".join(
        _effect_par(sp, kind, direction, dur, nid,
                    (first_node if i == 0 else "withEffect"))
        for i, sp in enumerate(spids))
    # The sub-step wrapper carries only the timing (cond delay); the nodeType lives
    # on the effects (leader afterEffect/clickEffect, members withEffect) — matching
    # how PowerPoint itself writes it.
    return (
        f'<p:par><p:cTn id="{sid}" fill="hold">'
        f'<p:stCondLst><p:cond delay="{abs_delay}"/></p:stCondLst>'
        f'<p:childTnLst>{effs}</p:childTnLst></p:cTn></p:par>')


def _emit(slide, steps, on_click):
    """Render an ordered list of steps -> the slide's <p:timing>. Each step is
    (spids, effect, direction, dur); steps play back-to-back (each at the cumulative
    sum of prior durations — James's verified cadence). The whole build auto-starts
    on slide entry unless on_click."""
    steps = [s for s in steps if s[0]]
    if not steps:
        return
    counter = [2]                       # ids 1,2 reserved for tmRoot + mainSeq

    def nid():
        counter[0] += 1
        return counter[0]

    cid = nid()                         # the single build-step container id
    parts, t = [], 0
    for i, (spids, effect, direction, dur) in enumerate(steps):
        first = "clickEffect" if (on_click and i == 0) else "afterEffect"
        parts.append(_substep_par(spids, effect, direction, dur, t, nid, first))
        t += dur
    container = (
        f'<p:par><p:cTn id="{cid}" fill="hold">'
        f'<p:stCondLst><p:cond delay="{"indefinite" if on_click else 0}"/></p:stCondLst>'
        f'<p:childTnLst>{"".join(parts)}</p:childTnLst></p:cTn></p:par>')
    xml = (
        f'<p:timing {nsdecls("p")}><p:tnLst>'
        f'<p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{container}</p:childTnLst>'
        f'</p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq>'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:tnLst></p:timing>')
    el = parse_xml(xml)                 # raises on malformed XML — never a silent corrupt save
    _remove_timing(slide)
    slide._element.append(el)


def load_template(layout):
    """Load a VERIFIED animation sequence for `layout` from templates/anim/<layout>.json,
    or None if there's no library entry yet (caller falls back to the auto-planner)."""
    if not layout:
        return None
    path = os.path.join(_TEMPLATE_DIR, f"{layout}.json")
    if not os.path.exists(path):
        return None
    try:
        with open(path, encoding="utf-8") as fh:
            return json.load(fh)
    except (OSError, ValueError):
        return None


def _resolve_dir(token, above):
    """Resolve a template direction token. Positional tokens depend on whether the
    unit sits above the slide's vertical middle: axisV = stretch grows from the axis
    side (bottom for above, top for below); outV = wipe reveals away from the axis
    (down for above, up for below)."""
    if token == "axisV":
        return "bottom" if above else "top"
    if token == "outV":
        return "down" if above else "up"
    return token


def _steps_from_template(slide, tmpl):
    """Map a verified template onto THIS rendered slide via role names. Per-unit
    blocks repeat for every unit present; positional directions resolve from each
    unit's geometry. Roles absent on the slide are skipped, so optional parts
    (a milestone with no phase chip) just drop out."""
    role2spid, top = {}, {}
    for sh in slide.shapes:
        if sh.name:
            role2spid[sh.name] = sh.shape_id
            b = _box(sh)
            if b:
                top[sh.name] = (b[1] + b[3]) / 2
    units = set()
    for name in role2spid:
        m = re.search(r"\.(\d+)$", name)
        if m:
            units.add(int(m.group(1)))
    n = max(units) if units else 0

    def above(u):
        cy = top.get(f"focal.card.{u}", top.get(f"focal.node.{u}"))
        return True if cy is None else cy < _SLIDE_MIDY

    def resolve(roles, ab):
        return [role2spid[r] for r in roles if r in role2spid]

    steps = []
    for st in tmpl.get("steps", []):
        if st.get("per_unit"):
            for u in range(1, n + 1):
                ab = above(u)
                for b in st.get("block", []):
                    roles = [r.replace("{u}", str(u)) for r in b["roles"]]
                    spids = resolve(roles, ab)
                    if spids:
                        steps.append((spids, b["effect"], _resolve_dir(b.get("dir"), ab), b["dur"]))
        else:
            spids = resolve(st["roles"], True)
            if spids:
                steps.append((spids, st["effect"], _resolve_dir(st.get("dir"), True), st["dur"]))
    return steps


def _is_text(shape):
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    except Exception:
        return False


def _side(shapes):
    """Which half of the slide a group sits on -> the edge it should enter from, so
    it never crosses over other content. Full-width groups read as centered."""
    boxes = [b for b in (_box(s) for s in shapes) if b]
    if not boxes:
        return "left"
    left = min(b[0] for b in boxes)
    right = max(b[2] for b in boxes)
    if left < _SLIDE_MID < right and (right - left) > _SLIDE_W * 0.6:
        return "up"                      # spans the slide -> drift up, no side bias
    return "left" if (left + right) / 2 < _SLIDE_MID else "right"


def _plan_substeps(slide, effect):
    """Turn the slide's clusters into an ordered list of sub-steps, each
    (spids, effect_kind, direction). One sub-step = one GROUP that enters together
    as a rigid unit. Recipe mirroring James's reference:
      - header (cluster 0): band/rule, then title/subtitle — both Float In (up);
      - every other multi-part cluster: its STRUCTURAL shapes enter as one group
        (Fly In from THE SIDE THEY SIT ON, so nothing crosses over), then its TEXT
        fades in right after;
      - a cluster that is just a marker + its text (a bullet) stays ONE group so the
        bullet and its text enter together, from the side it sits on.
    `effect` picks the shell motion: rise/float -> Float In, fly* -> Fly In, fade -> Fade."""
    shell = {"rise": "floatin", "fade": "fade", "flyleft": "flyin",
             "flyright": "flyin", "fly": "flyin"}.get(effect, "floatin")
    plan = []
    for ci, cluster in enumerate(cluster_shapes(slide.shapes)):
        nontext = [s for s in cluster if not _is_text(s)]
        text = [s for s in cluster if _is_text(s)]
        nt_box = [_box(s) for s in nontext]
        substantial = len(nontext) >= 2 or any(b and not _is_small(b) for b in nt_box)
        if ci == 0:
            # Header: structural band/rule, then title/subtitle — both Float In.
            if nontext:
                plan.append(([s.shape_id for s in nontext], "floatin", "up"))
            if text:
                plan.append(([s.shape_id for s in text], "floatin", "up"))
        elif nontext and text and substantial:
            plan.append(([s.shape_id for s in nontext], shell, _side(nontext)))
            plan.append(([s.shape_id for s in text], "fade", None))
        else:
            # Lone marker + text (bullet), or a single-kind cluster: ONE group,
            # entering together from the side it sits on. A wide/centered group
            # (a full-width bullet) drifts up gently rather than flying far.
            side = _side(cluster)
            if not nontext:
                kind = "fade"
            elif side == "up":
                kind = "floatin"
            else:
                kind = shell
            plan.append(([s.shape_id for s in cluster], kind, side))
    return plan


def apply(slide, effect="rise", speed="med", stagger_ms=None, on_click=False, layout=None):
    """Inject a grouped, sequenced entrance for `slide`.

    If the layout has a VERIFIED sequence in the library (templates/anim/<layout>.json),
    REPLAY it — map its role-based steps onto this slide and reproduce James's exact
    choreography for any content. Otherwise fall back to the geometry-based
    auto-planner (header, then each unit's shell + text). The build auto-starts on
    slide entry unless `on_click`. effect='none' removes any existing animation."""
    effect = (effect or "rise").lower()
    if effect == "none":
        _remove_timing(slide)
        return "none"
    if effect not in EFFECTS:
        raise ValueError("unknown animation %r (have: %s)" % (effect, ", ".join(EFFECTS)))
    if not list(slide.shapes):
        return effect
    tmpl = load_template(layout)
    if tmpl:
        steps = _steps_from_template(slide, tmpl)
        on_click = on_click or not tmpl.get("autostart", True)
    else:
        dur = SPEEDS.get(speed, 400)
        steps = [(spids, kind, direction, dur)
                 for (spids, kind, direction) in _plan_substeps(slide, effect)]
    _emit(slide, steps, on_click)
    return effect


def apply_all(prs, effect="rise", speed="med", stagger_ms=None, on_click=False):
    """Apply the same entrance to every slide in a presentation."""
    for s in prs.slides:
        apply(s, effect, speed, stagger_ms, on_click)
    return effect
